# -*- coding: utf-8 -*-
"""
File name: GenerateGraphs.py
Package: resq
Written by: Marie Jankujova - jankujova.marie@fnusa.cz on 11-2017
Version: v1.0
Version comment: The first version of script which is used to generate charts in the PPTX format. 
Description: This script is used to generate PPTX file from the dataframe with computed statistics. 
This script is used by class GeneratePresentation.py. 
As a result, the presentation with chart is produced. 
"""


import pandas as pd
import sys
import os
from datetime import datetime, date
import sqlite3

import xlsxwriter

import csv

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.oxml.xmlchemy import OxmlElement


class GenerateGraphs:
    """ The class generating presentation with graphs for general reports.

    :param dataframe: the dataframe with calculated statistics
    :type dataframe: pandas dataframe
    :param presentation: the opened presentation document
    :type presentation: Presentation object
    :param title: the title of the slide
    :type title: str
    :param column_name: the column name from the dataframe to be shown in the graph
    :type column_name: str
    :param graph_type: the type of graph to be generated
    :type graph_type: str
    :param number_of_series: the number of columns to be shown in the stacked barplot
    :type number_of_series: int
    :param legend: the list of names to be used in the legend in the stacked barplot
    :type legend: list
    :param country: the country name used in the first slide
    :type country: str
    """

    def __init__(self, dataframe, presentation, title, column_name, graph_type = None, number_of_series=0, legend=None, country=None):

        self.dataframe = dataframe
        self.presentation = presentation
        self.title = title
        self.column_name = column_name
        self.number_of_series = number_of_series
        self.legend = legend
        self.country_name = country
        self.font_name = 'Century Gothic'

        # If country name is Ukraine or Poland set as categories value Site IDs not Site Name due to huge amount of sites in graph
        # if (self.country_name in ['Ukraine', 'Poland'] and len(self.dataframe) > 2):
        #     self.categories_column = 'Site ID'
        # else:
        self.categories_column = 'Site Name'
            
        # Estimate font sizes based on number of sites included in the graph
        if (len(self.dataframe) > 60):
            self.category_font_size = Pt(4)
            self.data_label_font_size = Pt(4)
        elif (len(self.dataframe) > 50 and len(self.dataframe) <= 60):
            self.category_font_size = Pt(6)
            self.data_label_font_size = Pt(6)
        else:
            self.category_font_size = Pt(8)
            self.data_label_font_size = Pt(8)

        # Select graph which should be exported
        if (graph_type == 'stacked'):
            self._create_stacked_barplot(dataframe=self.dataframe, title=self.title, column_name=self.column_name, legend=self.legend, number_of_series=self.number_of_series)
        else:
            self._create_barplot(dataframe=self.dataframe, title=self.title, column_name=self.column_name)


    def _get_length_of_legend(self, legend):
        """ The function adjusting the number of letters in legend to quess the number of columns in the legend! 
        
        :param legend: the names of legend
        :type legend: list
        :returns: the adjusted number of letters
        """
        count = 0

        for i in legend:
            count = count + len(i)
        
        return count


    def _set_transparency(self, transparency, elm):
        """ The function set the transparency of the row. 

        :param transparency: the transparency in %
        :type transparency: int
        :param elm: the element which transparency should be changed
        :type elm: format.line.color._xFill
        """
        a = str(100 - transparency) + '196'
        
        alpha = OxmlElement('a:alpha')
        alpha.set('val', a)
        elm.srgbClr.append(alpha)


    def _create_barplot(self, dataframe, title, column_name):
        """ The function creating the normal barplot graph into the presentation based on the graph type. 
        
        :param df: the dataframe with data to be shown
        :type df: pandas dataframe
        :param title: the title of the graph
        :type title: str
        :param column_name: the column name to be displayed in the graph
        :type column_name: str
        """
        maximum = 0

        # If graph is in %, set maximum valut to 100. 
        if '%' in title.lower():
            maximum = 100
            if self.country_name == 'Czech Republic':
                values = [round(x, 0) for x in dataframe[column_name].tolist()]
        else:
            maximum = round((max(dataframe[column_name].tolist())),1)
            if self.country_name == 'Czech Republic':
                values = dataframe[column_name].tolist()

        # Add slide to presentation (layout 11 is our custom layout where only title 'Agency FB', color: RGBColor(43, 88, 173)  and size:24 is set)
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        # Get title object
        title_placeholders = slide.shapes.title
        # Set title
        title_placeholders.text = title

        chart_data = ChartData()
        chart_data.categories = dataframe[self.categories_column].tolist()
        chart_data.add_series(column_name, dataframe[column_name].tolist())

        # Add chart on slide
        specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
            }
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        # Get series of chart
        series = chart.series[0]
        # If graphs for whole country are generated, set for bar with country with red color
        # else set to blue color (same color as title uses)
        site_names = dataframe[self.categories_column].tolist()
        if (len(dataframe) > 2):
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                if (site_names[idx] == self.country_name):
                    fill.fore_color.rgb = RGBColor(128,0,0)
                else:
                    fill.fore_color.rgb = RGBColor(43, 88, 173)
        else:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(43, 88, 173) 

        # Get plot 
        plot = chart.plots[0]
        # Set for each bar same color
        plot.vary_by_categories = False
        # Show data labels and set font
        plot.has_data_labels = True
        # Change gap width
        plot.gap_width = 100

        
        data_labels = plot.data_labels
        data_labels.font.size = self.data_label_font_size
        data_labels.font.bold = True
        data_labels.font.name = self.font_name

        if 'Total Patients' in column_name or 'Median patient age' in column_name:
            value_axis = chart.value_axis
            value_axis.visible = False
            value_axis.has_major_gridlines = False
        else:
            # Value for x-axis (change font size, name, and other things)
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = self.category_font_size
            tick_labels.font.name = self.font_name
            
            # Don't show major gridlines
            value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
            value_axis.has_major_gridlines = False
            # Set range of axis
            value_axis.maximum_scale = maximum
            value_axis.minimum_scale = 0
       
        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
        solidFill = category_axis.format.line.color._xFill
        self._set_transparency(100, solidFill)

        # Delete tick marks
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.major_unit = 1
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name

    def _create_stacked_barplot(self, dataframe, title, column_name, legend, number_of_series):
        """ The function creating the normal barplot graph into the presentation based on the graph type. 
        
        :param dataframe: the dataframe with data to be shown
        :type dataframe: pandas dataframe
        :param title: the title of the graph
        :type title: str
        :param column_name: the first column to be displayed in the graph
        :type column_name: str
        :param legend: the list of legend names
        :type legend: list
        :param number_of_series: the number of series to be shown in the graph
        :type number_of_series: int
        """

        colors = {
            0: RGBColor(43, 88, 173), # dark blue
            1: RGBColor(237, 125, 49), # orange
            2: RGBColor(165, 165, 165), # gray
            3: RGBColor(255, 192, 0), # yellow
            #4: RGBColor(59, 100, 173), # blue
            4: RGBColor(136, 106, 159),
            5: RGBColor(98, 153, 62), # green
            6: RGBColor(151, 185, 224), # light blue
            7: RGBColor(241, 167, 138), # beige       
        }

        # Calculate length of legend (in case that legend is too long, make smaller font size)
        count = self._get_length_of_legend(legend)

        # Get column names of dataframe
        column_names = dataframe.columns.tolist()

        index = column_names.index(column_name)

        # Add new slide into presentation
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        title_placeholders = slide.shapes.title
        title_placeholders.text = title


        chart_data = ChartData()
        
            
        chart_data.categories = dataframe[self.categories_column].tolist()
        # Add data in each category 
        chart_data.add_series(legend[0], dataframe[column_names[index]].tolist()) 
        if (number_of_series >= 2):
            chart_data.add_series(legend[1], dataframe[column_names[index+1]].tolist())
        if (number_of_series >= 3):
            chart_data.add_series(legend[2], dataframe[column_names[index+2]].tolist())
        if (number_of_series >= 4):
            chart_data.add_series(legend[3], dataframe[column_names[index+3]].tolist())
        if (number_of_series >= 5):
            chart_data.add_series(legend[4], dataframe[column_names[index+4]].tolist())
        if (number_of_series >= 6):
            chart_data.add_series(legend[5], dataframe[column_names[index+5]].tolist())
        if (number_of_series >= 7):
            chart_data.add_series(legend[6], dataframe[column_names[index+6]].tolist())
        if (number_of_series >= 8):
            chart_data.add_series(legend[7], dataframe[column_names[index+7]].tolist())

        # add chart to slide --------------------
        specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
            }
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        # Change gap width
        plot = chart.plots[0]
        plot.gap_width = 100


        site_names = dataframe[self.categories_column].tolist()
       
        """
        if len(dataframe) <= 2:
            series = chart.series[0]
                
            # Set color of first series to dark blue
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = colors[0] 
        """

        for i in range(0, number_of_series):
            series = chart.series[i]
            
            # Set color of first series to dark blue
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = colors[i]  

            
            if  len(dataframe) > 2:
                # ---add an `a:alpha` child element---
                solidFill = fill.fore_color._xFill
                self._set_transparency(30, solidFill)

                # Change color of borders of series and transparency
                series.format.line.color.rgb = colors[i]
                solidFill = series.format.line.color._xFill
                self._set_transparency(70, solidFill)

                # Remove transparency from country point
                for idx, point in enumerate(series.points):
                    if (site_names[idx] == self.country_name):
                        point.format.line.color.rgb = colors[i]
                        # Get fill of point for country
                        fill = point.format.fill
                        fill.solid()
                        fill.fore_color.rgb = colors[i]
          
        # Value for x-axis (change font size, name, and other things)
        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = Pt(11)
        tick_labels.font.name = self.font_name

        value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
        
        if (len(dataframe) > 2):
            value_axis.has_major_gridlines = True
            #value_axis.major_gridlines.format.line.dash_style = MSO_LINE.DASH
            value_axis.major_gridlines.format.line.width = Pt(0.5)
            value_axis.major_gridlines.format.line.color.rgb = RGBColor(166, 166, 166) # Set color to gray (A6A6A6)

            # Set 100% transparency to value axis
            value_axis.format.line.color.rgb = RGBColor(0, 0, 0)
            solidFill = value_axis.format.line.color._xFill
            self._set_transparency(100, solidFill)
        else:
            value_axis.has_major_gridlines = False

        value_axis.maximum_scale = 100
        value_axis.minimum_scale = 0

        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        
        # Set 100% transparency to category axis
        category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
        solidFill = category_axis.format.line.color._xFill
        self._set_transparency(100, solidFill)

        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name
        category_labels.tickLblSkip = 1

        # Set legend 
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.name = self.font_name

        if (count > 180 or 'antithrombotics prescribed' in title.lower()):
            chart.legend.font.size = Pt(11)
        else:
            chart.legend.font.size = Pt(12)

class GenerateGraphsSites(GenerateGraphs):

    def __init__(self):
        super(GenerateGraphsSites, self).__init__(*args, **kwargs)

        if 'data' in kwargs.keys():
            self.data = kwargs['data']

    def __get_specs(self, ngraphs=1):
        """ Get specification for graphs (the position in the pptx) based on number of graphs placed on slide. """
        main_height = Cm(16.5)
        if ngraphs == 1:
            specs = {
                'height': main_height,
                'width': Cm(32),
                'left': Cm(0.7),
                'top': Cm(2)
            }
        elif ngraphs == 2:
            specs = OrderedDict()
            height = main_height
            width = Cm(15.26)
            top = Cm(2)
            left = {0: Cm(0.5), 1: Cm(17.5)}
            for i in range(0, ngraphs):
                specs.update([(i, {
                    'height': height,
                    'width': width,
                    'top': top,
                    'left': left[i]
                })])
        elif ngraphs == 3:
            specs = OrderedDict()
            height = {0: Cm(16.5), 1: Cm(8.25), 2: Cm(8.25)}
            width = Cm(15.26)
            left = {0: Cm(0.5), 1: Cm(17.5), 2: Cm(17.5)}
            top = {0: Cm(2), 1: Cm(2), 2: Cm(10.25)}
            for i in range(0, ngraphs):
                specs.update([(i, {
                    'height': height,
                    'width': width,
                    'top': top,
                    'left': left[i]
                })])
        elif ngraphs == 4:
            specs = OrderedDict()
            height = Cm(8.25)
            width = Cm(15.26)
            left = {0: Cm(0.5), 1: Cm(0.5), 2: Cm(17.5), 3: Cm(17.5)}
            top = {0: Cm(2), 1: Cm(10.25), 2: Cm(2), 3: Cm(10.25)}
            for i in range(0, ngraphs):
                specs.update([(i, {
                    'height': height,
                    'width': width,
                    'top': top,
                    'left': left[i]
                })])
        
        return specs      

    def _create_plot(self, df, title, specs, graph_type, legend=None, ix=0):   
        """ The function creating the new graph into the presentation based on the graph type. 
        
        :param df: the dataframe with data to be shown
        :type df: pandas dataframe
        :param title: the title of the graph
        :type title: str
        :param specs: the position settings
        :type specs: dictionary
        :param graph_type: the type of graph (normal or stacked)
        :type graph_type: str
        :param legend: the list of values in legend based on columns (only for stacked barplot)
        :type legend: list
        :param ix: the index which legend should be used 
        :type ix: int
        """

        if graph_type == "normal":
            # Get column names of dataframe
            column_names = df.columns.tolist()
            index = column_names.index(self.categories_column)   

            # 1st chart (left side) - nationally sample
            chart_data = ChartData()
            chart_data.categories = df[self.categories_column].tolist()
            chart_data.add_series(column_names[index+1], df[column_names[index+1]].tolist())

            chart = self.slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, specs['left'], specs['top'], specs['width'], specs['height'], chart_data).chart

            # Get series of chart
            series = chart.series[0]
            series.points[0].format.fill.solid()
            series.points[0].format.fill.fore_color.rgb = self.colors['blue']
            series.points[1].format.fill.solid()
            series.points[1].format.fill.fore_color.rgb = self.colors['orange']
            series.points[2].format.fill.solid()
            series.points[2].format.fill.fore_color.rgb = self.colors['green']
            
            # Get plot 
            plot = chart.plots[0]
            # Set for each bar same color
            plot.vary_by_categories = True
            # Show data labels and set font
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = self.data_label_font_size
            data_labels.font.bold = True
            data_labels.font.name = self.font_name

            chart_text = chart.chart_title.text_frame
            chart_text.text = title
            chart_text.paragraphs[0].font.name = self.font_name
            chart_text.paragraphs[0].font.size = Pt(14)
            chart_text.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

            # Value for x-axis (change font size, name, and other things)
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = self.category_font_size
            tick_labels.font.name = self.font_name

            # Don't show major gridlines
            #value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
            value_axis.has_major_gridlines = True
            value_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
            value_axis.major_gridlines.format.line.width = Pt(0.5)
            # Set range of axis
            value_axis.minimum_scale = 0
            values = df[column_names[index+1]].tolist()
            max_value = max(values)
            if '%' in title and max_value >= 90:
                value_axis.maximum_scale = 100
            else:
                value_axis.maximum_scale = math.ceil(max_value / 10.0) * 10
            value_axis.visible = False

            # Value for y-axis (change font size, name, and other things)
            category_axis = chart.category_axis
            # Delete tick marks
            category_axis.major_tick_mark = XL_TICK_MARK.NONE
            category_axis.major_unit = 1
            category_labels = category_axis.tick_labels
            category_labels.font.size = self.category_font_size
            category_labels.font.name = self.font_name

        # Create stacked barplot
        else: 
            # Get column names of dataframe
            column_names = df.columns.tolist()
            index = column_names.index(self.categories_column) + 1

            legend = self.legends[ix]
            number_of_series = len(legend)

            # 1st dataframe (nationally sample)
            chart_data = ChartData()
            chart_data.categories = df[self.categories_column].tolist()
            # Add data in each category
            chart_data.add_series(legend[0], df[column_names[index]].tolist())
            if (number_of_series >= 2):
                chart_data.add_series(legend[1], df[column_names[index+1]].tolist())
            if (number_of_series >= 3):
                chart_data.add_series(legend[2], df[column_names[index+2]].tolist())
            if (number_of_series >= 4):
                chart_data.add_series(legend[3], df[column_names[index+3]].tolist())
            if (number_of_series >= 5):
                chart_data.add_series(legend[4], df[column_names[index+4]].tolist())
            if (number_of_series >= 6):
                chart_data.add_series(legend[5], df[column_names[index+5]].tolist())
            if (number_of_series >= 7):
                chart_data.add_series(legend[6], df[column_names[index+6]].tolist())
            if (number_of_series >= 8):
                chart_data.add_series(legend[7], df[column_names[index+7]].tolist())

            chart = self.slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_STACKED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

            series = chart.series[0]
            # Get series of chart
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = self.colors['blue']
            if (number_of_series >= 2):
                chart.series[1].format.fill.solid()
                chart.series[1].format.fill.fore_color.rgb = self.colors['orange']
            if (number_of_series >= 3):
                chart.series[2].format.fill.solid()
                chart.series[2].format.fill.fore_color.rgb = self.colors['green']
            if (number_of_series >= 4):
                chart.series[3].format.fill.solid()
                chart.series[3].format.fill.fore_color.rgb = self.colors['grey']
            if (number_of_series >= 5):
                chart.series[4].format.fill.solid()
                chart.series[4].format.fill.fore_color.rgb = self.colors['violet']
            if (number_of_series >= 6):
                chart.series[5].format.fill.solid()
                chart.series[5].format.fill.fore_color.rgb = self.colors['yellow']

            # Value for x-axis (change font size, name, and other things)
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = Pt(11)
            tick_labels.font.name = self.font_name

            value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE

            # Don't show major gridlines
            #value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
            value_axis.has_major_gridlines = True
            value_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
            value_axis.major_gridlines.format.line.width = Pt(0.5)
            # Set range of axis
            value_axis.minimum_scale = 0
            value_axis.maximum_scale = 100
            value_axis.visible = True

             # Set 100% transparency to value axis
            value_axis.format.line.color.rgb = RGBColor(0, 0, 0)
            solidFill = value_axis.format.line.color._xFill
            self._set_transparency(100, solidFill)


            category_axis = chart.category_axis
            category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
            solidFill = category_axis.format.line.color._xFill
            self._set_transparency(100, solidFill)
            
            category_axis.major_tick_mark = XL_TICK_MARK.NONE
            category_labels = category_axis.tick_labels
            category_labels.font.size = self.category_font_size
            category_labels.font.name = self.font_name
            category_labels.tickLblSkip = 1
        
            # Set legend 
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.name = self.font_name
            chart.legend.font.size = Pt(14)

            chart_text = chart.chart_title.text_frame
            chart_text.text = title
            chart_text.paragraphs[0].font.size = Pt(18)
            chart_text.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

    def _create_barplot(self, dataframe, title, column_name):
        """ The function creating the normal barplot graph into the presentation based on the graph type. 
        
        :param df: the dataframe with data to be shown
        :type df: pandas dataframe
        :param title: the title of the graph
        :type title: str
        :param column_name: the column name to be displayed in the graph
        :type column_name: str
        """
        maximum = 0

        # If graph is in %, set maximum valut to 100. 
        if '%' in title.lower():
            maximum = 100
            if self.country_name == 'Czech Republic':
                values = [round(x, 0) for x in dataframe[column_name].tolist()]
        else:
            maximum = round((max(dataframe[column_name].tolist())),1)
            if self.country_name == 'Czech Republic':
                values = dataframe[column_name].tolist()

        # Add slide to presentation (layout 11 is our custom layout where only title 'Agency FB', color: RGBColor(43, 88, 173)  and size:24 is set)
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        # Get title object
        title_placeholders = slide.shapes.title
        # Set title
        title_placeholders.text = title

        chart_data = ChartData()
        chart_data.categories = dataframe[self.categories_column].tolist()
        chart_data.add_series(column_name, dataframe[column_name].tolist())

        # Add chart on slide
        specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
            }
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        # Get series of chart
        series = chart.series[0]
        # If graphs for whole country are generated, set for bar with country with red color
        # else set to blue color (same color as title uses)
        site_names = dataframe[self.categories_column].tolist()
        if (len(dataframe) > 2):
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                if (site_names[idx] == self.country_name):
                    fill.fore_color.rgb = RGBColor(128,0,0)
                else:
                    fill.fore_color.rgb = RGBColor(43, 88, 173)
        else:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(43, 88, 173) 

        # Get plot 
        plot = chart.plots[0]
        # Set for each bar same color
        plot.vary_by_categories = False
        # Show data labels and set font
        plot.has_data_labels = True
        # Change gap width
        plot.gap_width = 100

        
        data_labels = plot.data_labels
        data_labels.font.size = self.data_label_font_size
        data_labels.font.bold = True
        data_labels.font.name = self.font_name

        if 'Total Patients' in column_name or 'Median patient age' in column_name:
            value_axis = chart.value_axis
            value_axis.visible = False
            value_axis.has_major_gridlines = False
        else:
            # Value for x-axis (change font size, name, and other things)
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = self.category_font_size
            tick_labels.font.name = self.font_name
            
            # Don't show major gridlines
            value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
            value_axis.has_major_gridlines = False
            # Set range of axis
            value_axis.maximum_scale = maximum
            value_axis.minimum_scale = 0
       
        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
        solidFill = category_axis.format.line.color._xFill
        self._set_transparency(100, solidFill)

        # Delete tick marks
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.major_unit = 1
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name

    def _create_stacked_barplot(self, dataframe, title, column_name, legend, number_of_series):
        """ The function creating the normal barplot graph into the presentation based on the graph type. 
        
        :param dataframe: the dataframe with data to be shown
        :type dataframe: pandas dataframe
        :param title: the title of the graph
        :type title: str
        :param column_name: the first column to be displayed in the graph
        :type column_name: str
        :param legend: the list of legend names
        :type legend: list
        :param number_of_series: the number of series to be shown in the graph
        :type number_of_series: int
        """

        colors = {
            0: RGBColor(43, 88, 173), # dark blue
            1: RGBColor(237, 125, 49), # orange
            2: RGBColor(165, 165, 165), # gray
            3: RGBColor(255, 192, 0), # yellow
            #4: RGBColor(59, 100, 173), # blue
            4: RGBColor(136, 106, 159),
            5: RGBColor(98, 153, 62), # green
            6: RGBColor(151, 185, 224), # light blue
            7: RGBColor(241, 167, 138), # beige       
        }

        # Calculate length of legend (in case that legend is too long, make smaller font size)
        count = self._get_length_of_legend(legend)

        # Get column names of dataframe
        column_names = dataframe.columns.tolist()

        index = column_names.index(column_name)

        # Add new slide into presentation
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        title_placeholders = slide.shapes.title
        title_placeholders.text = title


        chart_data = ChartData()
        
            
        chart_data.categories = dataframe[self.categories_column].tolist()
        # Add data in each category 
        chart_data.add_series(legend[0], dataframe[column_names[index]].tolist()) 
        if (number_of_series >= 2):
            chart_data.add_series(legend[1], dataframe[column_names[index+1]].tolist())
        if (number_of_series >= 3):
            chart_data.add_series(legend[2], dataframe[column_names[index+2]].tolist())
        if (number_of_series >= 4):
            chart_data.add_series(legend[3], dataframe[column_names[index+3]].tolist())
        if (number_of_series >= 5):
            chart_data.add_series(legend[4], dataframe[column_names[index+4]].tolist())
        if (number_of_series >= 6):
            chart_data.add_series(legend[5], dataframe[column_names[index+5]].tolist())
        if (number_of_series >= 7):
            chart_data.add_series(legend[6], dataframe[column_names[index+6]].tolist())
        if (number_of_series >= 8):
            chart_data.add_series(legend[7], dataframe[column_names[index+7]].tolist())

        # add chart to slide --------------------
        specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
            }
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        # Change gap width
        plot = chart.plots[0]
        plot.gap_width = 100


        site_names = dataframe[self.categories_column].tolist()
       
        """
        if len(dataframe) <= 2:
            series = chart.series[0]
                
            # Set color of first series to dark blue
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = colors[0] 
        """

        for i in range(0, number_of_series):
            series = chart.series[i]
            
            # Set color of first series to dark blue
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = colors[i]  

            
            if  len(dataframe) > 2:
                # ---add an `a:alpha` child element---
                solidFill = fill.fore_color._xFill
                self._set_transparency(30, solidFill)

                # Change color of borders of series and transparency
                series.format.line.color.rgb = colors[i]
                solidFill = series.format.line.color._xFill
                self._set_transparency(70, solidFill)

                # Remove transparency from country point
                for idx, point in enumerate(series.points):
                    if (site_names[idx] == self.country_name):
                        point.format.line.color.rgb = colors[i]
                        # Get fill of point for country
                        fill = point.format.fill
                        fill.solid()
                        fill.fore_color.rgb = colors[i]
          
        # Value for x-axis (change font size, name, and other things)
        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = Pt(11)
        tick_labels.font.name = self.font_name

        value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
        
        if (len(dataframe) > 2):
            value_axis.has_major_gridlines = True
            #value_axis.major_gridlines.format.line.dash_style = MSO_LINE.DASH
            value_axis.major_gridlines.format.line.width = Pt(0.5)
            value_axis.major_gridlines.format.line.color.rgb = RGBColor(166, 166, 166) # Set color to gray (A6A6A6)

            # Set 100% transparency to value axis
            value_axis.format.line.color.rgb = RGBColor(0, 0, 0)
            solidFill = value_axis.format.line.color._xFill
            self._set_transparency(100, solidFill)
        else:
            value_axis.has_major_gridlines = False

        value_axis.maximum_scale = 100
        value_axis.minimum_scale = 0

        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        
        # Set 100% transparency to category axis
        category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
        solidFill = category_axis.format.line.color._xFill
        self._set_transparency(100, solidFill)

        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name
        category_labels.tickLblSkip = 1

        # Set legend 
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.name = self.font_name

        if (count > 180 or 'antithrombotics prescribed' in title.lower()):
            chart.legend.font.size = Pt(11)
        else:
            chart.legend.font.size = Pt(12)

    



class GenerateGraphsQuantiles:
    """ The class generating presentation with graphs for general reports.

    :param dataframe: the dataframe with calculated statistics
    :type dataframe: pandas dataframe
    :param presentation: the opened presentation document
    :type presentation: Presentation object
    :param title: the title of the slide
    :type title: str
    :param column_name: the column name from the dataframe to be shown in the graph
    :type column_name: str
    :param graph_type: the type of graph to be generated
    :type graph_type: str
    :param number_of_series: the number of columns to be shown in the stacked barplot
    :type number_of_series: int
    :param legend: the list of names to be used in the legend in the stacked barplot
    :type legend: list
    :param country: the country name used in the first slide
    :type country: str
    """

    def __init__(self, dataframe, presentation, title, column_name, graph_type = None, number_of_series=0, legend=None, country=None):

        self.dataframe = dataframe
        self.presentation = presentation
        self.title = title
        self.column_name = column_name
        self.number_of_series = number_of_series
        self.legend = legend
        self.country_name = country
        self.font_name = 'Century Gothic'

        # If country name is Ukraine or Poland set as categories value Site IDs not Site Name due to huge amount of sites in graph
        # if (self.country_name in ['Ukraine', 'Poland'] and len(self.dataframe) > 2):
        #     self.categories_column = 'Site ID'
        # else:
        self.categories_column = 'Site Name'
            
        # Estimate font sizes based on number of sites included in the graph
        if (len(self.dataframe) > 60):
            self.category_font_size = Pt(4)
            self.data_label_font_size = Pt(4)
        elif (len(self.dataframe) > 50 and len(self.dataframe) <= 60):
            self.category_font_size = Pt(6)
            self.data_label_font_size = Pt(6)
        else:
            self.category_font_size = Pt(8)
            self.data_label_font_size = Pt(8)

        # Select graph which should be exported
        self._create_barplot(dataframe=self.dataframe, title=self.title, column_name=self.column_name)

    def _get_length_of_legend(self, legend):
        """ The function adjusting the number of letters in legend to quess the number of columns in the legend! 
        
        :param legend: the names of legend
        :type legend: list
        :returns: the adjusted number of letters
        """
        count = 0

        for i in legend:
            count = count + len(i)
        
        return count


    def _set_transparency(self, transparency, elm):
        """ The function set the transparency of the row. 

        :param transparency: the transparency in %
        :type transparency: int
        :param elm: the element which transparency should be changed
        :type elm: format.line.color._xFill
        """
        a = str(100 - transparency) + '196'
        
        alpha = OxmlElement('a:alpha')
        alpha.set('val', a)
        elm.srgbClr.append(alpha)


    def _create_barplot(self, dataframe, title, column_name):
        """ The function creating the normal barplot graph into the presentation based on the graph type. 
        
        :param df: the dataframe with data to be shown
        :type df: pandas dataframe
        :param title: the title of the graph
        :type title: str
        :param column_name: the column name to be displayed in the graph
        :type column_name: str
        """
        maximum = 0

        # If graph is in %, set maximum valut to 100. 
        if '%' in title.lower():
            maximum = 100
            if self.country_name == 'Czech Republic':
                values = [round(x, 0) for x in dataframe[column_name].tolist()]
        else:
            maximum = round((max(dataframe[column_name].tolist())),1)
            if self.country_name == 'Czech Republic':
                values = dataframe[column_name].tolist()

        # Add slide to presentation (layout 11 is our custom layout where only title 'Agency FB', color: RGBColor(43, 88, 173)  and size:24 is set)
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        # Get title object
        title_placeholders = slide.shapes.title
        # Set title
        title_placeholders.text = title

        chart_data = ChartData()
        chart_data.categories = dataframe[self.categories_column].tolist()
        chart_data.add_series(column_name, dataframe[column_name].tolist())

        # Add chart on slide
        specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
            }
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        # Get series of chart
        series = chart.series[0]
        # If graphs for whole country are generated, set for bar with country with red color
        # else set to blue color (same color as title uses)
        site_names = dataframe[self.categories_column].tolist()
        if (len(dataframe) > 2):
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                if (site_names[idx] == self.country_name):
                    fill.fore_color.rgb = RGBColor(128,0,0)
                elif (site_names[idx] == 'Q1' or site_names[idx] == 'Q3'):
                    fill.fore_color.rgb = RGBColor(84,130,53)
                else:
                    fill.fore_color.rgb = RGBColor(43, 88, 173)
        else:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(43, 88, 173) 

        # Get plot 
        plot = chart.plots[0]
        # Set for each bar same color
        plot.vary_by_categories = False
        # Show data labels and set font
        plot.has_data_labels = True
        # Change gap width
        plot.gap_width = 100

        
        data_labels = plot.data_labels
        data_labels.font.size = self.data_label_font_size
        data_labels.font.bold = True
        data_labels.font.name = self.font_name

        if 'Total Patients' in column_name or 'Median patient age' in column_name:
            value_axis = chart.value_axis
            value_axis.visible = False
            value_axis.has_major_gridlines = False
        else:
            # Value for x-axis (change font size, name, and other things)
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = self.category_font_size
            tick_labels.font.name = self.font_name
            
            # Don't show major gridlines
            value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
            value_axis.has_major_gridlines = False
            # Set range of axis
            value_axis.maximum_scale = maximum
            value_axis.minimum_scale = 0
       
        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
        solidFill = category_axis.format.line.color._xFill
        self._set_transparency(100, solidFill)

        # Delete tick marks
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.major_unit = 1
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name

