# -*- coding: utf-8 -*-
"""
File name: GenerateNationalComparisonGraphs.py
Package: resq
Written by: Marie Jankujova - jankujova.marie@fnusa.cz on 02-2019
Version: v1.0
Version comment: The first version of script which create national comaprison reports. 
Description: This script is used to call class GenerateGraphs and produce graphs into presentation. 
"""


import pandas as pd
import numpy as np
import sys
import os
from datetime import datetime, date
import sqlite3

from resqdb.GenerateGraphs import GenerateGraphs
from resqdb.GenerateComparisonPresentation import GenerateYearsCompGraphs

import xlsxwriter

from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor

import pytz


class GenerateNationalComparisonGraphs:
    """ The class generating the presentation with graphs for national comparison reports. 

    :param df: the dataframe with calculated statistic for the period of time
    :type df: pandas dataframe
    :param fdf: the dataframe with calculated comparison statistics eg. 2016, 2017, 2018, etc.
    :type fdf: pandas dataframe
    :param outcome: the dataframe with outcome results (default: None)
    :type outcome: pandas dataframe
    :param country: `True` if country should be included in the results as site
    :type country: bool
    :param country_code: the country code (default: None)
    :type country_code: str
    :param split_sites: `True` if graphs for each site should be generated seperately
    :type split_sites: bool
    :param site: the site ID
    :type site: str
    :param report: the type of the report eg. quarter
    :type report: str
    :param quarter: the type of the period eg. Q1_2019
    :type quarter: str
    """

    def __init__(self, df, fdf, outcome=None, country=False, country_code=None, split_sites=False, site=None, report=None, quarter=None):

        self.df = df
        self.fdf = fdf
        self.country_code = country_code
        self.report = report
        self.quarter = quarter
        self.outcome = outcome

        # Get absolute path to the database.
        script_dir = os.path.dirname(__file__)

        master_pptx = "master.pptx"
        self.master = os.path.normpath(os.path.join(script_dir, "backgrounds", master_pptx))

        # Connect to database and get country name according to country code.
        def select_country(value):
            """ The function obtaining the name of country from the package pytz based on the country code. 

            :param value: the country code
            :type value: str
            :returns: the country name
            """
            country_name = pytz.country_names[value]
            return country_name

        # If country is used as site, the country name is selected from countries dictionary by country code. :) 
        if (country):
            self.country_name = select_country(self.country_code)
        else:
            self.country_name = None

        # If split_sites is True, then go through dataframe and generate graphs for each site (the country will be included as site in each file).
        site_ids = self.df['Site ID'].tolist()
        # Delete country name from side ids list.
        try:
            site_ids.remove(self.country_name)
        except:
            pass

        
        if site is not None:
            df = self.df[self.df['Site ID'].isin([site, self.country_name])].copy()
            self._generate_graphs(df=df, site_code=site)

        # Generate formatted statistics for all sites individualy + country as site is included
        if (split_sites) and site is None:
            for i in site_ids:
                df = self.df[self.df['Site ID'].isin([i, self.country_name])].copy()
                self._generate_graphs(df=df, site_code=i)
    
        # Produce formatted statistics for all sites + country as site
        if site is None:
            if outcome is not None:
                self._generate_graphs(df=self.df, fdf=self.fdf, outcome=self.outcome, site_code=country_code)
            else:
                self._generate_graphs(df=self.df, fdf=self.fdf, site_code=country_code)

    def _generate_graphs(self, df, fdf, outcome=None, site_code=None):
        """ The function generating the graph in the presentation. 
        
        :param df: the dataframe containing the general statistics for country
        :type df: pandas dataframe
        :param fdf: the dataframe containing the general statistics for country through period in each year
        :type fdf: pandas dataframe
        :param outcome: the outcome dataframe containing outcome information for the given period
        :type outcome: pandas dataframe
        :param site_code: the site ID
        :type site_code: str
        """
        
        prs = Presentation(self.master)

        first_slide = prs.slides[0]
        shape = first_slide.shapes[5]
        text_frame = shape.text_frame

        #first_slide_text = self.country_name + "\nData Summary"
        if site_code == 'CZ':   
            first_slide_text = "Přehled dat\nČeská republika\n2019"
        else:
            first_slide_text = "\nData Summary"

        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = first_slide_text

        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(26)
        font.color.rgb = RGBColor(250,250,250)

        # if (self.country_name in ['Ukraine', 'Poland'] and len(df) > 2):
        #     main_col = 'Site ID'
        # else:
        main_col = 'Site Name'

        years = ', '.join(map(str, self.fdf[main_col].tolist()))

        if site_code == "CZ":
            comp_title = "Dočasné trendy - {}".format(years)
        else:
            comp_title = "Temporal trends - {}".format(years)

        ########################
        #### TOTAL PATIENTS ####
        ########################
        column_name = 'Total Patients'
        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        country_patients = str(max(tmp_df[column_name].tolist()))
        if site_code == 'CZ':
            title = 'Celkový počet pacientů (n = {})'.format(country_patients)
        else:
            title = 'TOTAL PATIENTS (n = {})'.format(country_patients)
        if self.country_name is not None:
            tmp_df = tmp_df.loc[tmp_df[main_col] != self.country_name]

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##########################
        ### MEDIAN PATIENT AGE ###
        ##########################
        column_name = 'Median patient age'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            title = "Věk pacientů - medián"
        else:
            title = "MEDIAN PATIENT AGE"
        

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###################
        ## OUTCOME table ##
        ###################
        if outcome is not None:
            slide = prs.slides.add_slide(prs.slide_layouts[11])
            title_placeholders = slide.shapes.title
            title_placeholders.text = "VÝSLEDKY PODLE SKUPIN PACIENTŮ"

            # table_placeholder = slide.shapes[1]
            x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
            shape = slide.shapes.add_table(7, 3, x, y, cx, cy)
            table = shape.table

            # table = shape.table
            tmp_df = outcome[['Patient Group', 'n', 'Median discharge mRS']]


            columns = ['Patient Group', 'n', 'Median discharge mRS']
            for i in range(0, len(columns)):
                cell = table.cell(0, i)
                cell.text = columns[i]
            for index, row in tmp_df.iterrows():
                for i in range(0, len(row)):
                    cell = table.cell(index+1, i)
                    cell.text = str(row[i])


        #######################################
        #### Outcome for group of patients ####
        #######################################
        if outcome is not None:
            titles = []
            graph_types = []
            legends = []
            
            if site_code == 'CZ':
                #title = "VÝSLEDKY PODLE SKUPIN TYPU LÉČBY"
                title = "VÝSLEDKY PODLE SKUPIN TYPU PACIENTŮ"
                legend = ['domů', 'překlad v rámci stejného zdravotnického zařízení (ZZ)', 'překlad do jiného zdravotnického zařízení (ZZ)', 'zařízení sociální péče', 'zemřel/a']
            else:
                title = "Results for outcome groups"
                legend = ['% home', '% transferred within the same centre', '% transferred to another centre', '% social care facility', '% dead']
            legends.append(legend)
            graph_types.append("stacked")

            """
            tmp = outcome.loc[outcome['Patient Group'].isin(['IV tPA', 'IV tPA + TBY', 'TBY'])]
            df1 = tmp[['Patient Group', '% home', '% transferred within the same centre', '% transferred to another centre', '% social care facility', '% dead']]

            GenerateYearsCompGraphs(df=df1, presentation=prs, title=title, titles=[""], graph_types=graph_types, legends=legends, outcome=True)

            titles = []
            graph_types = []
            legends = []
            
            
            if site_code == 'CZ':
                title = "VÝSLEDKY PODLE SKUPIN TYPU CMP"
                legend = ['domů', 'překlad v rámci stejného zdravotnického zařízení (ZZ)', 'překlad do jiného zdravotnického zařízení (ZZ)', 'zařízení sociální péče', 'zemřel/a']
            else:
                title = "Results for outcome gruops"
                legend = ['% home', '% transferred within the same centre', '% transferred to another centre', '% social care facility', '% dead']
            legends.append(legend)
            graph_types.append("stacked")

            tmp = outcome.loc[outcome['Patient Group'].isin(['iCMP', 'ICH', 'SAK'])]
            df1 = tmp[['Patient Group', '% home', '% transferred within the same centre', '% transferred to another centre', '% social care facility', '% dead']]

            GenerateYearsCompGraphs(df=df1, presentation=prs, title=title, titles=[""], graph_types=graph_types, legends=legends, outcome=True)
            """
            outcome_df = outcome[['Patient Group', '% home', '% transferred within the same centre', '% transferred to another centre', '% social care facility', '% dead']]
            GenerateYearsCompGraphs(df=outcome_df, presentation=prs, title=title, titles=[""], graph_types=graph_types, legends=legends, outcome=True)


        ####################################
        #### Outcome per patients group ####
        ####################################
        
        #### Comparison graphs for trends per years - total patients and median patient age
        titles = []
        graph_types = []
        legend = []
        # Total patients and median age graphst)
        column_name = 'Total Patients'
        df1 = fdf[[main_col, column_name]]
        if site_code == 'CZ':
            titles.append("Celkem pacientů")
        else:
            titles.append("Total Patients")
        graph_types.append("normal")

        column_name = 'Median patient age'
        df2 = fdf[[main_col, column_name]]
        if site_code == 'CZ':
            titles.append("Věk pacientů - medián")
        else:
            titles.append("Median patient age")
        graph_types.append("normal")
        
        GenerateYearsCompGraphs(df=df1, df1=df2, presentation=prs, title=comp_title, titles=titles, graph_types=graph_types)

        ###########################
        ### GENDER DISTRIBUTION ###
        ###########################
        column_name = '% patients female'
        legend = ['Žena', 'Muž']

        tmp_df = df[[main_col, '% patients female', '% patients male']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)
        
        if site_code == 'CZ':
            title = "Rozdělení podle pohlaví"
        else:
            title = "GENDER DISTRIBUTION"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #######################
        ### DEPARTMENT TYPE ###
        #######################
        column_name = '% department type - neurology'
        
        #

        tmp_df = df[[main_col, '% department type - neurology', '% department type - neurosurgery', '% department type - anesthesiology/resuscitation/critical care', '% department type - internal medicine', '% department type - geriatrics', '% department type - Other']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            title = 'ODDĚLENÍ, KDE BYL PACIENT HOSPITALIZOVÁN (%)'
            legend = ['neurologie', 'neurochirurgie', 'ARO', 'interna', 'geriatrie', 'ostatní']
        else:
            title = "% DEPARTMENT TYPE ALLOCATION out of all cases" 
            legend = ['neurology', 'neurosurgery', 'anesthesiology resuscitation critical care', 'internal medicine', 'geriatrics', 'other']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')


        ###################################
        ### HOSPITALIZATION DESTINATION ###
        ###################################
        column_name = '% patients hospitalized in stroke unit / ICU'

        tmp_df = df[[main_col, '% patients hospitalized in stroke unit / ICU', '% patients hospitalized in monitored bed with telemetry', '% patients hospitalized in standard bed']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            title = "TYP LŮŽKA, NA KTERÉM BYL PACIENT HOSPITALIZOVÁN (%)"
            legend = ['JIP ', 'jiné monitorované lůžko (telemetrie)', 'standardní lůžko']
        else:
            title = "% HOSPITALIZATION DESTINATION out of all cases" 
            legend = ['stroke unit', 'monitored bed with telemetry', 'standard bed']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ################
        # COMPARISON ###
        ################
        legends = []
        titles = []
        graph_types = []
        ### DEPARTMENT TYPE ###
        if site_code == 'CZ':
            titles.append("Typ oddělení (%)")
            legend = ['neurologie', 'neurochirurgie', 'ARO', 'interna', 'geriatrie', 'ostatní']
        
        else:
            titles.append("Department type (%)")
            legend = ['neurology', 'neurosurgery', 'anesthesiology resuscitation critical care', 'internal medicine', 'geriatrics', 'other']
        legends.append(legend)
        graph_types.append("stacked")

        df1 = fdf[[main_col, '% department type - neurology', '% department type - neurosurgery', '% department type - anesthesiology/resuscitation/critical care', '% department type - internal medicine', '% department type - geriatrics', '% department type - Other']]

        ### HOSPITALIZATION DESTINATION ###
        if site_code == 'CZ':
            titles.append("Typ hospitalizace (%)")
            legend = ['JIP ', 'jiné monitorované lůžko (telemetrie)', 'standardní lůžko']
        else:
            titles.append("Hospitalization type (%)")
            legend = ['stroke unit', 'monitored bed with telemetry', 'standard bed']
        legends.append(legend)
        graph_types.append("stacked")

        df2 = fdf[[main_col, '% patients hospitalized in stroke unit / ICU', '% patients hospitalized in monitored bed with telemetry', '% patients hospitalized in standard bed']]
   
        '''
        ### REHABILIATION ###
        column_name = '% patients assessed for rehabilitation - Yes'
        df3 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == 'CZ':
            titles.append("Vyšetřena fyzioterapeutem (%)")
        else:
            titles.append("Assessed for rehabilitation (%)")
        graph_types.append("normal")
        '''

        GenerateYearsCompGraphs(df=df1, df1=df2, presentation=prs, title=comp_title, titles=titles, graph_types=graph_types, legends=legends)

        ###############
        # STROKE TYPE #
        ###############
        column_name = '% stroke type - ischemic stroke'

        tmp_df = df[[main_col, '% stroke type - ischemic stroke', '% stroke type - transient ischemic attack', '% stroke type - intracerebral hemorrhage', '% stroke type - subarrachnoid hemorrhage', '% stroke type - cerebral venous thrombosis', '% stroke type - undetermined stroke']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            title = "TYP CMP ze všech případů (%)"
            legend = ['ischemická CMP', 'tranzitorní ischemická ataka - TIA', 'intracerebrální krvácení', 'subarachnoidální krvácení', 'mozková žilní trombóza', 'neurčená']
        else:
            title = "% STROKE TYPE out of all cases"
            legend = ['ischemic', 'transient ischemic attack', 'intracerebral hemorrhage', 'subarrachnoid hemorrhage', 'cerebral venous thrombosis', 'undetermined']


        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ######################
        # CONSIOUSNESS LEVEL #
        ######################
        '''
        column_name = 'alert_all_perc'
        legend = ['alert', 'drowsy', 'comatose']

        tmp_df = df[[main_col, 'alert_all_perc', 'drowsy_all_perc', 'comatose_all_perc']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            title = "% STUPEŇ VĚDOMÍ u iCMP, ICH, CVT, SAK"
        else:
            title = "% CONSCIOUSNESS LEVEL for IS, ICH, CVT, SAH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')
        '''
        ###################
        # NIHSS PERFORMED #
        ###################
        column_name = '% NIHSS - Performed'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            #title = "% pacientů s iCMP, ICH nebo CVT, u nichž bylo provedeno NIHSS"
            title = "% pacientů s iCMP nebo ICH, u nichž bylo provedeno NIHSS"
        else:
            title = "% NIHSS PERFORMED for IS, ICH, CVT"
            
            

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ######################
        # NIHSS MEDIAN SCORE #
        ######################
        column_name = 'NIHSS median score'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        if site_code == 'CZ':
            title = "NIHSS SKÓRE - medián"
        else:
            title = "NIHSS median score"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ####################
        # CT/MRI performed #
        ####################
        column_name = '% CT/MRI - performed'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            title = "% pacientů s iCMP, ICH, CVT nebo TIA, kterým bylo provedeno CT/MRI"
        else:
            title = "% CT/MRI PERFORMED for IS, ICH, CVT, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################
        # CT/MRI PERFORMED WITHIN 1 HOUR #
        ##################################
        column_name = '% CT/MRI - Performed within 1 hour after admission'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            title = "% CT/MRI provedených DO HODINY OD PŘIJETÍ u iCMP, ICH, CVT nebo TIA "
        else:
            title = "% CT/MRI PERFORMED WITHIN 1 HOUR AFTER ADMISSION for IS, ICH, CVT, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ####################
        # VASCULAR IMAGING #
        ####################
        sorted_by = ['vascular_imaging_cta_norm', 'vascular_imaging_mra_norm', 'vascular_imaging_dsa_norm']
        column_name = 'vascular_imaging_cta_norm'
        #legend = ['CTA', 'MRA', 'DSA', 'none']
        legend = ['CTA', 'MRA', 'DSA']

        #tmp_df = df[[main_col, 'vascular_imaging_cta_norm', 'vascular_imaging_mra_norm', 'vascular_imaging_dsa_norm', 'vascular_imaging_none_norm']]
        tmp_df = df[[main_col, 'vascular_imaging_cta_norm', 'vascular_imaging_mra_norm', 'vascular_imaging_dsa_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == 'CZ':
            title = "% VASKULÁRNÍCH ZOBRAZOVÁNÍ provedených u ICH nebo SAK"
        else:
            title = "% VASCULAR IMAGING PERFORMED for ICH, SAH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ############################
        # RECANALIZATION TREATMENT #
        ############################
        column_name = '% recanalization procedures - IV tPa'
        
        tmp_df = df[[main_col, '% patients recanalized', '% recanalization procedures - IV tPa', '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone', '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment']]
        tmp_df = tmp_df.sort_values(['% patients recanalized'], ascending = True)

        if site_code == 'CZ':
            title = "% všech REKANALIZAČNÍCH VÝKONŮ u iCMP"
            legend = ['IV tPA – IC / KCC', 'IV tPA + endovaskulární výkon – KCC', 'pouze endovaskulární výkon – KCC', 'IV tPA + odeslán do jiného centra k endovaskulárnímu výkonu – IC']
        else:
            title = "% RECANALIZATION PROCEDURES for IS"
            legend = ['IV tPa', 'IV tPa + endovascular treatment', 'endovascular treatment', 'IV tPa + another centre for endovascular treatment']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ############################
        # RECANALIZATION TREATMENT #
        ############################
        column_name = '% recanalization procedures - IV tPa'
        

        tmp_df = df[[main_col, '% recanalization procedures - IV tPa', '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            title = "% PACIENTŮ LÉČENÝCH IV tPa u iCMP"
            legend = ['IV tPA – IC / KCC', 'IV tPA + endovaskulární výkon – KCC', 'IV tPA + odeslán do jiného centra k endovaskulárnímu výkonu – IC']
        else:
            title = "% IV tPa for IS"
            legend = ['IV tPA – IC / KCC', 'IV tPa + endovascular treatment', 'IV tPa + another centre for endovascular treatment']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ####################################################
        # RECANALIZATION TREATMENT IN COMPREHENSIVE CENTRE #
        ####################################################
        sorted_by = ['% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone']
        column_name = '% recanalization procedures - IV tPa + endovascular treatment'
        

        tmp_df = df[[main_col, '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == 'CZ':
            title = "% ENDOVASKULÁRNÍCH VÝKONŮ V KOMPLEXNÍCH CENTRECH u iCMP"
            legend = ['IV tPA + endovaskulární výkon – KCC', 'pouze endovaskulární výkon – KCC']
        else:
            title = "% RECANALIZATION PROCEDURES IN COMPREHENSIVE CENTRES for IS"
            legend = ['IV tPa + endovascular treatment', 'endovascular treatment']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################################################
        # RECANALIZATION TREATMENT IN REFFERED TO ANOTHER CENTRE #
        ##########################################################
        sorted_by = ['% recanalization procedures - IV tPa + referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment and hospitalization continues at the referred to centre', '% recanalization procedures - Referred for endovascular treatment and patient is returned to the initial centre']
        column_name = '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment'
        

        tmp_df = df[[main_col, '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment and hospitalization continues at the referred to centre', '% recanalization procedures - Referred for endovascular treatment and patient is returned to the initial centre']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == 'CZ':
            title = "% PACIENTŮ PŘELOŽENÝCH K ENDOVASKULÁRNÍMU VÝKONU DO JINÉHO CENTRA Z PRIMÁRNÍHO CENTRA u iCMP"
            legend = ['IV tPA + odeslán do jiného centra k endovaskulárnímu výkonu – IC', 'pouze odeslán do jiného centra k endovaskulárnímu výkonu – IC', 'pouze odeslán k endovaskulárnímu výkonu a hospitalizace pokračuje – KCC', 'pouze odeslán k endovaskulárnímu výkonu a pacient přeložen zpět do jiného centra – KCC']
        else:
            title = "% PATIENTS TRANSFERRED TO ANOTHER CENTRE FOR RECANALIZATION PROCEDURES FROM PRIMARY CENTRE for IS"
            legend = ['IV tPa + another centre for endovascular treatment', 'another centre for endovascular treatment', 'another centre for endovascular treatment and hospitalization continues', 'another centre for endovascular treatment and returned to the initial centre']
            

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #### Stroke types ####

        legends = []
        titles = []
        graph_types = []

        # STROKE TYPE #
        
        
        if site_code == 'CZ':
            titles.append("Typ CMP (%)")
            legend = ['ischemická CMP', 'tranzitorní ischemická ataka - TIA', 'intracerebrální krvácení', 'subarachnoidální krvácení', 'mozková žilní trombóza', 'neurčená']
        else:
            titles.append("Stroke type (%)")
            legend = ['ischemic', 'transient ischemic attack', 'intracerebral hemorrhage', 'subarrachnoid hemorrhage', 'cerebral venous thrombosis', 'undetermined']
        legends.append(legend)
        graph_types.append("stacked")

        df1 = fdf[[main_col, '% stroke type - ischemic stroke', '% stroke type - transient ischemic attack', '% stroke type - intracerebral hemorrhage', '% stroke type - subarrachnoid hemorrhage', '% stroke type - cerebral venous thrombosis', '% stroke type - undetermined stroke']]

        # CT/MRI performed #
        column_name = '% CT/MRI - performed'
        df2 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("CT/MRI (%)")
        graph_types.append("normal")

        # patients recanalized
        column_name = '% patients recanalized'
        df3 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == 'CZ':
            titles.append("% rekanalizovaných pacientů")
        else:
            titles.append("% patients recanalized")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df1, df1=df2, df2=df3, presentation=prs, title=comp_title, titles=titles, graph_types=graph_types, legends=legends)


        ################
        # % median DTN #
        ################
        column_name = 'Median DTN (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        if site_code == 'CZ':
            title = "MEDIÁN DOOR-TO-NEEDLE TIME (v minutách) u trombolyzovaných pacientů"
        else:
            title = "MEDIAN DOOR-TO-NEEDLE TIME (minutes) for thrombolyzed patients"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###############
        #  median DTG #
        ###############
        column_name = 'Median DTG (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        if site_code == 'CZ':
            title = "MEDIÁN DOOR-TO-GROIN TIME (v minutách) u pacientů, kteří podstoupili endovaskulární zákrok v komplexním centru"
        else:
            title = "MEDIAN DOOR-TO-GROIN TIME (minutes) for patients receiving endovascular treatment in a comprehensive centre"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ####################
        #  median TBY DIDO #
        ####################
        column_name = 'Median TBY DIDO (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        if site_code == 'CZ':
            title = "MEDIÁN DOOR-IN-DOOR-OUT TIME (v minutách) u pacientů transportovaných z primárního centra do jiného centra k rekanalizační terapii"
        else:
            title = "MEDIAN DOOR-IN-DOOR-OUT TIME (minutes) for patients referred from a primary centre to another centre for recanalization therapy"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)


        ##################
        ### COMPARISON ###
        ##################
        legends = []
        titles = []
        graph_types = []

        # RECANALIZATION PROCEDURES #
        if site_code == 'CZ':
            titles.append("Rekanalizační výkony (%)")
            legend =['IV tPA – IC / KCC', 'IV tPA + endovaskulární výkon – KCC', 'pouze endovaskulární výkon – KCC', 'IV tPA + odeslán do jiného centra k endovaskulárnímu výkonu – IC']
        else:
            titles.append("Recanalization procedures (%)")
            legend = ['IV tPa', 'IV tPa + endovascular treatment', 'endovascular treatment', 'IV tPa + another centre for endovascular treatment']

        legends.append(legend)
        
        df1 = fdf[[main_col, '% recanalization procedures - IV tPa', '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone', '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment']]
        graph_types.append("stacked")

        # MEDIAN DTN #
        column_name = 'Median DTN (minutes)'
        df2 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == 'CZ':
            #titles.append("Medián DTN (v minutách)")
            titles.append("Medián Door-To-Needle time (v minutách)")
        else:
            titles.append("Median DTN (minutes)")
        graph_types.append("normal")

        # MEDIAN DTN #
        column_name = 'Median DTG (minutes)'
        df3 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == 'CZ':
            #titles.append("Medián DTG (v minutách)")
            titles.append("Medián Door-To-Groin time (v minutách)")
        else:
            titles.append("Median DTG (minutes)")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df1, df1=df2, df2=df3, presentation=prs, title=comp_title, titles=titles, graph_types=graph_types, legends=legends)

        #######################
        # dysphagia screening #
        #######################
        column_name = '% dysphagia screening - Guss test'
        column_names = ['% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre']

        tmp_df = df[[main_col, '% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre']]
        tmp_df = tmp_df.sort_values(column_names, ascending = True)

        if site_code == 'CZ':
            #title = "% SCREENINGŮ DYSFAGIE provedených u iCMP, ICH nebo CVT"
            title = "% SCREENINGŮ DYSFAGIE provedených u iCMP nebo ICH"
            legend = ['GUSS test', 'Jiný test', 'Jiné pracoviště']
        else:
            title = "% DYSPHAGIA SCREENING PERFORMED for IS, ICH, CVT"
            legend = ['GUSS test', 'Other test', 'Another centre']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ######################################
        #  dypshagia screening within 1 hour #
        ######################################
        column_name = '% dysphagia screening time - Within first 24 hours'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            title = "% SCREENINGŮ DYSFAGIE BĚHEM 24 HODIN OD PŘIJETÍ"
        else:
            title = "% DYSPHAGIA SCREENING TIME WITHIN FIRST 24 HOURS AFTER ADMISSION"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        
        ###########################
        # PUT ON VENTILATOR - YES #
        ###########################
        column_name = '% patients put on ventilator - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            #title = "% PACIENTŮ NA UPV u iCMP, ICH nebo CVT"
            title = "% PACIENTŮ NA UPV u iCMP nebo ICH"
        else:
            title = "% PATIENTS PUT ON VENTILATOR for IS, ICH, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###################
        # HEMICRANEOCTOMY #
        ###################
        column_name = '% hemicraniectomy - Yes'
        

        tmp_df = df[[main_col, '% hemicraniectomy - Yes', '% hemicraniectomy - Referred to another centre']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% HEMIKRANIEKTOMIÍ provedených u iCMP"
            legend = ['Ano', 'Odeslán/a do jiného centra']
        else:
            title = "% HEMICRANIECTOMY PERFORMED for IS"
            legend = ['Yes', 'Referred to another centre']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        '''
        ################
        # NEUROSURGERY #
        ################
        column_name = '% neurosurgery - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% NEUROCHIRURGICKÝCH ZÁKROKŮ U PACIENTŮ S ICH"
        else:
            title = "% NEUROSURGERY PERFORMED for ICH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name) 

        #######################################
        # NEUROSURGERY TYPE (FIRST 3 OPTIONS) #
        #######################################
        sorted_by = ['% neurosurgery type - intracranial hematoma evacuation', '% neurosurgery type - external ventricular drainage', '% neurosurgery type - decompressive craniectomy']
        column_name = '% neurosurgery type - intracranial hematoma evacuation'
        legend = ['intracranial hematoma evacuation', 'external ventricular drainage', 'decompressive craniectomy']

        tmp_df = df[[main_col, '% neurosurgery type - intracranial hematoma evacuation', '% neurosurgery type - external ventricular drainage', '% neurosurgery type - decompressive craniectomy']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == "CZ":
            title = "% TYPŮ NEUROCHIRURGICKÉHO ZÁKROKU u PACIENTŮ s ICH v komplexních centrech"
        else:
            title = "% NEUROSURGERY TYPE PERFORMED for ICH in comprehensive centres"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        
        ###################################
        # NEUROSURGERY TYPE (LAST OPTION) #
        ###################################
        column_name = '% neurosurgery type - Referred to another centre'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% PACIENTŮ ODESLANÝCH DO JINÉHO CENTRA K NEUROCHIRURGICKÉMU ZÁKROKU"
        else:
            title = "% PATIENTS REFERRED TO ANOTHER CENTRE FOR NEUROSURGERY"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)'''

        ###################
        # BLEEDING REASON #
        ###################
        sorted_by = ['bleeding_arterial_hypertension_perc_norm', 'bleeding_aneurysm_perc_norm', 'bleeding_arterio_venous_malformation_perc_norm', 'bleeding_anticoagulation_therapy_perc_norm', 'bleeding_amyloid_angiopathy_perc_norm', 'bleeding_other_perc_norm']
        column_name = 'bleeding_arterial_hypertension_perc_norm'
        

        tmp_df = df[[main_col, 'bleeding_arterial_hypertension_perc_norm', 'bleeding_aneurysm_perc_norm', 'bleeding_arterio_venous_malformation_perc_norm', 'bleeding_anticoagulation_therapy_perc_norm', 'bleeding_amyloid_angiopathy_perc_norm', 'bleeding_other_perc_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == "CZ":
            title = "% PŘÍČINY KRVÁCENÍ u ICH"
            legend = [' arteriální hypertenze', 'aneurysma', 'arteriovenózní malformace', 'antikoagulační terapie', 'amyloidní angiopatie', 'jiné / neznámé']
        else:
            title = "% BLEEDING REASON for ICH"
            legend = ['arterial hypertension', 'aneurysm', 'arterio-venous malformation', 'anticoagulation therapy', 'amyloid angiopathy', 'other']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################
        # INTERVENTION PERFORMED #
        ##########################
        #sorted_by = ['intervention_endovascular_perc_norm', 'intervention_neurosurgical_perc_norm', 'intervention_other_perc_norm', 'intervention_referred_perc_norm', 'intervention_none_perc_norm']
        sorted_by = ['intervention_endovascular_perc_norm', 'intervention_neurosurgical_perc_norm', 'intervention_other_perc_norm', 'intervention_referred_perc_norm']
        column_name = 'intervention_endovascular_perc_norm'
        

        tmp_df = df[[main_col, 'intervention_endovascular_perc_norm', 'intervention_neurosurgical_perc_norm', 'intervention_other_perc_norm', 'intervention_referred_perc_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == "CZ":
            title = "% INTERVENCÍ PROVEDENÝCH u SAK"
            legend = ['Endovaskulární (coiling)', 'Neurochirurgická (clipping)', 'Jiný neurochirurgický výkon', 'Odeslán k intervenci do jiného centra']
        else:
            title = "% INTERVENTION PERFORMED for SAH"
            legend = ['Endovascular (coiling)', 'Neurosurgical (clipping)', 'Other neurosurgical treatment', 'Patient referred to another centre']


        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        
        '''
        ###############################
        # VENOUS THROMBOSIS TREATMENT #
        ###############################
        sorted_by = ['vt_treatment_anticoagulation_perc_norm', 'vt_treatment_thrombectomy_perc_norm', 'vt_treatment_local_thrombolysis_perc_norm', 'vt_treatment_local_neurological_treatment_perc_norm']
        column_name = 'vt_treatment_anticoagulation_perc_norm'
        legend = ['anticoagulation', 'thrombectomy', 'local thrombolysis', 'neurosurgical treatment']

        tmp_df = df[[main_col, 'vt_treatment_anticoagulation_perc_norm', 'vt_treatment_thrombectomy_perc_norm', 'vt_treatment_local_thrombolysis_perc_norm', 'vt_treatment_local_neurological_treatment_perc_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == "CZ":
            title = "% TYPŮ LÉČBY VENÓZNÍ TROMBÓZY U CVT"
        else:
            title = "% VENOUS THROMBOSIS TREATMENT for CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')
        '''

        ##########################################
        # % PATIENTS ASSESSED FOR REHABILITATION #
        ##########################################
        column_name = '% patients assessed for rehabilitation - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == 'CZ':
            title = "% VYŠETŘENÍ FYZIOTERAPEUTEM ze všech případů"
        else: 
            title = "% REHABILITATION ASSESSMENT out of all cases"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ################################
        # ATRIAL FIBRILLATION DETECTED #
        ################################
        sorted_by = ['% afib/flutter - Detected during hospitalization', '% afib/flutter - Newly-detected at admission', '% afib/flutter - Known']
        column_name = '% afib/flutter - Detected during hospitalization'
        

        tmp_df = df[[main_col, '% afib/flutter - Detected during hospitalization', '% afib/flutter - Newly-detected at admission', '% afib/flutter - Known']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == "CZ":
            title = "% FIBRILACÍ SÍNÍ DETEKOVANÝCH  u iCMP nebo TIA"
            legend = ['zjištěná během hospitalizace', 'nově zjištěná při přijetí', 'známá FS']
        else:
            title = "% ATRIAL FIBRILLATION DETECTED for IS, TIA"
            legend = ['detected during hospitalization', 'newly-detected at admission', 'known aFib']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ########################################
        # ATRIAL FIBRILLATION DETECTION METHOD #
        ########################################
        sorted_by = ['% afib detection method - Telemetry with monitor allowing automatic detection of aFib', '% afib detection method - Telemetry without monitor allowing automatic detection of aFib', '% afib detection method - Holter-type monitoring', '% afib detection method - EKG monitoring in an ICU bed with automatic detection of aFib', '% afib detection method - EKG monitoring in an ICU bed without automatic detection of aFib']
        column_name = '% afib detection method - Telemetry with monitor allowing automatic detection of aFib'
        

        tmp_df = df[[main_col, '% afib detection method - Telemetry with monitor allowing automatic detection of aFib', '% afib detection method - Telemetry without monitor allowing automatic detection of aFib', '% afib detection method - Holter-type monitoring', '% afib detection method - EKG monitoring in an ICU bed with automatic detection of aFib', '% afib detection method - EKG monitoring in an ICU bed without automatic detection of aFib']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == 'CZ':
            title = "% METOD DETEKCE FIBRILACÍ SÍNÍ"
            legend = ['Telemetrie s automatickou detekci FS', 'Telemetrie bez automatické detekce FS', 'Holterovo monitorování', 'EKG monitorace na JIP lůžku s automatickou detekcí FS', 'EKG monitorace na JIP lůžku bez automatické detekce FS']
        else:
            title = "% ATRIAL FIBRILLATION DETECTION METHOD"
            legend = ['Telemetry with monitoring', 'Telemetry without monitoring', 'Holter-type monitoring', 'EKG monitoring in an ICU bed with automatic detection of aFib', 'EKG monitoring in an ICU bed without automatic detection of aFib']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #######################################
        # AMBULATORY HEART RHYTHM RECOMMENDED #
        #######################################
        column_name = '% other afib detection method - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            #title = "% DOPORUČENÉ AMBULANTNÍ MONITORACE SRDEČNÍHO RYTMU u iCMP nebo TIA bez detekované fibrilace"
            title = "% pacientů s iCMP nebo TIA DOPORUČENÝCH NA AMBULANTNÍ MONITORACI SRDEČNÍHO RYTMU, u kterých nebyla detekována fibrilace síní"
        else:
            title = "% AMBULATORY HEART RHYTHM RECOMMENDED for IS, TIA without AFib detection"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################
        # CAROTID ARTERIES IMAGING - YES #
        ##################################
        column_name = '% carotid arteries imaging - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% PROVEDENÉHO ZOBRAZOVÁNÍ KAROTID u iCMP, TIA"
        else:
            title = "% CAROTID ARTERIES IMAGING PERFORMED for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##############
        # COMPARISON #
        ##############
        legends = []
        titles = []
        graph_types = []
        # DYSPHAGIA SCREENING #
        column_name = '% dysphagia screening - Guss test'
        
        if site_code == "CZ":
            titles.append("Screening dysfagie (%)")
            legend = ['GUSS test', 'Jiný test', 'Jiné pracoviště']
        else:
            titles.append("Dysphagia screening (%)")
            legend = ['GUSS test', 'Other test', 'Another centre']
        legends.append(legend)
        graph_types.append("stacked")

        df1 = fdf[[main_col, '% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre']]

        # CAROTID ARTERIES IMAGING #
        column_name = '% carotid arteries imaging - Yes'
        df2 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == "CZ":
            titles.append("Zobrazení karotid (%)")
        else:
            titles.append("Carotid arteries imaging (%)")
        graph_types.append("normal")

        # % RECOMMENDED TO A CEREBROVASCULAR EXPERT - RECOMMENDED #
        column_name = '% recommended to a cerebrovascular expert - Recommended'
        df3 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == "CZ":
            titles.append("% doporučených k vyšetření cerebrovaskulárním odborníkem")
        else:
            titles.append("Recommended to a cerebrovascular \nexpert (%)")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df1, df1=df2, df2=df3, presentation=prs, title=comp_title, titles=titles, graph_types=graph_types, legends=legends)
   
        ##################################################
        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB #
        ##################################################
        column_name = '% patients prescribed anticoagulants with aFib'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            #title = "% PACIENTŮ S FIBRILACÍ SÍNÍ A PŘEDEPSANÝMI ANTIKOAGULANTY U iCMP, TIA"
            title = "POUŽITÍ ANTIKOAGULAČNÍ LÉČBY U PACIENTŮ S FIBRILACÍ SÍNÍ a iCMP nebo TIA (%)"
        else:
            title = "% PATIENTS WITH AFIB, PRESCRIBED ANTICOAGULANTS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##############################
        # ANTICOAGULANTS PRESCRIBED  #
        ##############################
        column_name = '% patients receiving Vit. K antagonist'
        legend = ['Vitamin K', 'Dabigatran', 'Rivaroxaban', 'Apixaban', 'Edoxaban', 'LMWH, heparin - profylaktická dávka', 'LMWH, heparin - antikoagulační, terapeutická dávka']

        tmp_df = df[[main_col, '% patients prescribed anticoagulants with aFib', '% patients receiving Vit. K antagonist', '% patients receiving dabigatran', '% patients receiving rivaroxaban', '% patients receiving apixaban', '% patients receiving edoxaban', '% patients receiving LMWH or heparin in prophylactic dose', '% patients receiving LMWH or heparin in full anticoagulant dose']]

        tmp_df = tmp_df.sort_values(['% patients prescribed anticoagulants with aFib'], ascending = True)

        if site_code == "CZ":
            title = "POUŽITÍ ANTIKOAGULAČNÍ LÉČBY U PACIENTŮ S FIBRILACÍ SÍNÍ a iCMP nebo TIA (%)"
        else:
            title = "% PATIENTS WITH AFIB, PRESCRIBED ANTICOAGULANTS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #########################################################
        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB (HOME) #
        #########################################################
        column_name = '% afib patients discharged home with anticoagulants'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            #title = "% PACIENTŮ PROPUŠTĚNÝCH DOMŮ, S FIBRILACÍ SÍNÍ A PŘEDEPSANÝMI ANTIKOAGULANTY u iCMP, TIA"
            title = "% pacientů s PŘEDEPSANÝMI ANTIKOAGULANTY u iCMP nebo TIA pro FIBRILACI SÍNÍ, pokud byli PROPUŠTĚNI DOMŮ"
        else:
            title = "% PATIENTS DISCHARGED HOME WITH AFIB, PRESCRIBED ANTICOAGULANTS \nfor IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################################
        # PATIENTS PRESCRIBED ANTIPLATELETS WITHOUT AFIB #
        ##################################################
        column_name = '% patients prescribed antiplatelets without aFib'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
           # title = "% PACIENTŮ BEZ FIBRILACÍ SÍNÍ A PŘEDEPSANOU ANTOAGREGACÍ u iCMP, TIA"
            title = "% pacientů LÉČENÝCH PROTIDESTIČKOVÝMI LÉKY pro iCMP nebo TIA, pokud neměli FIBRILACI SÍNÍ"
           # title = "% PACIENTŮ LÉČENÝCH PROTIDESTIČKOVÝMI LÉKY bez FIBRILACE SÍNÍ u iCMP nebo TIA"
        else:
            title = "% PATIENTS WITHOUT AFIB, PRESCRIBED ANTIPLATELETS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        '''        
        ###################################
        # NOT PRESCRIBED, BUT RECOMMENDED #
        ###################################
        column_name = '% patients not prescribed antithrombotics, but recommended'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% PACIENTŮ BEZ PŘEDEPSANÝCH, ALE S DOPORUČENÝMI ANTITROMBOTIKY u iCMP, TIA"
        else:
            title = "% PATIENTS NOT PRESCRIBED, BUT RECOMMENDED ANTITHROMBOTICS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)
        '''
        '''

        ############
        # WITH CVT #
        ############
        ##############################
        # ANTITHROMBOTICS PRESCRIBED #
        ##############################
        column_name = '% patients receiving antiplatelets with CVT'
        legend = ['Antiplatelets', 'Vitamin K', 'Dabigatran', 'Rivaroxaban', 'Apixaban', 'Edoxaban', 'LMWH or heparin in prophylactic dose', 'LMWH or heparin in anticoagulant dose']

        tmp_df = df[[main_col, '% patients prescribed antithrombotics with CVT', '% patients receiving antiplatelets with CVT', '% patients receiving Vit. K antagonist with CVT', '% patients receiving dabigatran with CVT', '% patients receiving rivaroxaban with CVT', '% patients receiving apixaban with CVT', '% patients receiving edoxaban with CVT', '% patients receiving LMWH or heparin in prophylactic dose with CVT', '% patients receiving LMWH or heparin in full anticoagulant dose with CVT']]

        tmp_df = tmp_df.sort_values(['% patients prescribed antithrombotics with CVT'], ascending = True)

        if site_code == "CZ":
            title = "% PŘEDEPSANÉ ANTITROMBOTICKÉ LÉČBY u iCMP, TIA, CVT"
        else:
            title = "% ANTITHROMBOTICS PRESCRIBED for IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##################################################
        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB #
        ##################################################
        column_name = '% patients prescribed anticoagulants with aFib with CVT'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% PACIENTŮ S FIBRILACÍ SÍNÍ A PŘEDEPSANÝMI ANTIKOAGULANTY u iCMP, TIA, CVT"
        else:
            title = "% PATIENTS WITH AFIB, PRESCRIBED ANTICOAGULANTS for IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################################
        # PATIENTS PRESCRIBED ANTIPLATELETS WITHOUT AFIB #
        ##################################################
        column_name = '% patients prescribed antiplatelets without aFib with CVT'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% PACIENTŮ BEZ FIBRILACE SÍNÍ, S PŘEDEPSANOU ANTIAGREGACÍ u iCMP, TIA, CVT"
        else:
            title = "% PATIENTS WITHOUT AFIB, PRESCRIBED ANTIPLATELETS for IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        
        ###################################
        # NOT PRESCRIBED, BUT RECOMMENDED #
        ###################################
        column_name = '% patients not prescribed antithrombotics, but recommended with CVT'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% PACIENTŮ S POUZE DOPORUČENOU LÉČBOU ANTITROMBOTIKY u iCMP, TIA, CVT"
        else:
            title = "% PATIENTS NOT PRESCRIBED, BUT RECOMMENDED ANTITHROMBOTICS \nfor IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        '''
        ################
        # COMPARISON ###
        ################
        legends = []
        titles = []
        graph_types = []
        # PRESCRIBED ANTIPLATELETS #
        column_name = '% patients prescribed antiplatelets without aFib'
        df1 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == "CZ":
            titles.append("% pacientů s antiagregací bez fibrilace síní")
        else:
            titles.append("% patients prescribed antiplatelets without aFib")
        graph_types.append("normal")

        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB #
        column_name = '% patients prescribed anticoagulants with aFib'
        df2 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == "CZ":
            titles.append("% pacientů s předepsanými antikoagulanty a s fibrilací síní")
        else:
            titles.append("% patients prescribed anticoagulants with aFib")
        graph_types.append("normal")

        #  % PATIENTS PRESCRIBED ANTITHROMBOTICS WITH AFIB  #
        column_name = '% patients prescribed antithrombotics with aFib'
        df3 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == "CZ":
            titles.append("% pacientů s předepsanými antitrombotiky a  s fibrilací síní")
        else:
            titles.append("% patients prescribed antithrombotics with aFib")
        graph_types.append("normal")

        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB (HOME) #
        column_name = '% afib patients discharged home with anticoagulants'
        df4 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == "CZ":
            titles.append("% pacientů s fibrilací síní propuštěných domů s antikoagulanty")
        else:
            titles.append("% aFib patients discharged home with anticoagulants")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df1, df1=df2, df2=df3, df3=df4, presentation=prs, title=comp_title, titles=titles, graph_types=graph_types, legends=legends)

        ##########################
        # DISCHARGE WITH STATINS #
        ##########################
        column_name = '% patients prescribed statins - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            #title = "% PROPUŠTĚNÝCH PACIENTŮ SE STATINY u iCMP nebo TIA"
            title = "% pacientů LÉČENÝCH STATINY u iCMP nebo TIA, pokud byli PROPUŠTĚNI DOMŮ"
        else:
            title = "% DISCHARGED WITH STATINS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #############################
        # CAROTID STENOSIS DETECTED #
        #############################
        column_name = '% carotid stenosis - >70%'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% STENÓZY KAROTID NAD 70% u iCMP nebo TIA"
        else:
            title = "% CAROTID STENOSIS OF OVER 70 PERCENT for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################################
        # % CAROTID STENOSIS FOLLOWUP - YES, BUT PLANNED #
        ##################################################
        sorted_by = ['% carotid stenosis followup - Yes, but planned', '% carotid stenosis followup - Referred to another centre']
        column_name = '% carotid stenosis followup - Yes, but planned'
        

        tmp_df = df[[main_col, '% carotid stenosis followup - Yes, but planned', '% carotid stenosis followup - Referred to another centre']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == "CZ":
            title = "% ENDARTEREKTOMIÍ NEBO ANGIOPLASTIK / STENTINGŮ PROVEDENÝCH NEBO PLÁNOVANÝCH u iCMP, TIA s ICA STENÓZOU > 70%"
            legend = ['Ano nebo plánováno', 'Odeslán/a do jiného centra']
        else:
            title = "% ENDARTERECTOMY OR ANGIOPLASTY / STENTING DONE OR PLANNED for IS, TIA with ICA STENOSIS > 70%"
            legend = ['Yes or planned', 'Referred to another centre']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #############################################################
        # % ANTIHYPERTENSIVE MEDICATION PRESCRIBED out of all cases #
        #############################################################
        column_name = '% prescribed antihypertensives - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% PŘEDEPSANÉ ANTIHYPERTENZNÍ LÉČBY ze všech případů"
        else:
            title = "% ANTIHYPERTENSIVE MEDICATION PRESCRIBED out of all cases"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##############
        # COMPARISON #
        ##############
        legends = []
        titles = []
        graph_types = []
        # DISCHARGE WITH STATINS #
        column_name = '% patients prescribed statins - Yes'
        df1 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == "CZ":
            titles.append("% pacientů s předepsanými statiny")
        else:
            titles.append("% patients prescribed statins")
        graph_types.append("normal")

        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB #
        column_name = '% prescribed antihypertensives - Yes'
        df2 = fdf[[main_col, column_name]]
        legend = []
        legends.append(legend)
        if site_code == "CZ":
            titles.append("% pacientů s předepsanými antihypertenzívy")
        else:
            titles.append("% patients prescribed antihypertensives")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df1, df1=df2, presentation=prs, title=comp_title, titles=titles, graph_types=graph_types, legends=legends)

        ######################################################
        # % RECOMMENDED TO A SMOKING CESSATION PROGRAM - YES #
        ######################################################
        column_name = '% recommended to a smoking cessation program - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% KUŘÁKŮ, JIMŽ BYL DOPORUČEN PROGRAM K ODVYKÁNÍ KOUŘENÍ"
        else:
            title = "% RECOMMENDED TO A SMOKING CESSATION PROGRAM out of smokers"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###########################################################
        # % RECOMMENDED TO A CEREBROVASCULAR EXPERT - RECOMMENDED #
        ###########################################################
        column_name = '% recommended to a cerebrovascular expert - Recommended'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "% PACIENTŮ ODESLANÝCH K CEREBROVASKULÁRNÍMU ODBORNÍKOVI ze všech případů"
        else:
            title = "% RECOMMENDED TO A CEREBROVASCULAR EXPERT out of all cases"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #########################
        # DISCHARGE DESTINATION #
        #########################
        sorted_by = ['% discharge destination - Home', '% discharge destination - Transferred within the same centre', '% discharge destination - Transferred to another centre', '% discharge destination - Social care facility', '% discharge destination - Dead']
        column_name = '% discharge destination - Home'
        

        tmp_df = df[[main_col, '% discharge destination - Home', '% discharge destination - Transferred within the same centre', '% discharge destination - Transferred to another centre', '% discharge destination - Social care facility', '% discharge destination - Dead']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == "CZ":
            title = "% MÍSTO URČENÍ U PROPUŠTĚNÍ / PŘEKLADU - DISCHARGE DESTINATION"
            legend = ['domů', 'překlad v rámci stejného zdravotnického zařízení (ZZ)', 'překlad do jiného zdravotnického zařízení (ZZ)', 'zařízení sociální péče', 'zemřel/a']
        else:
            title = "% DISCHARGE DESTINATION"
            legend = ['home', 'transferred within the same centre', 'transferred to another centre', 'social care facility', 'dead']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ############################################################################
        # DISCHARGE DESTINATION - DEPARTMENT TRANSFERRED TO WITHIN THE SAME CENTRE #
        ############################################################################
        sorted_by = ['% transferred within the same centre - Acute rehabilitation', '% transferred within the same centre - Post-care bed', '% transferred within the same centre - Another department']
        column_name = '% transferred within the same centre - Acute rehabilitation'
        

        tmp_df = df[[main_col, '% transferred within the same centre - Acute rehabilitation', '% transferred within the same centre - Post-care bed', '% transferred within the same centre - Another department']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == "CZ":
            title = "% PACIENTŮ PŘELOŽENÝCH V RÁMCI STEJNÉHO CENTRA"
            legend = ['Oddělení akutní rehabilitace', 'Lůžko následné péče', 'Jiné oddělení']
        else:
            title = "% DISCHARGE DESTINATION - PATIENT TRANSFERRED WITHIN THE SAME CENTRE"
            legend = ['Acute rehabilitation', 'Post-care bed', 'Another department']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ####################################################
        # % DISCHARGE DESTINATION - CENTRE TRANSFERRED TO  #
        ####################################################
        sorted_by = ['% transferred to another centre - Stroke centre', '% transferred to another centre - Comprehensive stroke centre', '% transferred to another centre - Another hospital']
        column_name = '% transferred to another centre - Stroke centre'
        

        tmp_df = df[[main_col, '% transferred to another centre - Stroke centre', '% transferred to another centre - Comprehensive stroke centre', '% transferred to another centre - Another hospital']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == "CZ":
            title = "% PACIENTŮ PŘELOŽENÝCH DO JINÉHO CENTRA"
            legend = ['Iktové centrum', 'Komplexní cerebrovaskulární centrum', 'Jiná nemocnice']
        else:
            title = "% DISCHARGE DESTINATION - PATIENT TRANSFERRED TO ANOTHER CENTRE"
            legend = ['Stroke centre', 'Comprehensive stroke centre', 'Another hospital']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #################################################################################
        # % DISCHARGE DESTINATION - DEPARTMENT TRANSFERRED TO WITHIN TO ANOTHER CENTRE  #
        #################################################################################
        sorted_by = ['% department transferred to within another centre - Acute rehabilitation', '% department transferred to within another centre - Post-care bed', '% department transferred to within another centre - Neurology', '% department transferred to within another centre - Another department']
        column_name = '% department transferred to within another centre - Acute rehabilitation'
        

        tmp_df = df[[main_col, '% department transferred to within another centre - Acute rehabilitation', '% department transferred to within another centre - Post-care bed', '% department transferred to within another centre - Neurology', '% department transferred to within another centre - Another department']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        if site_code == "CZ":
            title = "% PACIENTŮ PŘELOŽENÝCH DO JINÉHO CENTRA PODLE TYPU CÍLOVÉHO ODDĚLENÍ"
            legend = ['Oddělení akutní rehabilitace', 'Lůžko následné péče', 'Neurologie', 'Jiné oddělení']
        else:
            title = "% DISCHARGE DESTINATION - PATIENT TRANSFERRED TO ANOTHER CENTRE (DEPARTMENT)"
            legend = ['Acute rehabilitation', 'Post-care bed', 'Neurology', 'Another department']

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ########################
        # MEDIAN DISCHARGE mRS #
        ########################
        column_name = 'Median discharge mRS'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "MEDIÁN mRS U PROPUŠTĚNÍ"
        else:
            title = "MEDIAN DISCHARGE MRS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###############################
        # MEDIAN HOSPITAL STAY (DAYS) #
        ###############################
        column_name = 'Median hospital stay (days)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        if site_code == "CZ":
            title = "DÉLKA POBYTU V NEMOCNICI (VE DNECH)"
        else:
            title = "MEDIAN HOSPITAL STAY (DAYS)"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        # set pptx output name (for cz it'll be presentation_CZ.pptx)
        working_dir = os.getcwd()
        pptx = self.report + "_" + site_code + "_" + self.quarter + "_national_comparison.pptx"
        presentation_path = os.path.normpath(os.path.join(working_dir, pptx))

    


        prs.save(presentation_path)


        


