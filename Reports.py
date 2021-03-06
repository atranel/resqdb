# -*- coding: utf-8 -*-
"""
Created on May 20, 2019

@author: Marie Jankujova
"""

import sys
import os
from datetime import datetime, date, timedelta, time
import sqlite3
import pandas as pd
import numpy as np
from numpy import inf
import pytz
import logging
import scipy.stats as st
from scipy.stats import sem, t
from scipy import mean
from resqdb.Calculation import FilterDataset
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.oxml.xmlchemy import OxmlElement
import xlsxwriter
import json
from pptx.oxml.table import CT_Table
from pptx.enum.text import PP_ALIGN
import statistics


class Reports:
    """ The class generating the SITS-like reports for recanalization procedure. 

    :param df: the preprocessed data
    :type df: pandas dataframe
    :param year: the year for which reports are generated
    :type year: int
    :param month: the last month to be included in the calculation
    :type month: int
    :param country: the country code (only CZ for now)
    :type country: str

    """
    def __init__(self, df, year, month, country):
    
        # create dataframe with regions, to each region assign population and hospitals
        path = os.path.join(os.path.dirname(__file__), 'tmp', 'regions.json')
        with open(path, 'r', encoding='utf-8') as json_file:
            self.regions = json.load(json_file)

        # Create dataframe with hospitals who do thrombectomy
        self.hospitals_mt = ['CZ_034', # FN Brno
            'CZ_013', # Ceske Budejovice
            'CZ_017', # Na Homolce
            'CZ_039', # Usti nad Labem
            'CZ_006', # Vitkovice
            'CZ_010', # Hradec Kralove
            'CZ_036', # Liberec
            'CZ_009', # VFN Neurologie Praha
            'CZ_002', # Motol
            'CZ_018', # FN Ostrava
            'CZ_001', # St. Anne Brno
            'CZ_041', # Ustredni vojenska nemocnice
            'CZ_025', # Plzen
            'CZ_042', # Olomouc
            'CZ_008', # Kralovske Vinohrady,
            'CZ'
        ]

        self.year = year
        self.country = country
        self.month = month
        self.country_name = 'Česká republika'

        debug = 'reports_debug_' + datetime.now().strftime('%d-%m-%Y') + '.log'
        # Create log file in the working folder
        log_file = os.path.join(os.getcwd(), debug)
        logging.basicConfig(filename=log_file,
                            filemode='a',
                            format='%(asctime)s,%(msecs)d %(name)s %(levelname)s %(message)s',
                            datefmt='%H:%M:%S',
                            level=logging.DEBUG)

        # Get only dataframe for selected country
        # Rename 'RES-Q reports name' column to 'Site Name'
        if 'RES-Q reports name' in df.columns:
            df.rename(columns={'Site Name': 'Site Name Old', 'RES-Q reports name': 'Site Name'}, inplace=True)

        fd_ojb = FilterDataset(df=df, country=self.country)
        df = fd_ojb.fdf.copy()
        df = df.loc[df['Protocol ID'] != 'CZ_052'].copy()
        #df.to_csv("test.csv", sep=",")

        dev_form = df.loc[df['crf_parent_name'] == 'F_RESQ_IVT_TBY_1565_DEVCZ10'].copy()
        dev_form = dev_form[['Site Name', 'crf_parent_name', 'Subject ID', 'HOSPITAL_DATE', 'DISCHARGE_DATE']]
        dev_form = dev_form.astype(str)
        dev_form.to_csv("development_form.csv", sep=",", encoding='utf-8', index=False)

        # Get all sites which have more than 5 patients in the development form
        development_form_counts = df.loc[df['crf_parent_name'] == 'F_RESQ_IVT_TBY_1565_DEVCZ10'].groupby(['Protocol ID', 'Site Name', 'crf_parent_name']).size().reset_index().rename(columns={0: 'n'})
        development_form_counts.to_csv("development_form_counts.csv", sep=",", encoding='utf-8', index=False)
        """
        self.development_forms_counts = df.loc[df['crf_parent_name'] == 'F_RESQ_IVT_TBY_1565_DEVCZ10'].groupby(['Protocol ID', 'Site Name', 'crf_parent_name']).size().reset_index().rename(columns={0: 'n'})
        # Get sites for which patients in development should be removed
        to_remove = self.development_forms_counts.loc[self.development_forms_counts['n'] <= 5, 'Protocol ID'].tolist()
        # Get indexes which should be removed
        indexes = df.index[(df['Protocol ID'].isin(to_remove)) & (df['crf_parent_name'] == 'F_RESQ_IVT_TBY_1565_DEVCZ10')].tolist()
        df.drop(indexes, inplace=True)
        df.reset_index(drop=True, inplace=True)
        """

        tmp_country_df = df.copy()
        tmp_country_df['Site Name'] = self.country_name
        tmp_country_df['Protocol ID'] = 'CZ'
        df = df.append(tmp_country_df, ignore_index=True, sort=False)
        
        self.country_df = df
        self.df = df.copy()

        # Get site names to hospitals_mt
        self.site_id_mapped_to_site_name = self.df[self.df['Protocol ID'].isin(self.hospitals_mt)][['Protocol ID', 'Site Name']].drop_duplicates(subset='Protocol ID', keep='first').reset_index()
        
        self.site_id_mapped_to_site_name.drop(['index'], inplace=True, axis=1)
        
        # Filter dataframes per month
        self.filtered_dfs = self.filter_dataframe()
        self.names = list(self.filtered_dfs.keys())
        self.incorrect_ivtpa = {}
        self.incorrect_tby = {}
        self.thrombolysis_stats_df = self.calculate_thrombolysis()
        self.thrombectomy_stats_df = self.calculate_thrombectomy()
        self.statistic_region_dfs = self.calculate_statistic_per_region()

        self.save_excel()


    def filter_dataframe(self):
        """ The function filtering the preprocessed data for each month. 

        :returns: the dictionary of filtered dataframes
        """
        dfs = {}
        # Get current date
        current_year = datetime.today().year
        current_month = self.month

        # Filter dataframe per month
        for month in range(1, current_month + 1):
            if current_month == 12:
                start_date = datetime(self.year, current_month, 1, 0, 0)
                end_date = datetime(self.year, current_month, 31, 0, 0) - timedelta(days=1)
            else:
                start_date = datetime(self.year, month, 1, 0, 0)
                end_date = datetime(self.year, (month % 12 + 1), 1, 0, 0) - timedelta(days=1)

            # Create object FilterDataset
            fd_ojb = FilterDataset(df=self.df, country=self.country, date1=start_date, date2=end_date)
            df = fd_ojb.fdf.copy()
            df['Protocol ID']
            df = df.loc[~df['Protocol ID'].isin(['CZ_052'])].copy()


            # Add dataframe into dictionary
            dfs[month] = df

        # Filter dataframe for whole year
        start_date = datetime(self.year, 1, 1, 0, 0)
        # End date from current_month
        if current_month == 12:
            end_date = datetime(self.year, current_month, 31, 0, 0) - timedelta(days=1)
        else:
            end_date = datetime(self.year, (current_month % 12 + 1), 1, 0, 0) - timedelta(days=1)
        fd_obj = FilterDataset(df=self.df, country=self.country, date1=start_date, date2=end_date)
        df = fd_obj.fdf.copy()
        df = df.loc[~df['Protocol ID'].isin(['CZ_052'])].copy()

        dfs[str(self.year)] = df
        
        return dfs

    
    def count_patients(self, df, statistic):
        """ The function calculating the number of patients grouped by Protocol ID. 

        :param df: the dataframe with preprocessed data
        :type df: pandas dataframe
        :param statistic: the dataframe with calculated statistics
        :type statistic: pandas dataframe
        :returns: the column with patient numbers
        """
        tmp = df.groupby(['Protocol ID']).size().reset_index(name='total_patients')
        tmp_df = statistic.merge(tmp, how='outer')
        tmp_df.fillna(0, inplace=True)

        return tmp_df['total_patients']

    
    def time_diff(self, visit_date, hospital_date):
        """ The function calculating the difference in minutes between two dates. 

        :param visit_date: the last seen normal date
        :type visit_date: date
        :param hospital_date: the date of hospitalization
        :type hospital_date: date
        :returns: the difference in minutes
        """
        if type(visit_date) is pd.Timestamp and type(hospital_date) is pd.Timestamp:
            time_diff = hospital_date - visit_date
            # Convert difference to minutes
            total_minutes = time_diff.total_seconds() / 60.0
        else:
            total_minutes = 0

        if total_minutes < 0 or total_minutes > 40000:
            total_minutes = 0
        
        return total_minutes

    
    def get_region(self, site_id):
        """ The function returning the region based on Site ID. 

        :param site_id: the site ID
        :type site_id: str
        :returns: the name of the region
        """
        for key, value in self.regions.items():
            if site_id in value['hospitals']:
                return key           

    
    def get_incorrect_times(self, admission_time, recan_time, maximum):
        """ The function checking if time was incorrectly entered. Based on negative values, higher values than realistic value or zero value. 

        :param admission_time: the time of the admission/hospitalization
        :type admission_time: time
        :param recan_time: the time of recanalization procedures (needle time/groin time)
        :type recan_time: time
        :param maximum: the realistic time for the recanalizaiton procedure
        :type maximum: int
        :returns: `True` if the condition was not met
        """
       
        timeformat = '%H:%M:%S'
        if admission_time is None or recan_time is None or pd.isnull(admission_time) or pd.isnull(recan_time):
            incorrect = True
        elif admission_time == 0 or recan_time == 0:
            incorrect = True
        else:
            if isinstance(admission_time, time) and isinstance(recan_time, time):
                tdelta = datetime.combine(date.today(), recan_time) - datetime.combine(date.today(), admission_time)
            elif isinstance(admission_time, time):
                tdelta = datetime.strptime(recan_time, timeformat) - datetime.combine(date.today(), admission_time)
            elif isinstance(recan_time, time):
                tdelta = datetime.strptime(recan_time, timeformat) - datetime.strptime(admission_time, timeformat)
            else:
                tdelta = datetime.strptime(recan_time, timeformat) - datetime.strptime(admission_time, timeformat)	
            tdelta_min = tdelta.total_seconds()/60.0

            if tdelta_min < -1000:
                # Add one day if time was after midnight (1 day = 86400 seconds = 1440 minutes)
                tdelta_min = tdelta_min + 1440

            if tdelta_min <= 0 or tdelta_min > maximum:
                incorrect = True
            else:
                incorrect = False

        return incorrect

    def calculate_thrombolysis(self):
        """ The function calculating the result statistic for patients who have recieved the thrombolysis. """
        stats_dfs = {}

        # Iterate through filtered dataframes
        for name, df in self.filtered_dfs.items():
             
            statistic = self.country_df.groupby(['Protocol ID', 'Site Name']).size().reset_index(name="Total Patients")			# Get Protocol IDs and Total Patients
            # Calculate IVtPa median
            ischemic_cmp = df[df['STROKE_TYPE'].isin([1])].copy()		
            thrombolysis_df = ischemic_cmp.loc[ischemic_cmp['IVT_DONE'].isin([1])].copy()	

            # print('Thrombolysis {} - {}'.format(name, len(thrombolysis_df)))
            # only patients with ischemic stroke
            # thrombolysis_df = ischemic_cmp[ischemic_cmp['RECANALIZATION_PROCEDURES'].isin([2,3,5])].copy() 	# only patients with ischemic stroke who underwent recanalizaiton procedure (IVtPa, IVtPa + TBY, IVtPa + referred for TBY)
            if thrombolysis_df.empty:
                statistic['Median DTN (minutes)'] = 0
                statistic['# IVT'] = 0
                statistic['Median last seen normal'] = 0
                statistic['# incorrect IVtPa times'] = 0
                statistic['% incorrect IVtPa times'] = 0
            else:
                thrombolysis_df.fillna(0, inplace=True)
                #thrombolysis_df['IVTPA'] = thrombolysis_df['IVT_ONLY_NEEDLE_TIME'] + thrombolysis_df['IVT_ONLY_NEEDLE_TIME_MIN'] + thrombolysis_df['IVT_TBY_NEEDLE_TIME'] + thrombolysis_df['IVT_TBY_NEEDLE_TIME_MIN'] + thrombolysis_df['IVT_TBY_REFER_NEEDLE_TIME'] + thrombolysis_df['IVT_TBY_REFER_NEEDLE_TIME_MIN']       			# get one column with all needle times

                statistic['Total patients undergone IVT'] = self.count_patients(df=thrombolysis_df, statistic=statistic)

                # Apr 22, 2020 - exclude patients if hospital stroke and times for IVT as timestamps
                thrombolysis_df = thrombolysis_df.loc[
                    ~(thrombolysis_df['HOSPITAL_STROKE_IVT_TIMESTAMPS'].isin([1]))
                ].copy()

                # thrombolysis_df = thrombolysis_df.loc[
                #     ~thrombolysis_df['HOSPITAL_STROKE'].isin([1])
                # ].copy() 
                # Get number of incorrectly entered times
                thrombolysis_df['INCORRECT_TIMES'] = False
                thrombolysis_df['INCORRECT_TIMES'] = thrombolysis_df.apply(lambda x: self.get_incorrect_times(x['IVT_ONLY_ADMISSION_TIME'], x['IVT_ONLY_BOLUS_TIME'], 400) if x['RECANALIZATION_PROCEDURES'] == 2 and x['IVT_ONLY'] == 2 else x['INCORRECT_TIMES'], axis=1)
                thrombolysis_df['INCORRECT_TIMES'] = thrombolysis_df.apply(lambda x: self.get_incorrect_times(x['IVT_TBY_ADMISSION_TIME'], x['IVT_TBY_BOLUS_TIME'], 400) if x['RECANALIZATION_PROCEDURES'] == 3 and x['IVT_TBY'] == 2 else x['INCORRECT_TIMES'], axis=1)
                thrombolysis_df['INCORRECT_TIMES'] = thrombolysis_df.apply(lambda x: self.get_incorrect_times(x['IVT_TBY_REFER_ADMISSION_TIME'], x['IVT_TBY_REFER_BOLUS_TIME'], 400) if x['RECANALIZATION_PROCEDURES'] == 5 and x['IVT_TBY_REFER'] == 2 else x['INCORRECT_TIMES'], axis=1)
                thrombolysis_df['INCORRECT_TIMES'] = thrombolysis_df.apply(lambda x: True if (x['IVTPA'] <= 0 or x['IVTPA'] > 400) and x['IVT_ONLY'] == 1 else x['INCORRECT_TIMES'], axis=1)
                thrombolysis_df['INCORRECT_TIMES'] = thrombolysis_df.apply(lambda x: True if (x['IVTPA'] <= 0 or x['IVTPA'] > 400) and x['IVT_TBY'] == 1 else x['INCORRECT_TIMES'], axis=1)
                thrombolysis_df['INCORRECT_TIMES'] = thrombolysis_df.apply(lambda x: True if (x['IVTPA'] <= 0 or x['IVTPA'] > 400) and x['IVT_TBY_REFER'] == 1 else x['INCORRECT_TIMES'], axis=1)

                incorrect_ivtpa_times = thrombolysis_df[
                    (thrombolysis_df['INCORRECT_TIMES'] == True) & 
                    (thrombolysis_df['HOSPITAL_STROKE_TBY_TIMESTAMPS'] != 1)
                    ].copy()
                incorrect_ivtpa_times_save = incorrect_ivtpa_times.loc[
                    incorrect_ivtpa_times['Protocol ID'] != "CZ"].copy()
                incorrect_ivtpa_times_save.to_csv('incorrect_ivtpa_times.csv', sep=',')

                thrombolysis = thrombolysis_df[(thrombolysis_df['IVTPA'] > 0) & (thrombolysis_df['IVTPA'] <= 400)].copy()

                if thrombolysis.empty:
                    statistic['Median DTN (minutes)'] = 0
                    statistic['# IVT'] = 0
                    statistic['Median last seen normal'] = 0
                    if incorrect_ivtpa_times.empty:
                        statistic['# incorrect IVtPa times'] = 0
                        statistic['% incorrect IVtPa times'] = 0
                    else:
                        statistic['# incorrect IVtPa times'] = self.count_patients(df=incorrect_ivtpa_times, statistic=statistic)
                        statistic['% incorrect IVtPa times'] = round((statistic['# incorrect IVtPa times'] / statistic['Total patients undergone IVT'])*100, 2)
                else:
                    thrombolysis_grouped = thrombolysis.groupby(['Protocol ID']).IVTPA.agg(['median']).rename(columns={'median': 'Median DTN (minutes)'}).reset_index() # calculate median DTN per site
                    statistic = statistic.merge(thrombolysis_grouped, how='outer') # Merge with statistic dataframe

                    # Get number of IVTs on IC/KCC
                    # statistic['# IVT'] = self.count_patients(df=thrombolysis, statistic=statistic)
                    statistic['# IVT'] = self.count_patients(df=thrombolysis, statistic=statistic)
                    statistic.loc[statistic['Protocol ID'] == 'CZ', '# IVT'] = int(statistics.mean(statistic.loc[statistic['Protocol ID'] != 'CZ']['# IVT'].tolist()))

                    # Get difference in minutes between hospitalization and last visit
                    #thrombolysis['LAST_SEEN_NORMAL'] = thrombolysis.apply(lambda x: self.time_diff(x['VISIT_TIMESTAMP'], x['HOSPITAL_TIMESTAMP']), axis=1)
                    #thrombolysis['LAST_SEEN_NORMAL'].fillna(0, inplace=True)
                    last_seen_normal_grouped = thrombolysis[thrombolysis['LAST_SEEN_NORMAL'] != 0].groupby(['Protocol ID']).LAST_SEEN_NORMAL.agg(['median']).rename(columns={'median': 'Median last seen normal'}).reset_index()
                    statistic = statistic.merge(last_seen_normal_grouped, how='outer') # Merge with statistic dataframe

                    if incorrect_ivtpa_times.empty:
                        statistic['# incorrect IVtPa times'] = 0
                        statistic['% incorrect IVtPa times'] = 0
                    else:
                        statistic['# incorrect IVtPa times'] = self.count_patients(df=incorrect_ivtpa_times, statistic=statistic)
                        statistic['% incorrect IVtPa times'] = round((statistic['# incorrect IVtPa times'] / statistic['Total patients undergone IVT'])*100, 2)

                statistic.loc[statistic['Protocol ID'] == 'CZ', 'Total patients undergone IVT'] = int(statistics.mean(statistic.loc[statistic['Protocol ID'] != 'CZ']['Total patients undergone IVT'].tolist()))

            statistic.fillna(0, inplace=True)

            #statistic.to_csv(str(name) + '.csv', sep=',')

            stats_dfs[name] = statistic
            self.incorrect_ivtpa[name] = incorrect_ivtpa_times_save
        
        return stats_dfs

    def calculate_thrombectomy(self):
        """ The function calculating the result statistic for patients who have recieved the thrombectomy. """
        stats_dfs = {}

        # Iterate through filtered dataframes
        for name, df in self.filtered_dfs.items():

            first_hosp_mapping = {
                'CZ_004': 'Municipal Hospital Ostrava - Neurologické oddělení',
                'CZ_024': 'Krajská zdravotní, a.s. - Nemocnice Chomutov, o.z.',
                'CZ_025': 'Faculty Hospital Plzen',
                'CZ_026': 'Hospital Teplice',
                'CZ_041': 'Central military hospital - Praha 6'
            }
            
            # Calculate IVtPa median
            # thrombectomy_df = df[(df['Protocol ID'].isin(self.hospitals_mt)) & df['RECANALIZATION_PROCEDURES'].isin([3,4]) & df['STROKE_TYPE'].isin([1])].copy()
            df['TBY_DONE'], df['INCLUDE_MEDIAN'] = zip(*df.apply(lambda x: (1, True) if (x['RECANALIZATION_PROCEDURES'] in [7,8] and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_CZ') else (x['TBY_DONE'], True), axis=1))
            #df['TBY_DONE'] = df.apply(lambda x: 1 if (x['RECANALIZATION_PROCEDURES'] in [7,8] and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_CZ') else x['TBY_DONE'], axis=1)
            #df['INCLUDE_MEDIAN'] = df.apply(lambda x: True if (x['RECANALIZATION_PROCEDURES'] in [7,8] and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_CZ') else True, axis=1)

            #df['TBY_DONE'], df['INCLUDE_MEDIAN'] = zip(*df.apply(lambda x: (1, True) if x['RECANALIZATION_PROCEDURES'] in [7] and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_1565_DEVCZ10' and x['Protocol ID'] == 'CZ_041' else (x['TBY_DONE'], True), axis=1))
            #df['TBY'] = df.apply(lambda x: x['TBY_REFER_ALL_DIDO_TIME'] if x['RECANALIZATION_PROCEDURES'] in [7] and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_1565_DEVCZ10' and x['Protocol ID'] == 'CZ_041' else x['TBY'], axis=1)

            #df['TBY'] = df.apply(lambda x: x['TBY_REFER_ALL_DIDO_TIME'] if x['RECANALIZATION_PROCEDURES'] in [7] and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_CZ' and x['Protocol ID'] == 'CZ_041' else x['TBY'], axis=1)
            
            df['FIRST_HOSPITAL'] = df.apply(lambda x: 1 if (x['crf_parent_name'] == 'F_RESQ_IVT_TBY_1565_DEVCZ10' and x['Protocol ID'] in first_hosp_mapping.keys() and (x['FIRST_ARRIVAL_HOSP'] == 'unknown' or x['FIRST_ARRIVAL_HOSP'] == first_hosp_mapping[x['Protocol ID']])) else x['FIRST_HOSPITAL'], axis=1)

            df['FIRST_HOSPITAL'] = df.apply(lambda x: 2 if (x['crf_parent_name'] == 'F_RESQ_IVT_TBY_1565_DEVCZ10' and x['Protocol ID'] in first_hosp_mapping.keys() and (x['FIRST_ARRIVAL_HOSP'] != 'unknown' or x['FIRST_ARRIVAL_HOSP'] != first_hosp_mapping[x['Protocol ID']])) else x['FIRST_HOSPITAL'], axis=1)

            thrombectomy_df = df[(df['Protocol ID'].isin(self.hospitals_mt)) & (df['TBY_DONE'].isin([1])) & (df['STROKE_TYPE'].isin([1]))].copy()
            thrombectomy_df.fillna(0, inplace=True)
            statistic = self.site_id_mapped_to_site_name.copy()
            
            if thrombectomy_df.empty:
                statistic['# TBY'] = 0
                statistic['Median DTG (minutes)'] = 0
                statistic['# incorrect TBY times'] = 0
                statistic['% incorrect TBY times'] = 0
                statistic['Median DTG (minutes) - first hospital'] = 0
                statistic['Median DTG (minutes) - second hospital'] = 0
            else:
                # Median DTG
                # thrombectomy_df['TBY'] = thrombectomy_df['TBY_ONLY_GROIN_PUNCTURE_TIME'] + thrombectomy_df['TBY_ONLY_GROIN_TIME_MIN'] + thrombectomy_df['IVT_TBY_GROIN_TIME'] + thrombectomy_df['IVT_TBY_GROIN_TIME_MIN']  # get TBY times in one column

                thrombectomy_df['INCORRECT_TIMES'] = False
                thrombectomy_df['INCORRECT_TIMES'] = thrombectomy_df.apply(
                    lambda x: self.get_incorrect_times(x['IVT_TBY_ADMISSION_TIME'], x['IVT_TBY_GROIN_PUNCTURE_TIME'], 700) if x['RECANALIZATION_PROCEDURES'] == 3 and x['IVT_TBY'] == 2 else x['INCORRECT_TIMES'], axis=1)
                thrombectomy_df['INCORRECT_TIMES'] = thrombectomy_df.apply(
                    lambda x: self.get_incorrect_times(x['TBY_ONLY_ADMISSION_TIME'], x['TBY_ONLY_PUNCTURE_TIME'], 700) if x['RECANALIZATION_PROCEDURES'] == 4 and x['TBY_ONLY'] == 2 else x['INCORRECT_TIMES'], axis=1)
                # Add also if tby_refer_all and tby_refer_lim has been selected, but also version of ivt/tby form has to be checked
                thrombectomy_df['INCORRECT_TIMES'] = thrombectomy_df.apply(
                    lambda x: self.get_incorrect_times(x['TBY_REFER_ALL_ADMISSION_TIME'], x['TBY_REFER_ALL_BOLUS_TIME'], 700) if x['RECANALIZATION_PROCEDURES'] == 7 and x['TBY_REFER_ALL'] == 2 and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_CZ_2' else x['INCORRECT_TIMES'], axis=1)
                thrombectomy_df['INCORRECT_TIMES'] = thrombectomy_df.apply(
                    lambda x: self.get_incorrect_times(x['TBY_REFER_LIM_ADMISSION_TIME'], x['TBY_REFER_LIM_BOLUS_TIME'], 700) if x['RECANALIZATION_PROCEDURES'] == 8 and x['TBY_REFER_ALL'] == 2 and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_CZ_2' else x['INCORRECT_TIMES'], axis=1)

                thrombectomy_df['INCORRECT_TIMES'] = thrombectomy_df.apply(
                    lambda x: True if (x['TBY'] <= 0 or x['TBY'] > 700) and x['IVT_TBY'] == 1 else x['INCORRECT_TIMES'], axis=1)
                thrombectomy_df['INCORRECT_TIMES'] = thrombectomy_df.apply(
                    lambda x: True if (x['TBY'] <= 0 or x['TBY'] > 700) and x['TBY_ONLY'] == 1 else x['INCORRECT_TIMES'], axis=1)
                thrombectomy_df['INCORRECT_TIMES'] = thrombectomy_df.apply(
                    lambda x: True if (x['TBY'] <= 0 or x['TBY'] > 700) and x['TBY_REFER_ALL'] == 1 and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_CZ_2' else x['INCORRECT_TIMES'], axis=1)
                thrombectomy_df['INCORRECT_TIMES'] = thrombectomy_df.apply(
                    lambda x: True if (x['TBY'] <= 0 or x['TBY'] > 700) and x['TBY_REFER_ALL'] == 1 and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_1565_DEVCZ10' else x['INCORRECT_TIMES'], axis=1)
                thrombectomy_df['INCORRECT_TIMES'] = thrombectomy_df.apply(
                    lambda x: True if (x['TBY'] <= 0 or x['TBY'] > 700) and x['TBY_REFER_LIM'] == 1 and x['crf_parent_name'] == 'F_RESQ_IVT_TBY_CZ_2' else x['INCORRECT_TIMES'], axis=1)


                #incorrect_tby_times = thrombectomy_df.loc[(thrombectomy_df['INCORRECT_TIMES'] == True) & (~thrombectomy_df['HOSPITAL_STROKE'].isin([1]))].copy()
                # Aug 04, 2020
                incorrect_tby_times = thrombectomy_df.loc[
                    (thrombectomy_df['INCORRECT_TIMES'] == True) & 
                    (thrombectomy_df['HOSPITAL_STROKE_TBY_TIMESTAMPS'] != 1)].copy()

                statistic['Total patients undergone TBY'] = self.count_patients(df=thrombectomy_df, statistic=statistic)
                incorrect_tby_times_save = incorrect_tby_times.loc[incorrect_tby_times['Protocol ID'] != "CZ"].copy()
                incorrect_tby_times_save.to_csv('incorrect_tby_times.csv', sep=',')
                
                #thrombectomy_df.to_csv('thrombectomy_{}.csv'.format(name), sep=',')
                included_in_median = thrombectomy_df[thrombectomy_df['INCLUDE_MEDIAN'] == True].copy()
                included_in_median.to_csv('included_in_median.csv', sep=',')
                thrombectomy = included_in_median[
                    (included_in_median['TBY'] > 0) & 
                    (included_in_median['TBY'] < 700)
                ].copy()

                # Apr 22, 2020 - exclude patients if hospital stroke and times for TBY as timestamps
                thrombectomy = thrombectomy.loc[
                    ~(thrombectomy['HOSPITAL_STROKE_TBY_TIMESTAMPS'].isin([1]))
                ].copy()

                # thrombectomy = thrombectomy.loc[~thrombectomy['HOSPITAL_STROKE'].isin([1])].copy()

                if thrombectomy.empty:
                    statistic['# TBY'] = 0
                    statistic['Median DTG (minutes)'] = 0
                    if incorrect_tby_times.empty:
                        statistic['# incorrect TBY times'] = 0
                        statistic['% incorrect TBY times'] = 0
                    else:
                        statistic['# incorrect TBY times'] = self.count_patients(df=incorrect_tby_times, statistic=statistic)
                        statistic['% incorrect TBY times'] = round((statistic['# incorrect TBY times'] / statistic['Total patients undergone TBY'])*100, 2)
                    statistic['Median DTG (minutes) - first hospital'] = 0
                    statistic['Median DTG (minutes) - second hospital'] = 0
                else:
                    # Total patients
                    # total_patients = thrombectomy.groupby(['Protocol ID']).size().reset_index(name="# TBY")
                    total_patients = thrombectomy_df.groupby(['Protocol ID']).size().reset_index(name="# TBY")
                    statistic = statistic.merge(total_patients, on='Protocol ID', how='outer') # Merge with statistic dataframe
                    statistic.fillna(0, inplace=True)
                    statistic.loc[statistic['Protocol ID'] == 'CZ', '# TBY'] = int(statistics.mean(statistic.loc[statistic['Protocol ID'] != 'CZ']['# TBY'].tolist()))
                    statistic.fillna(0, inplace=True)

                    


                    thrombectomy_grouped = thrombectomy.groupby(['Protocol ID']).TBY.agg(['median']).rename(columns={'median': 'Median DTG (minutes)'}).reset_index()
                    statistic = statistic.merge(thrombectomy_grouped, how='outer') # Merge with statistic dataframe

                    if incorrect_tby_times.empty:
                        statistic['# incorrect TBY times'] = 0
                        statistic['% incorrect TBY times'] = 0
                    else:
                        statistic['# incorrect TBY times'] = self.count_patients(df=incorrect_tby_times, statistic=statistic)
                        statistic['% incorrect TBY times'] = round((statistic['# incorrect TBY times'] / statistic['Total patients undergone TBY'])*100, 2)

                        statistic.loc[statistic['Protocol ID'] == 'CZ', '# incorrect TBY times'] = statistic.loc[statistic['Protocol ID'] != 'CZ']['# incorrect TBY times'].sum(axis=0, skipna=True)
                        statistic.loc[statistic['Protocol ID'] == 'CZ', '% incorrect TBY times'] = round((statistic['# incorrect TBY times'] / statistic['Total patients undergone TBY'])*100, 2)
                        
                    # Median DTG for first hospital arrival
                    thrombectomy_first = thrombectomy[thrombectomy['FIRST_HOSPITAL'] == 1].copy()
                    if thrombectomy_first.empty:
                        statistic['Median DTG (minutes) - first hospital'] = 0
                    else:
                        # thrombectomy_first['TBY'] = thrombectomy_first['TBY_ONLY_GROIN_PUNCTURE_TIME'] + thrombectomy_first['TBY_ONLY_GROIN_TIME_MIN'] + thrombectomy_first['IVT_TBY_GROIN_TIME'] + thrombectomy_first['IVT_TBY_GROIN_TIME_MIN']  # get TBY times in one column
                        thrombectomy_first_grouped = thrombectomy_first.groupby(['Protocol ID']).TBY.agg(['median']).rename(columns={'median': 'Median DTG (minutes) - first hospital'}).reset_index()
                        statistic = statistic.merge(thrombectomy_first_grouped, how='outer') # Merge with statistic dataframe

                    # Median DTG for secondary hospital
                    thrombectomy_second = thrombectomy[thrombectomy['FIRST_HOSPITAL'] == 2].copy()
                    if thrombectomy_second.empty:
                        statistic['Median DTG (minutes) - second hospital'] = 0
                    else:
                        # thrombectomy_second['TBY'] = thrombectomy_second['TBY_ONLY_GROIN_PUNCTURE_TIME'] + thrombectomy_second['TBY_ONLY_GROIN_TIME_MIN'] + thrombectomy_second['IVT_TBY_GROIN_TIME'] + thrombectomy_second['IVT_TBY_GROIN_TIME_MIN']  # get TBY times in one column
                        thrombectomy_second_grouped = thrombectomy_second.groupby(['Protocol ID']).TBY.agg(['median']).rename(columns={'median': 'Median DTG (minutes) - second hospital'}).reset_index()
                        statistic = statistic.merge(thrombectomy_second_grouped, how='outer') # Merge with statistic dataframe
                
                
                statistic.loc[statistic['Protocol ID'] == 'CZ', 'Total patients undergone TBY'] = int(statistics.mean(statistic.loc[statistic['Protocol ID'] != 'CZ']['Total patients undergone TBY'].tolist()))
                


            statistic.fillna(0, inplace=True)

            stats_dfs[name] = statistic
            self.incorrect_tby[name] = incorrect_tby_times_save

        return stats_dfs

    def calculate_statistic_per_region(self):
        """ The function calculating the result statistic for recanalization procedures per regions. """
        stats_dfs = {}

        # Iterate through filtered dataframes
        for name, df in self.filtered_dfs.items():
            

            # Calculate IVtPa median
            ischemic_cmp = df[df['STROKE_TYPE'].isin([1])].copy() 													# only patients with ischemic stroke
            # thrombolysis = ischemic_cmp[ischemic_cmp['RECANALIZATION_PROCEDURES'].isin([2,3,5])].copy() 	# only patients with ischemic stroke who underwent recanalizaiton procedure (IVtPa, IVtPa + TBY, IVtPa + referred for TBY)
            thrombolysis = ischemic_cmp[ischemic_cmp['IVT_DONE'].isin([1])].copy()

            region_total_patients = pd.DataFrame(list(self.regions.keys()), columns=['Site Name'])

            if thrombolysis.empty:
                region_total_patients['Total patients'] = 0
                region_total_patients['# IVT per population'] = 0
            else:
                # Get results per region
                thrombolysis['Site Name'] = thrombolysis.apply(lambda x: self.get_region(x['Protocol ID']), axis=1)
                total_patients =  thrombolysis.groupby(['Site Name']).size().reset_index(name='Total patients')
                region_total_patients = region_total_patients.merge(total_patients, on='Site Name', how='outer')
                region_total_patients.fillna(0, inplace=True)

                region_total_patients['# IVT per population'] = region_total_patients.apply(lambda x: round((x['Total patients']/self.regions[x['Site Name']]['population'])*100000, 2) if x['Total patients'] > 0 else 0, axis=1)

                region_total_patients.loc[region_total_patients['Site Name'] == self.country_name, 'Total patients'] = int(statistics.mean(region_total_patients.loc[region_total_patients['Site Name'] != self.country_name]['Total patients'].tolist()))

            stats_dfs[name] = region_total_patients
        
        return stats_dfs

    def save_excel(self):
        """ The function generating the Excel file with intermediate data used for generating SITS-like reports. The excel file contains three sheets for each period (thrombolysis, thrombectomy, per region). """

        # Create workbook
        output_filename = "SITSlike_reports_stats_" + datetime.now().strftime('%d-%m-%Y') + ".xlsx"
        workbook = xlsxwriter.Workbook(output_filename)
        logging.info('Preprocessed data: The workbook was created.')
        # Create worksheets
        sheets = {}
        
        for name, df in self.thrombolysis_stats_df.items():
            if name == str(self.year):
                sheet_name = "thrombolysis_" + str(self.year)
            else:
                month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                sheet_name = "thrombolysis_" + month_name + "_" + str(self.year)
            
            sheet = workbook.add_worksheet(sheet_name)
            
            values = df.values.tolist()
            nrow = len(df)

            columns = df.columns.tolist()
            ncol = len(df.columns)


            # Create header
            col = []
            for j in range(0, ncol):
                tmp = {}
                tmp['header'] = df.columns.tolist()[j]
                col.append(tmp)
            
            options = {'data': values,
                   'header_row': True,
                   'columns': col,
                   'style': 'Table Style Light 1'
                }
            
            sheet.add_table(0, 0, nrow, ncol - 1, options)
            sheet.set_column(0, 10, 30)
            logging.info('Statistics: {0} sheet was added into excel file!'.format(sheet_name))

        for name, df in self.thrombectomy_stats_df.items():
            if name == str(self.year):
                sheet_name = "thrombectomy_" + str(self.year)
            else:
                month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                sheet_name = "thrombectomy_" + month_name + "_" + str(self.year)
            
            sheet = workbook.add_worksheet(sheet_name)
            
            df.fillna(0, inplace=True)
            values = df.values.tolist()
            nrow = len(df)

            columns = df.columns.tolist()
            ncol = len(df.columns)


            # Create header
            col = []
            for j in range(0, ncol):
                tmp = {}
                tmp['header'] = df.columns.tolist()[j]
                col.append(tmp)
            
            options = {'data': values,
                   'header_row': True,
                   'columns': col,
                   'style': 'Table Style Light 1'
                }
            
            sheet.add_table(0, 0, nrow, ncol - 1, options)
            sheet.set_column(0, 10, 30)
            logging.info('Statistics: {0} sheet was added into excel file!'.format(sheet_name))
        
        for name, df in self.statistic_region_dfs.items():
            if name == str(self.year):
                sheet_name = "region_" + str(self.year)
            else:
                month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                sheet_name = "region_" + month_name + "_" + str(self.year)
            
            sheet = workbook.add_worksheet(sheet_name)
            
            values = df.values.tolist()
            nrow = len(df)

            columns = df.columns.tolist()
            ncol = len(df.columns)


            # Create header
            col = []
            for j in range(0, ncol):
                tmp = {}
                tmp['header'] = df.columns.tolist()[j]
                col.append(tmp)
            
            options = {'data': values,
                   'header_row': True,
                   'columns': col,
                   'style': 'Table Style Light 1'
                }
            
            sheet.add_table(0, 0, nrow, ncol - 1, options)
            sheet.set_column(0, 10, 30)
            logging.info('Statistics: {0} sheet was added into excel file!'.format(sheet_name))
            
        workbook.close()

class GeneratePresentation(Reports):
    """ The class generating graphs in the presentation. """

    def _generate_graphs(self):
        """ The functin generating graphs in the presentation. For each month is generated seperated presentation with the graphs. The last month is included in the cumulative presentation. """
        
        df_names = self.names.copy()
        # Delete last item from list of names (the whole year)
        # del df_names[-1]

        for i in df_names:  
            if i == str(self.year):
                wanted_keys = [i]
                dictfilt = lambda x, y: dict([ (i,x[i]) for i in x if i in set(y) ])

                # master_pptx = self.country_code + ".pptx"
                script_dir = os.path.dirname(__file__) #<-- absolute dir the script is in
                master_pptx = "master.pptx"
                self.master = os.path.normpath(os.path.join(script_dir, "backgrounds", master_pptx))

                # If country is used as site, the country name is selected from countries dictionary by country code. :) 
                '''if self.country == 'UZB':
                    self.country = 'UZ'
                self.country_name = pytz.country_names[self.country]
                '''
                self.country_name = 'Česká republika'

                prs = Presentation(self.master)

                first_slide = prs.slides[0]
                shape = first_slide.shapes[5]
                text_frame = shape.text_frame

                first_slide_text = self.country_name + "\nReports\n"

                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = first_slide_text

                font = run.font
                font.name = 'Century Gothic'
                font.size = Pt(24)
                font.color.rgb = RGBColor(250,250,250)

                text_frame_sub = shape.text_frame

                first_month = datetime(self.year, 1, 1, 0, 0).strftime("%b")
                if self.month != 1:
                    if self.month == 12:
                        last_month = (datetime(self.year, self.month, 31, 0, 0)).strftime("%b")
                    else:
                        last_month = (datetime(self.year, (self.month % 12 + 1), 1, 0, 0) - timedelta(days=1)).strftime("%b")
                else:
                    last_month = ""

                if self.month == 1:
                    text_month = " ".join([first_month, str(self.year)])
                else:
                    text_month = first_month + " - " + last_month + "\n" + str(self.year)

                p = text_frame_sub.paragraphs[0]
                run = p.add_run()
                run.text = text_month

                font = run.font
                font.name = 'Century Gothic'
                font.size = Pt(18)
                font.color.rgb = RGBColor(250,250,250)

                main_col = 'Site Name'
                

                thrombolysis_stats_df = dictfilt(self.thrombolysis_stats_df, wanted_keys)
                statistic_region_dfs = dictfilt(self.statistic_region_dfs, wanted_keys)
                thrombectomy_stats_df = dictfilt(self.thrombectomy_stats_df, wanted_keys)
                # Iterate through dictionaries with statistics
                for name, df in thrombolysis_stats_df.items():

                    # MEDIAN DNT
                    column_name = 'Median DTN (minutes)'
                    axis_title = "Čas [min]"
                    content = ["Parametr medián DOOR-TO-NEEDLE TIME je čas, který odráží kvalitu nemocničního managementu.", "Tento čas musí zahrnovat všechen čas, který uplyne od překročení pacienta prvních dvěří nemocnice až po zahájení léčby."]

                    tmp_df_orig = df[[main_col, column_name]]
                    tmp_df_zeros = tmp_df_orig[tmp_df_orig[column_name] == 0]
                    tmp_df_not_zeros = tmp_df_orig[tmp_df_orig[column_name] != 0].sort_values([column_name], ascending=False)
                    tmp_df = tmp_df_zeros.append(tmp_df_not_zeros, ignore_index=False, sort=False)
                    
                    if name == str(self.year):
                        if last_month == "":
                            title = "Medián door-to-needle time pro intravenózní trombolýzu - " + first_month + " " + str(self.year)
                        else:
                            title = "Medián door-to-needle time pro intravenózní trombolýzu - " + first_month + "-" + last_month + " " + str(self.year)
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "Medián door-to-needle time pro intravenózní trombolýzu - " + month_name + " " + str(self.year)

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, coloring=True, content=content)

                # Iterate through dictionaries with statistics
                for name, df in thrombolysis_stats_df.items():
                    # MEDIAN DGT
                    column_name = 'Total patients undergone IVT'
                    axis_title = 'Počet trombolýz'
                    tmp_df = df[[main_col, column_name]].sort_values([column_name], ascending=True)
                    #tmp_df = tmp_df.loc[tmp_df['Site Name'] != self.country_name]
                    total_pts = round(sum(tmp_df.loc[tmp_df['Site Name'] != self.country_name][column_name].tolist()))
                    
                    if name == str(self.year):
                        if last_month == "":
                            title = "Počet IVT na IC/KCC - {} {} (n={})".format(first_month, self.year, total_pts) 
                        else:
                            title = "Počet IVT na IC/KCC - {}-{} {} (n={})".format(first_month, last_month, self.year, total_pts) 
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "Počet IVT na IC/KCC - {} {} (n={})".format(month_name, self.year, total_pts) 

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title)

                # Iterate through dictionaries with statistics
                for name, df in thrombolysis_stats_df.items():
                    # MEDIAN last seen normal
                    column_name = 'Median last seen normal'
                    axis_title = "Čas [min]"
                    tmp_df_orig = df[[main_col, column_name]]
                    tmp_df_zeros = tmp_df_orig[tmp_df_orig[column_name] == 0]
                    tmp_df_not_zeros = tmp_df_orig[tmp_df_orig[column_name] != 0].sort_values([column_name], ascending=False)
                    tmp_df = tmp_df_zeros.append(tmp_df_not_zeros, ignore_index=False, sort=False)
                    
                    if name == str(self.year):
                        if last_month == "":
                            title = "Medián viděn naposledy zdráv (=začátek symptomů) - příjezd do nemocnice - " + first_month + " " + str(self.year)
                        else:
                            title = "Medián viděn naposledy zdráv (=začátek symptomů) - příjezd do nemocnice - " + first_month + "-" + last_month + " " + str(self.year)
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "Medián viděn naposledy zdráv (=začátek symptomů) - příjezd do nemocnice - " + month_name + " " + str(self.year)

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title)

                # Number of IVT per region
                for name, df in statistic_region_dfs.items():
                    column_name = 'Total patients'
                    tmp_df = df.sort_values([column_name], ascending=True)

                    total_pts = round(sum(tmp_df.loc[tmp_df['Site Name'] != self.country_name][column_name].tolist()))
                    #total_pts = sum(tmp_df[column_name].tolist())

                    if name == str(self.year):
                        if last_month == "":
                            title = "Počet IVT provedených v jednotlivých krajích - {} {} (n={})".format(first_month, self.year, total_pts)
                        else:
                            title = "Počet IVT provedených v jednotlivých krajích - {}-{} {} (n={})".format(first_month, last_month, self.year, total_pts)
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "Počet IVT provedených v jednotlivých krajích - {} {} (n={})".format(month_name, self.year, total_pts)

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, region=True)

                # IVT per population
                for name, df in statistic_region_dfs.items():
                    column_name = '# IVT per population'
                    tmp_df = df.sort_values([column_name], ascending=True)

                    if name == str(self.year):
                        if last_month == "":
                            title = "Počet IVT na 100 000 obyvatel jednotlivých krajů - " + first_month + " " + str(self.year)
                        else:   
                            title = "Počet IVT na 100 000 obyvatel jednotlivých krajů - " + first_month + "-" + last_month + " " + str(self.year)
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "Počet IVT na 100 000 obyvatel jednotlivých krajů - " + month_name + " " + str(self.year)

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, region=True)

                # Iterate through dictionaries with statistics
                for name, df in thrombolysis_stats_df.items():
                    # incorrect times
                    column_name = '% incorrect IVtPa times'
                    axis_title = 'Procento [%]'
                    tmp_df = df[[main_col, column_name]].sort_values([column_name], ascending=False)
                    
                    if name == str(self.year):
                        if last_month == "":
                            title = "% nezadaných nebo chybně zadaných údajů pro DNT - " + first_month + " " + str(self.year)
                        else:
                            title = "% nezadaných nebo chybně zadaných údajů pro DNT - " + first_month + "-" + last_month + " " + str(self.year)
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "% nezadaných nebo chybně zadaných údajů pro DNT - " + month_name + " " + str(self.year)

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, incorrect=True, maximum=100)
                # Iterate through dictionaries with statistics
                
                for name, df in thrombolysis_stats_df.items():
                    # Generate table for incorrect IVTPA
                    df = self.incorrect_ivtpa[name]
                    title = 'Pacienti, kteří mají nesprávně zadané údaje pro výpočet DNT'
                    content = ['Údaj DTN je brán jako nesprávný pokud je čas v minutách: ', 
                        '\ta) menší nebo roven 0 nebo',
                        '\tb) větší než 400.',
                        'Ve většině případů se jedná o chybu, kdy čas léčby předchází čas hospitalizace.',
                        '\n',
                    ]
                    """
                    if not self.development_forms_counts.loc[self.development_forms_counts['n'] > 5].empty:
                        content.append('V případě chybného údaje ve formuláři F_RESQ_IVT_TBY_1565_DEVCZ10 již bohužel není možné udělat změny, jednalo se o testovací verzi a tento formulář se již nepoužívá. Pacienti jsou započítáni do celkového počtu pouze v případě, jestliže celkový počet pacientů v tomto formuláři byl větší než 5.')
                        content.append('Více jak 5 pacientů mají v tomto formuláři tyto nemocnice: {}.'.format(', '.join(self.development_forms_counts.loc[(self.development_forms_counts['n'] > 5) & (self.development_forms_counts['Site Name'] != self.country_name), 'Site Name'].tolist())))
                        content.append('V případě zájmu o upravení a upřesnění dat z této verze nás prosím kontaktuje na qualityregistry@fnusa.cz.')
                        """
                    GenerateTable(df=df, presentation=prs, title=title, content=content)
                    
                

                for name, df in thrombectomy_stats_df.items():
                    # Median DTG
                    column_name = 'Median DTG (minutes)'
                    axis_title = "Čas [min]"
                    content = ["Parametr medián DOOR-TO-GROIN TIME je čas, který odráží kvalitu nemocničního managementu.", "Tento čas musí zahrnovat všechen čas, který uplyne od překročení pacienta prvních dvěří nemocnice až po vpich do třísla."]

                    tmp_df_orig = df[[main_col, column_name]]
                    tmp_df_zeros = tmp_df_orig[tmp_df_orig[column_name] == 0]
                    tmp_df_not_zeros = tmp_df_orig[tmp_df_orig[column_name] != 0].sort_values([column_name], ascending=False)
                    tmp_df = tmp_df_zeros.append(tmp_df_not_zeros, ignore_index=False, sort=False)
                    
                    if name == str(self.year):
                        if last_month == "":
                            title = "Medián door-to-groin time - " + first_month + " " + str(self.year)
                        else:
                            title = "Medián door-to-groin time - " + first_month + "-" + last_month + " " + str(self.year)
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "Medián door-to-groin time - " + month_name + " " + str(self.year)

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, content=content)
                    

                for name, df in thrombectomy_stats_df.items():
                    # Median DTG
                    column_name = 'Median DTG (minutes) - first hospital'
                    axis_title = "Čas [min]"
                    content = ["Parametr medián DOOR-TO-GROIN TIME je čas, který odráží kvalitu nemocničního managementu.", "Tento čas musí zahrnovat všechen čas, který uplyne od překročení pacienta prvních dvěří nemocnice až po vpich do třísla."]

                    tmp_df_orig = df[[main_col, column_name]]
                    tmp_df_zeros = tmp_df_orig[tmp_df_orig[column_name] == 0]
                    tmp_df_not_zeros = tmp_df_orig[tmp_df_orig[column_name] != 0].sort_values([column_name], ascending=False)
                    tmp_df = tmp_df_zeros.append(tmp_df_not_zeros, ignore_index=False, sort=False)
                    
                    if name == str(self.year):
                        if last_month == "":
                            title = "Medián door-to-groin time - Primární příjem k intervenci MT - " + first_month + " " + str(self.year)
                        else:
                            title = "Medián door-to-groin time - Primární příjem k intervenci MT - " + first_month + "-" + last_month + " " + str(self.year)
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "Medián door-to-groin time - Primární příjem k intervenci MT - " + month_name + " " + str(self.year)

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, content=content)

                for name, df in thrombectomy_stats_df.items():
                    # Median DTG
                    column_name = 'Median DTG (minutes) - second hospital'
                    axis_title = "Čas [min]"
                    content = ["Parametr medián DOOR-TO-GROIN TIME je čas, který odráží kvalitu nemocničního managementu.", "Tento čas musí zahrnovat všechen čas, který uplyne od překročení pacienta prvních dvěří nemocnice až po vpich do třísla."]

                    tmp_df_orig = df[[main_col, column_name]]
                    tmp_df_zeros = tmp_df_orig[tmp_df_orig[column_name] == 0]
                    tmp_df_not_zeros = tmp_df_orig[tmp_df_orig[column_name] != 0].sort_values([column_name], ascending=False)
                    tmp_df = tmp_df_zeros.append(tmp_df_not_zeros, ignore_index=False, sort=False)
                    
                    if name == str(self.year):
                        if last_month == "":
                            title = "Medián door-to-groin time - Sekundární příjem k intervenci MT - " + first_month + " " + str(self.year)
                        else:
                            title = "Medián door-to-groin time - Sekundární příjem k intervenci MT - " + first_month + "-" + last_month + " " + str(self.year)
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "Medián door-to-groin time - Sekundární příjem k intervenci MT - " + month_name + " " + str(self.year)

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, content=content)
                
                for name, df in thrombectomy_stats_df.items():
                    # Median DTG
                    column_name = '# TBY'
                    axis_title = 'Počet MT'
                    tmp_df = df[[main_col, column_name]].sort_values([column_name], ascending=True)
                    # tmp_df = tmp_df.loc[tmp_df['Site Name'] != self.country_name]
                    # total_pts = sum(tmp_df[column_name].tolist())
                    total_pts = sum(tmp_df.loc[tmp_df['Site Name'] != self.country_name, column_name].tolist())

                    if name == str(self.year):
                        if last_month == "":
                            title = "Počet MT na nemocnici - {} {} (n={})".format(first_month, self.year, total_pts)
                        else:
                            title = "Počet MT na nemocnici - {}-{} {} (n={})".format(first_month, last_month, self.year, total_pts)
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "Počet MT na nemocnici - {} {} (n={})".format(month_name, self.year, total_pts)

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title)

                for name, df in thrombectomy_stats_df.items():
                    # incorrect times
                    column_name = '% incorrect TBY times'
                    axis_title = 'Procento [%]'
                    tmp_df = df[[main_col, column_name]].sort_values([column_name], ascending=False)
                    
                    if name == str(self.year):
                        if last_month == "":
                            title = "% nezadaných nebo chybně zadaných údajů pro DGT - " + first_month + " " + str(self.year)
                        else:
                            title = "% nezadaných nebo chybně zadaných údajů pro DGT - " + first_month + "-" + last_month + " " + str(self.year)
                    else:
                        month_name = datetime(self.year, name, 1, 0, 0).strftime("%b")
                        title = "% nezadaných nebo chybně zadaných údajů pro DGT - " + month_name + " " + str(self.year)

                    GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, incorrect=True, maximum=100)

                
                for name, df in thrombectomy_stats_df.items():
                    # Save table with incorrect DTG
                    df = self.incorrect_tby[name]
                    title = 'Pacienti, kteří mají nesprávně zadané údaje pro výpočet DTG'
                    content = ['Údaj DTG je brán jako nesprávný pokud je čas v minutách: ', 
                        '\ta) menší nebo roven 0 nebo',
                        '\tb) větší než 700.',
                        'Ve většině případů se jedná o chybu, kdy čas léčby předchází čas hospitalizace.', 
                    ]
                    GenerateTable(df=df, presentation=prs, title=title, content=content)
                

                # set pptx output name (for cz it'll be presentation_CZ.pptx)
                working_dir = os.getcwd()
                if self.month == 1:
                    pptx = str(self.year) + "_RES-Q_report.pptx"
                else:
                    if self.month < 10:
                        month = "0" + str(self.month)
                    else:
                        month = str(self.month)
                    pptx = str(self.year) + "_01_" + month + "_RES-Q_report.pptx"
                #pptx = self.country + "_" + str(self.year) + ".pptx"
                presentation_path = os.path.normpath(os.path.join(working_dir, pptx))

                prs.save(presentation_path)

            else:
                # master_pptx = self.country_code + ".pptx"
                script_dir = os.path.dirname(__file__) #<-- absolute dir the script is in
                master_pptx = "master.pptx"
                self.master = os.path.normpath(os.path.join(script_dir, "backgrounds", master_pptx))

                # If country is used as site, the country name is selected from countries dictionary by country code. :) 
                '''if self.country == 'UZB':
                    self.country = 'UZ'
                self.country_name = pytz.country_names[self.country]'''
                self.country_name = "Česká republika"

                prs = Presentation(self.master)

                first_slide = prs.slides[0]
                shape = first_slide.shapes[5]
                text_frame = shape.text_frame

                first_slide_text = self.country_name + "\nReports\n"

                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = first_slide_text

                font = run.font
                font.name = 'Century Gothic'
                font.size = Pt(24)
                font.color.rgb = RGBColor(250,250,250)

                text_frame_sub = shape.text_frame

                text_month = date(1900, self.month, 1).strftime('%B') + ' ' + str(self.year)

                p = text_frame_sub.paragraphs[0]
                run = p.add_run()
                run.text = text_month

                font = run.font
                font.name = 'Century Gothic'
                font.size = Pt(18)
                font.color.rgb = RGBColor(250,250,250)

                main_col = 'Site Name'
                first_month = datetime(self.year, 1, 1, 0, 0).strftime("%b")
                if (self.month == 12):
                    last_month = (datetime(self.year, self.month, 31, 1, 0, 0)).strftime("%b")
                else:
                    last_month = (datetime(self.year, self.month + 1, 1, 0, 0) - timedelta(days=1)).strftime("%b")

                # Iterate through dictionaries with statistics
                df = self.thrombolysis_stats_df[i]
                # MEDIAN DNT
                column_name = 'Median DTN (minutes)'
                axis_title = "Čas [min]"
                content = ["Parametr medián DOOR-TO-NEEDLE TIME je čas, který odráží kvalitu nemocničního managementu.", "Tento čas musí zahrnovat všechen čas, který uplyne od překročení pacienta prvních dvěří nemocnice až po zahájení léčby."]

                tmp_df_orig = df[[main_col, column_name]]
                tmp_df_zeros = tmp_df_orig[tmp_df_orig[column_name] == 0]
                tmp_df_not_zeros = tmp_df_orig[tmp_df_orig[column_name] != 0].sort_values([column_name], ascending=False)
                tmp_df = tmp_df_zeros.append(tmp_df_not_zeros, ignore_index=False, sort=False)

                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "Medián door-to-needle time pro intravenózní trombolýzu - " + month_name + " " + str(self.year)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, coloring=True, content=content)


                column_name = 'Total patients undergone IVT'
                axis_title = 'Počet trombolýz'
                tmp_df = df[[main_col, column_name]].sort_values([column_name], ascending=True)
                #tmp_df = tmp_df.loc[tmp_df['Site Name'] != self.country_name]
                total_pts = round(sum(tmp_df.loc[tmp_df['Site Name'] != self.country_name][column_name].tolist()))

                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "Počet IVT na IC/KCC - {} {} (n={})".format(month_name, self.year, total_pts)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title)

                # MEDIAN last seen normal
                column_name = 'Median last seen normal'
                axis_title = "Čas [min]"
                tmp_df_orig = df[[main_col, column_name]]
                tmp_df_zeros = tmp_df_orig[tmp_df_orig[column_name] == 0]
                tmp_df_not_zeros = tmp_df_orig[tmp_df_orig[column_name] != 0].sort_values([column_name], ascending=False)
                tmp_df = tmp_df_zeros.append(tmp_df_not_zeros, ignore_index=False, sort=False)

                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "Medián viděn naposledy zdráv (=začátek symptomů) - příjezd do nemocnice - " + month_name + " " + str(self.year)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title)

                # Number of IVT per region
                df = self.statistic_region_dfs[i]
                column_name = 'Total patients'
                tmp_df = df.sort_values([column_name], ascending=True)
                total_pts = round(sum(tmp_df.loc[tmp_df['Site Name'] != self.country_name][column_name].tolist()))
                #total_pts = sum(tmp_df[column_name].tolist())

                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "Počet IVT provedených v jednotlivých krajích - {} {} (n={})".format(month_name, self.year, total_pts)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, region=True)

                # IVT per population
                column_name = '# IVT per population'
                tmp_df = df.sort_values([column_name], ascending=True)

                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "Počet IVT na 100 000 obyvatel jednotlivých krajů - " + month_name + " " + str(self.year)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, region=True)

                # Iterate through dictionaries with statistics
                df = self.thrombolysis_stats_df[i]
                # incorrect times
                column_name = '% incorrect IVtPa times'
                axis_title = 'Procento [%]'
                tmp_df = df[[main_col, column_name]].sort_values([column_name], ascending=False)
   
                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "% nezadaných nebo chybně zadaných údajů pro DNT - " + month_name + " " + str(self.year)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, incorrect=True, maximum=100)

                # Generate table for incorrect IVTPA
                df = self.incorrect_ivtpa[i]
                title = 'Pacienti, kteří mají nesprávně zadané údaje pro výpočet DNT'
                content = ['Údaj DNT je brán jako nesprávný, pokud je čas v minutách: ', 
                    '\ta) menší nebo roven 0 nebo',
                    '\tb) větší než 400.',
                    'Ve většině případů se jedná o chybu, kdy čas léčby předchází čas hospitalizace.', 
                ]
                GenerateTable(df=df, presentation=prs, title=title, content=content)

                # Iterate through dictionaries with statistics
                df = self.thrombectomy_stats_df[i]

                column_name = 'Median DTG (minutes)'
                axis_title = "Čas [min]"
                content = ["Parametr medián DOOR-TO-GROIN TIME je čas, který odráží kvalitu nemocničního managementu.", "Tento čas musí zahrnovat všechen čas, který uplyne od překročení pacienta prvních dvěří nemocnice až po vpich do třísla."]

                tmp_df_orig = df[[main_col, column_name]]
                tmp_df_zeros = tmp_df_orig[tmp_df_orig[column_name] == 0]
                tmp_df_not_zeros = tmp_df_orig[tmp_df_orig[column_name] != 0].sort_values([column_name], ascending=False)
                tmp_df = tmp_df_zeros.append(tmp_df_not_zeros, ignore_index=False, sort=False)

                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "Medián door-to-groin time - " + month_name + " " + str(self.year)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, content=content)
                
                # Median DTG
                column_name = 'Median DTG (minutes) - first hospital'
                axis_title = "Čas [min]"
                content = ["Parametr medián DOOR-TO-GROIN TIME je čas, který odráží kvalitu nemocničního managementu.", "Tento čas musí zahrnovat všechen čas, který uplyne od překročení pacienta prvních dvěří nemocnice až po vpich do třísla."]

                tmp_df_orig = df[[main_col, column_name]]
                tmp_df_zeros = tmp_df_orig[tmp_df_orig[column_name] == 0]
                tmp_df_not_zeros = tmp_df_orig[tmp_df_orig[column_name] != 0].sort_values([column_name], ascending=False)
                tmp_df = tmp_df_zeros.append(tmp_df_not_zeros, ignore_index=False, sort=False)
                
                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "Medián door-to-groin time - Primární příjem k intervenci MT - " + month_name + " " + str(self.year)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, content=content)

                # Median DTG
                column_name = 'Median DTG (minutes) - second hospital'
                axis_title = "Čas [min]"
                content = ["Parametr medián DOOR-TO-GROIN TIME je čas, který odráží kvalitu nemocničního managementu.", "Tento čas musí zahrnovat všechen čas, který uplyne od překročení pacienta prvních dvěří nemocnice až po vpich do třísla."]

                tmp_df_orig = df[[main_col, column_name]]
                tmp_df_zeros = tmp_df_orig[tmp_df_orig[column_name] == 0]
                tmp_df_not_zeros = tmp_df_orig[tmp_df_orig[column_name] != 0].sort_values([column_name], ascending=False)
                tmp_df = tmp_df_zeros.append(tmp_df_not_zeros, ignore_index=False, sort=False)

                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "Medián door-to-groin time - Sekundární příjem k intervenci MT - " + month_name + " " + str(self.year)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, content=content)

                # Median DTG
                column_name = '# TBY'
                axis_title = 'Počet MT'
                tmp_df = df[[main_col, column_name]].sort_values([column_name], ascending=True)
                # tmp_df = tmp_df.loc[tmp_df['Site Name'] != self.country_name]
                # total_pts = sum(tmp_df[column_name].tolist())
                total_pts = sum(tmp_df.loc[tmp_df['Site Name'] != self.country_name, column_name].tolist())
                
                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "Počet MT na nemocnici - {} {} (n={})".format(month_name, self.year, total_pts)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title)

                # incorrect times
                column_name = '% incorrect TBY times'
                axis_title = 'Procento [%]'
                tmp_df = df[[main_col, column_name]].sort_values([column_name], ascending=False)

                month_name = datetime(self.year, i, 1, 0, 0).strftime("%b")
                title = "% nezadaných nebo chybně zadaných údajů pro DGT - " + month_name + " " + str(self.year)

                GenerateGraphs(df=tmp_df, presentation=prs, title=title, column_name=column_name, country_name=self.country_name, axis_name=axis_title, incorrect=True, maximum=100)

                # Save table with incorrect DTG
                df = self.incorrect_tby[i]
                title = 'Pacienti, kteří mají nesprávně zadané údaje pro výpočet DTG'
                content = ['Údaj DTG je brán jako nesprávný, pokud je čas v minutách: ', 
                    '\ta) menší nebo roven 0 nebo',
                    '\tb) větší než 700.',
                    'Ve většině případů se jedná o chybu, kdy čas léčby předchází čas hospitalizace.', 
                ]
                GenerateTable(df=df, presentation=prs, title=title, content=content)

                # set pptx output name (for cz it'll be presentation_CZ.pptx)
                working_dir = os.getcwd()
                if self.month < 10:
                    month = '0' + str(self.month)
                else:
                    month = str(self.month)
                pptx = str(self.year) + "_" + month + "_RES-Q_report.pptx"
                #pptx = self.country + "_" + month_name + ".pptx"
                presentation_path = os.path.normpath(os.path.join(working_dir, pptx))

                prs.save(presentation_path)  

            
    def generate_presentation(self):
        """ The function calling the :func:`self._generate_graphs`. """

        self._generate_graphs()


class GenerateGraphs:
    """ The class generating graphs into presentation and called inside the :class:`resqdb.Reports.GeneratePresentation`. 

    :param df: the dataframe with calculated statistics
    :type df: pandas dataframe
    :param presentation: the opened document (pptx)
    :type presentation: Presentation object
    :param title: the title of the slide
    :type title: str
    :param column_name: the name of column which should be used in the graph (for stacked graph represent the first column to get index where the data included in the graph starts)
    :type column_name: str
    :param country_name: the country name
    :type coutnry_name: str
    :param axis_name: the label of x-axis
    :type axis_name: str
    :param coloring: `True` if rows should be colored by number, else `False`
    :type coloring: bool
    :param region: `True` if region graphs should be generated (coloring issue)
    :type region: bool
    :param incorrect: `True` if graphs displaying the incorrect times are generated
    :type incorrect: bool
    :param maximum: maximum value of x-axis for some graph
    :type maximum: int
    :param content: the small guide text displayed on the slide next to graph, each paragraphs is new value in list
    :type content: list 
    """
    def __init__(self, df, presentation, title, column_name, country_name, axis_name=None, coloring=False, region=False, incorrect=False, maximum=0, content=None):

        self.dataframe = df
        self.presentation = presentation
        self.title = title
        self.column_name = column_name
        self.font_name = 'Century Gothic'
        self.categories_column = 'Site Name'
        self.country_name = country_name
        self.coloring = coloring
        self.region = region
        self.incorrect = incorrect
        self.maximum = maximum
        self.axis_name = axis_name
        self.content = content

        # Estimate font sizes based on number of sites included in the graph
        if (len(self.dataframe) > 15):
            self.category_font_size = Pt(10)
            self.data_label_font_size = Pt(8)
        else:
            self.category_font_size = Pt(11)
            self.data_label_font_size = Pt(11)

        self._create_barplot()

    def _set_transparency(self, transparency, elm):
        """ The function set the transparency of the row. 

        :param transparency: the transparency in %
        :type transparency: int
        :param elm: the element which transparency should be changed
        :type elm: format.line.color._xFill
        """
        a = str(100 - transparency) + '196'
        
        alpha = OxmlElement('a:alpha')
        alpha.set('val', a)
        elm.srgbClr.append(alpha)

    def _create_barplot(self):
        """ The function creating the new graph into the presentation based on the graph type. """

        colors = {
            'yellow': RGBColor(255, 192, 0), 
            'green': RGBColor(98, 153, 62), 
            'crimsom': RGBColor(220, 20, 60), 
            'blue': RGBColor(43, 88, 173),
            'wine_red': RGBColor(134, 0, 0)
        }

        site_names = self.dataframe[self.categories_column].tolist()
        values = self.dataframe[self.column_name].tolist()

        # Add slide to presentation (layout 11 is our custom layout where only title 'Agency FB', color: RGBColor(43, 88, 173)  and size:24 is set)
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        # Get title object
        title_placeholders = slide.shapes.title
        # Set title
        title_placeholders.text = self.title

        # Add textbox explanation
        if self.content is not None:
            len_df = len(self.dataframe[self.dataframe[self.column_name] > 0])
            if len_df < 12:
                left = Cm(24.7)
                top = Cm(12)
                width = Cm(8)
                height = Cm(5)
            else:
                left = Cm(24.7)
                top = Cm(2)
                width = Cm(8)
                height = Cm(5)

            # Add textbox with explanation
            txBox = slide.shapes.add_textbox(left, top, width, height)
            txBox.text_frame.clear()
            txBox.text_frame.word_wrap = True
            for i in range(0, len(self.content)):
                if i == 0:
                    p = txBox.text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = self.content[i]
                else:
                    p = txBox.text_frame.add_paragraph()
                    run = p.add_run()
                    run.text = self.content[i]
            
            for paragraph in txBox.text_frame.paragraphs:
                paragraph.line_spacing = Pt(18)
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(10.5)
                    run.font.name = self.font_name

        chart_data = ChartData()
        chart_data.categories = self.dataframe[self.categories_column].tolist()
        chart_data.add_series(self.column_name, self.dataframe[self.column_name].tolist())

        # Add chart on slide
        specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
            }
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        # Get series of chart
        series = chart.series[0]

       
        if self.coloring:
            # Coloring for median values - <= 20 green, > 20 and <= 30 yellow, else crimsom
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                value = values[idx]
                
                if (site_names[idx] == self.country_name):
                    fill.fore_color.rgb = colors['wine_red']                    
                elif (value > 0 and value <= 20):
                    fill.fore_color.rgb = colors['green']
                elif (value > 20 and value <= 30):
                    fill.fore_color.rgb = colors['yellow']
                else:
                    fill.fore_color.rgb = colors['crimsom']
        elif self.region:
            # The lowest value colored red, the biggest value colored green
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                if idx == values.count(0):
                    fill.fore_color.rgb = colors['crimsom']
                elif (site_names[idx] == self.country_name):
                    fill.fore_color.rgb = colors['yellow']  
                elif idx == (len(values) - 1):
                    fill.fore_color.rgb = colors['green']
                else:
                    fill.fore_color.rgb = colors['blue']
        elif self.incorrect:
            # Set red color for incorrect values
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                if (site_names[idx] == self.country_name):
                    fill.fore_color.rgb = colors['wine_red']   
                else:
                    fill.fore_color.rgb = colors['crimsom']
        else:
            # Blue color for the remaining values 
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                if (site_names[idx] == self.country_name):
                    fill.fore_color.rgb = colors['wine_red']     
                else:
                    fill.fore_color.rgb = colors['blue']

        # Get plot 
        plot = chart.plots[0]
        # Set for each bar same color
        plot.vary_by_categories = False
        # Show data labels and set font
        plot.has_data_labels = True
        # Change gap width
        plot.gap_width = 100

        
        data_labels = plot.data_labels
        data_labels.font.size = self.data_label_font_size
        data_labels.font.bold = True
        data_labels.font.name = self.font_name

        # Value for x-axis (change font size, name, and other things)
        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = self.category_font_size
        tick_labels.font.name = self.font_name
        
        # Don't show major gridlines
        value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
        value_axis.has_major_gridlines = False
        # Set range of axis
        if self.maximum != 0:
            value_axis.maximum_scale = self.maximum
        value_axis.minimum_scale = 0

        if self.axis_name is not None:
            value_axis.has_title = True
            value_axis.axis_title.text_frame.text = self.axis_name
            for paragraph in value_axis.axis_title.text_frame.paragraphs:
                paragraph.font.size = self.category_font_size
                paragraph.font.name = self.font_name

        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
        solidFill = category_axis.format.line.color._xFill
        self._set_transparency(100, solidFill)

        # Delete tick marks
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.major_unit = 1
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name


class GenerateTable:

    """ The class creates table on the slide in the presentation.

        :param df: the dataframe with calculated statistics
        :type df: pandas dataframe
        :param presentation: the opened document (pptx)
        :type presentation: Presentation object
        :param title: the title of the slide
        :type title: str
        """
    def __init__(self, df, presentation, title, content):

        df = df.loc[df['crf_parent_name'] != 'F_RESQ_IVT_TBY_1565_DEVCZ10']
        self.dataframe = df[['Site Name', 'Subject ID']].copy().sort_values(['Site Name'], ascending=True).reset_index(drop=True)
        self.presentation = presentation
        self.title = title
        self.content = content

        parts = len(self.dataframe)//25 # number of parts to which df should be split
        modulo = len(self.dataframe)%25 # if modulo is not zero, then there will be one more slide
        if modulo != 0:
            parts = parts + 1

        for i in range(0, parts):
            if (i + 1) == parts:
                start = i * 25
                df = self.dataframe[start:].copy().reset_index(drop=True)
            else:
                start = i * 25
                end = (i + 1) * 25
                df = self.dataframe[start:end].copy().reset_index(drop=True)

            if (i % 2) != 0:
                self._create_table(df, new_slide=False, left=False)
            else:
                self.slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
                self._create_table(df)

    @property
    def slide(self):
        return self._slide

    @slide.setter
    def slide(self, value):
        self._slide = value

    def _create_table(self, df, new_slide=True, left=True):
        """ The function generating new table in the presentation. 
        
        :param df: the dataframe that should be shown in the table
        :type df: DataFrame
        :param new_slide: True if new slide should be created else False (default is True)
        :type new_slide: boolean
        :param left: True if the table should be shown on the left else False (default is True)
        :type left: boolean
        """      
        if new_slide:
            title_placeholders = self.slide.shapes.title
            title_placeholders.text = self.title

            left = Cm(24)
            top = Cm(2)
            width = Cm(8)
            height = Cm(4)

            # Add textbox with explanation
            txBox = self.slide.shapes.add_textbox(left, top, width, height)
            txBox.text_frame.clear()
            txBox.text_frame.word_wrap = True
            
            for i in range(0, len(self.content)):
                if i == 0:
                    p = txBox.text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = self.content[i]
                else:
                    p = txBox.text_frame.add_paragraph()
                    run = p.add_run()
                    run.text = self.content[i]
            
            for paragraph in txBox.text_frame.paragraphs:
                paragraph.line_spacing = Pt(18)
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(10)

        if left:
            # table_placeholder = slide.shapes[1]
            left, top, width, height = Cm(1), Cm(2), Cm(11), Cm(1)
            shape = self.slide.shapes.add_table(len(df) + 1, len(df.columns), left, top, width, height)
        else:
            left, top, width, height = Cm(12.5), Cm(2), Cm(11), Cm(1)
            shape = self.slide.shapes.add_table(len(df) + 1, len(df.columns), left, top, width, height)
        
        # set table look
        # Change table style (https://github.com/scanny/python-pptx/issues/27)
        style_id = '{7DF18680-E054-41AD-8BC1-D1AEF772440D}'
        tbl = shape._element.graphic.graphicData.tbl
        tbl[0][-1].text = style_id

        table = shape.table
        # table = shape.table

        # Set column names
        columns = ['Nemocnice', 'Subject ID']
        for i in range(0, len(columns)):
            cell = table.cell(0, i) # Get cell in the first row
            cell.text = columns[i]
            for paragraphs in cell.text_frame.paragraphs:
                for run in paragraphs.runs:
                    run.font.size = Pt(8)
        for index, row in df.iterrows():
            for i in range(0, len(row)):
                cell = table.cell(index + 1, i)
                cell.text = str(row[i])
                for paragraphs in cell.text_frame.paragraphs:
                    for run in paragraphs.runs:
                        run.font.size = Pt(8)


            


