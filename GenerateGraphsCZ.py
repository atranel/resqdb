# -*- coding: utf-8 -*-
"""
File name: GenerateGraphs.py
Package: resq
Written by: Marie Jankujova - jankujova.marie@fnusa.cz on 11-2017
Version: v1.0
Version comment: The first version of script which is used to generate charts in the PPTX format. 
Description: This script is used to generate PPTX file from the dataframe with computed statistics. 
This script is used by class GeneratePresentation.py. 
As a result, the presentation with chart is produced. 
"""


import pandas as pd
import sys
import os
from datetime import datetime, date
import sqlite3

import xlsxwriter
import csv

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE

class GenerateGraphs:
    """This class is used to generate our typical presentation with graphs. 

    Arguments:
        dataframe - dataframe with calculated statistics
        presentation - opened pptx document
        title - name of slide
        column_name - name of column name which should be used in graph (for more columns, the first column from all of them)
        graph_type - set which graph should be generated (normal, stacked or grouped) (default = normal)
        number_of_series - set number of series of graphs - this value is equal to length of legend (default = 0)
        legend - list of legend names (default = None)
        country - name of country (if dataset was filtered based on country) (default = None)
    """
    def __init__(self, dataframe, presentation, title, column_name, graph_type = None, number_of_series=0, legend=None, country=None):

        self.dataframe = dataframe
        self.presentation = presentation
        self.title = title
        self.column_name = column_name
        self.number_of_series = number_of_series
        self.legend = legend
        self.country_name = country
        self.font_name = 'Roboto'

        # if (self.country_name in ['Ukraine', 'Poland'] and len(self.dataframe) > 2):
        #     self.categories_column = 'Site ID'
        # else:
        self.categories_column = 'Site Name'
            
        if (len(self.dataframe) > 60):
            self.category_font_size = Pt(6)
            self.data_label_font_size = Pt(6)
        elif (len(self.dataframe) > 50 and len(self.dataframe) <= 60):
            self.category_font_size = Pt(8)
            self.data_label_font_size = Pt(8)
        else:
            self.category_font_size = Pt(10)
            self.data_label_font_size = Pt(11)

        if (graph_type == 'stacked'):
            self._create_stacked_barplot(dataframe=self.dataframe, title=self.title, column_name=self.column_name, legend=self.legend, number_of_series=self.number_of_series)
        else:
            self._create_barplot(dataframe=self.dataframe, title=self.title, column_name=self.column_name)


    def _get_length_of_legend(self, legend):
        count = 0

        for i in legend:
            count = count + len(i)
        
        return count

    def _create_barplot(self, dataframe, title, column_name):
        """Create normal barplot

        Arguments:
            dataframe - dataframe with statistics
            title - title of slide
            column_name - name of column which is included in graph
        """
        maximum = 0

        # If graph is in %, set maximum valut to 100. 
        if '%' in title.lower():
            maximum = 100
            values = [round(x, 0) for x in dataframe[column_name].tolist()]
        else:
            maximum = round((max(dataframe[column_name].tolist())),1)
            values = dataframe[column_name].tolist()
            

        # Add slide to presentation (layout 11 is our custom layout where only title 'Agency FB', color: RGBColor(43, 88, 173)  and size:24 is set)
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        # Get title object
        title_placeholders = slide.shapes.title
        # Set title
        title_placeholders.text = title.upper()

        chart_data = ChartData()
        chart_data.categories = dataframe[self.categories_column].tolist()
        chart_data.add_series(column_name, values)

        # Add chart on slide
        specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
            }
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        # Get series of chart
        series = chart.series[0]
        # If graphs for whole country are generated, set for bar with country with red color
        # else set to blue color (same color as title uses)
        site_names = dataframe[self.categories_column].tolist()
        if (len(dataframe) > 2):
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                if (site_names[idx] == self.country_name):
                    fill.fore_color.rgb = RGBColor(128,0,0)
                else:
                    fill.fore_color.rgb = RGBColor(43, 88, 173)
        else:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(43, 88, 173) 

        # Get plot 
        plot = chart.plots[0]
        # Set for each bar same color
        plot.vary_by_categories = False
        # Show data labels and set font
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = self.data_label_font_size
        data_labels.font.bold = True
        data_labels.font.name = self.font_name

        if 'Total Patients' in column_name or 'Median patient age' in column_name:
            value_axis = chart.value_axis
            value_axis.visible = False
            value_axis.has_major_gridlines = False
        else:
            # Value for x-axis (change font size, name, and other things)
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = self.category_font_size
            tick_labels.font.name = self.font_name
            
            # Don't show major gridlines
            value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
            value_axis.has_major_gridlines = False
            # Set range of axis
            value_axis.maximum_scale = maximum
            value_axis.minimum_scale = 0
            
       
        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        # Delete tick marks
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.major_unit = 1
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name

    def _create_stacked_barplot(self, dataframe, title, column_name, legend, number_of_series):
        """Create stacked barplot

        Arguments:
            dataframe - dataframe with statistics
            title - title of slide
            column_name - name of column (name of fist column used for graph)
            legen - list of legend names
            number_of_series - number of columns included in graph
        """

        # Calculate length of legend (in case that legend is too long, make smaller font size)
        count = self._get_length_of_legend(legend)

        # Get column names of dataframe
        column_names = dataframe.columns.tolist()

        index = column_names.index(column_name)

        # Add new slide into presentation
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        title_placeholders = slide.shapes.title
        title_placeholders.text = title.upper()


        chart_data = ChartData()
        
            
        chart_data.categories = dataframe[self.categories_column].tolist()
        # Add data in each category 
        chart_data.add_series(legend[0], dataframe[column_names[index]].tolist()) 
        if (number_of_series >= 2):
            chart_data.add_series(legend[1], dataframe[column_names[index+1]].tolist())
        if (number_of_series >= 3):
            chart_data.add_series(legend[2], dataframe[column_names[index+2]].tolist())
        if (number_of_series >= 4):
            chart_data.add_series(legend[3], dataframe[column_names[index+3]].tolist())
        if (number_of_series >= 5):
            chart_data.add_series(legend[4], dataframe[column_names[index+4]].tolist())
        if (number_of_series >= 6):
            chart_data.add_series(legend[5], dataframe[column_names[index+5]].tolist())
        if (number_of_series >= 7):
            chart_data.add_series(legend[6], dataframe[column_names[index+6]].tolist())
        if (number_of_series >= 8):
            chart_data.add_series(legend[7], dataframe[column_names[index+7]].tolist())

        # add chart to slide --------------------
        specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
            }
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        series = chart.series[0]
        # If graphs for whole country are generated, set for bar with country with red color
        # else set to blue color (same color as title uses)
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(43, 88, 173)         

        if (number_of_series >= 5):
            series = chart.series[4]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(80, 137, 188)  
          
        # Value for x-axis (change font size, name, and other things)
        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = Pt(11)
        tick_labels.font.name = self.font_name

        value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
        
        if (len(dataframe) > 2):
            value_axis.has_major_gridlines = True
            value_axis.major_gridlines.format.line.dash_style = MSO_LINE.DASH
            value_axis.major_gridlines.format.line.width = Pt(0.5)
        else:
            value_axis.has_major_gridlines = False
        value_axis.maximum_scale = 100
        value_axis.minimum_scale = 0

        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name
        category_labels.tickLblSkip = 1

        # Set legend 
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.name = self.font_name

        if (count > 180 or 'antithrombotics prescribed' in title.lower()):
            chart.legend.font.size = Pt(11)
        else:
            chart.legend.font.size = Pt(12)
