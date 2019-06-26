# -*- coding: utf-8 -*-
"""
File name: GenerateComparisonPresentation.py
Package: resq
Written by: Marie Jankujova - jankujova.marie@fnusa.cz on 11-2018
Version: v1.0
Version comment: The first version of script which generate comaparison into presentation in PPTX format.  
Description: This script is used to generate comparison charts into presentation. 
"""


import pandas as pd
import sys
import os
import sqlite3
from datetime import datetime, date
import math
import pytz

import xlsxwriter


import csv

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.oxml.xmlchemy import OxmlElement

class GeneratePeriodCompPresentation:
    """ The class generating comparison graphs for nationally samples between two periods of times. 

    :param ndf1: the dataframe containing calculated statistics for the first period
    :type ndf1: pandas dataframe
    :param ndf2: the dataframe containing calculated statistics for the second period
    :type ndf2: pandas dataframe
    :param name1: the name of the 1st period, eg. 2017
    :type name1: str
    :param name2: the name of the 2nd period, eg. 2017
    :type name2: str
    """

    def __init__(self, ndf1, ndf2, name1, name2):

        self.ndf1 = ndf1
        self.ndf2 = ndf2
        self.name1 = name1
        self.name2 = name2

        # Get absolute path to the database.
        script_dir = os.path.dirname(__file__) 
        filename = "resq.db"
        self.abs_db_path = os.path.join(script_dir, "database", filename)

        master_pptx = "countries_comparison.pptx"
        self.master = os.path.normpath(os.path.join(script_dir, "backgrounds", master_pptx))

        self._generate_graphs()
    
    def _generate_graphs(self):
        """ The function generating graphs into the presentation. The final name of the presentation is comparison_two_periods.pptx. """
        prs = Presentation(self.master) # Read template presentation

        first_slide = prs.slides[0]
        shape = first_slide.shapes[5]
        text_frame = shape.text_frame

        first_slide_text = "Data Comparison" # Set title
    
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = first_slide_text

        font = run.font
        font.name = 'Calibri Light'
        font.size = Pt(26)
        font.color.rgb = RGBColor(250,250,250)

        # if (self.country_name in ['Ukraine', 'Poland'] and len(df) > 2):
        #     main_col = 'Site ID'
        # else:
        main_col = 'Site Name'

        ########################
        #### TOTAL PATIENTS ####
        ########################
        column_name = 'Total Patients'
        # Nationally df
        tmp_ndf1 = self.ndf1[[main_col, column_name]]
        tmp_ndf1 = tmp_ndf1.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_ndf2 = self.ndf2[[main_col, column_name]]
        tmp_ndf2 = tmp_ndf2.sort_values([column_name], ascending=True)
        # Merge dataframes
        tmp_df = pd.merge(tmp_ndf1, tmp_ndf2, how='right', on=['Site Name'])
        tmp_df.fillna(0, inplace=True)
        tmp_df = tmp_df.sort_values([main_col], ascending=True)
        
        title = "Total patients enrolled"

        legend = [self.name1, self.name2]
        
        GeneratePeriodCompGraph(df=tmp_df, presentation=prs, column_name=column_name, title=title, number_of_series=len(legend), legend=legend)

        ###########################################
        #### PATIENTS RECEIVING RECANALIZATION ####
        ###########################################
        column_name = '% patients recanalized'
        # Nationally df
        tmp_ndf1 = self.ndf1[[main_col, column_name]]
        tmp_ndf1 = tmp_ndf1.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_ndf2 = self.ndf2[[main_col, column_name]]
        tmp_ndf2 = tmp_ndf2.sort_values([column_name], ascending=True)
        # Merge dataframes
        tmp_df = pd.merge(tmp_ndf1, tmp_ndf2, how='right', on=['Site Name'])
        tmp_df.fillna(0, inplace=True)
        tmp_df = tmp_df.sort_values([main_col], ascending=True)

        title = "% patients receiving recanalization procedures"
        subtitle = "Calculated out of number of IS"

        legend = [self.name1, self.name2]
        
        GeneratePeriodCompGraph(df=tmp_df, presentation=prs, column_name=column_name, title=title, subtitle=subtitle, number_of_series=len(legend), legend=legend)

        ##############################
        #### MEDIAN DTN (MINUTES) ####
        ##############################
        column_name = 'Median DTN (minutes)'
        # Nationally df
        tmp_ndf1 = self.ndf1[[main_col, column_name]]
        tmp_ndf1 = tmp_ndf1.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_ndf2 = self.ndf2[[main_col, column_name]]
        tmp_ndf2 = tmp_ndf2.sort_values([column_name], ascending=True)
        # Merge dataframes
        tmp_df = pd.merge(tmp_ndf1, tmp_ndf2, how='right', on=['Site Name'])
        tmp_df.fillna(0, inplace=True)
        tmp_df = tmp_df.sort_values([main_col], ascending=True)

        title = "Median Door-to-Needle Time (DTN), in minutes"

        legend = [self.name1, self.name2]
        
        GeneratePeriodCompGraph(df=tmp_df, presentation=prs, column_name=column_name, title=title, number_of_series=len(legend), legend=legend)

        ##############################
        #### MEDIAN DTG (MINUTES) ####
        ##############################
        column_name = 'Median DTG (minutes)'
        # Nationally df
        tmp_ndf1 = self.ndf1[[main_col, column_name]]
        tmp_ndf1 = tmp_ndf1.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_ndf2 = self.ndf2[[main_col, column_name]]
        tmp_ndf2 = tmp_ndf2.sort_values([column_name], ascending=True)
        # Merge dataframes
        tmp_df = pd.merge(tmp_ndf1, tmp_ndf2, how='right', on=['Site Name'])
        tmp_df.fillna(0, inplace=True)
        tmp_df = tmp_df.sort_values([main_col], ascending=True)

        title = "Median Door-to-Needle Time (DTG), in minutes"

        legend = [self.name1, self.name2]
        
        GeneratePeriodCompGraph(df=tmp_df, presentation=prs, column_name=column_name, title=title, number_of_series=len(legend), legend=legend)
        

        # set pptx output name (for cz it'll be presentation_CZ.pptx)
        working_dir = os.getcwd()
        pptx = "comparison_two_periods.pptx"
        presentation_path = os.path.normpath(os.path.join(working_dir, pptx))

        prs.save(presentation_path)


class GeneratePeriodCompGraph:
    """ The class generating comparison graphs for given periods . 

    :param df: the temporary dataframe created in :class:`resqdb.GenerateComparisonPresentation.GeneratePeriodCompPresentation` class
    :type df: pandas dataframe
    :param presentation: the presentation opened in :class:`resqdb.GenerateComparisonPresentation.GeneratePeriodCompPresentation` class
    :type presentation: Presentation object
    :param column_name: the name of column which data should be shown in the graph
    :type column_name: str
    :param title: the title of the graph
    :type title: str
    :param subtitle: the subtitle of the graph
    :type subtitle: str
    :param number_of_series: the number of columns to be shown (stacked graphs)
    :type number_of_series: int 
    :param legend: the legend if the graph is stacked
    :type legend: list of strings
    """

    def __init__(self, df, presentation, column_name, title, subtitle="", number_of_series=0, legend=None):

        self.df = df        
        self.presentation = presentation
        self.title = title
        self.column_name = column_name
        self.number_of_series = number_of_series
        self.font_name = 'Roboto'
        self.legend = legend
        self.subtitle = subtitle

        self.category_font_size = Pt(10)
        self.data_label_font_size = Pt(11)
        self.categories_column = 'Site Name'

        self._create_column_clustered_barplot()
        

    def _get_length_of_legend(self, legend):
        """ The function adjusting the number of letters in legend to quess the number of columns in the legend! 
        
        :param legend: the names of legend
        :type legend: list
        :returns: the adjusted number of letters
        """
        count = 0

        for i in legend:
            count = count + len(i)
        
        return count

    def _create_column_clustered_barplot(self):
        """ The function creating the clustered barplot. """
        maximum = 0
        
        column_names = self.df.columns.tolist()
        index = column_names.index(self.categories_column)

        # Add slide to presentation (layout 11 is our custom layout where only title 'Agency FB', color: RGBColor(43, 88, 173)  and size:24 is set)
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        # Get title object
        title_placeholders = slide.shapes.title
        title_placeholders.text = self.title

        # If subtitle is not set, remove placeholder for the subtitle from page. 
        if self.subtitle == "":
            subtitle = slide.placeholders[1]
            sp = subtitle.element
            sp.getparent().remove(sp)
        else:
            subtitle = slide.placeholders[1]
            subtitle.text = self.subtitle

        # 1st chart (left side) - nationally sample
        chart_data = ChartData()
        chart_data.categories = self.df[self.categories_column].tolist()
        # Add data in each category 
        chart_data.add_series(self.legend[0], self.df[column_names[index+1]].tolist()) 
        if (self.number_of_series >= 2):
            chart_data.add_series(self.legend[1], self.df[column_names[index+2]].tolist())

        # Set margins.
        specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
        }

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        # Get series of chart
        series = chart.series[0]

        # Get plot 
        plot = chart.plots[0]
        plot.gap_width = 220
        plot.overlap = -25
        # Set for each bar same color
        plot.vary_by_categories = False

        # Value for x-axis (change font size, name, and other things)
        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = self.category_font_size
        tick_labels.font.name = self.font_name

        # Don't show major gridlines
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
        value_axis.major_gridlines.format.line.width = Pt(0.5)
        value_axis.major_tick_mark = XL_TICK_MARK.NONE
        value_axis.format.line.color.rgb = RGBColor(217, 217, 217)
        # Set range of axis
        #value_axis.maximum_scale = ndf_maximum
        value_axis.minimum_scale = 0

        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        category_axis.format.line.color.rgb = RGBColor(217, 217, 217)
        # Delete tick marks
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        #category_axis.major_unit = 1
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name

        # Set legend 
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = self.font_name


class GenerateCountriesCompPresentation:     
    """ The class creating presentation with the comparison between national samples and site samples in one period. 

    :param ndf: the calculated statistics for the national samples
    :type ndf: pandas dataframe
    :param sldf: the calculated statistics for the site samples
    :type sldf: pandas dataframe
    :param name: the name of the period, eg. 2017
    :type name: str
    :param samples: the list of countries which should be colored by different color in the main graphs
    :type samples: list of string
    
    """
    def __init__(self, ndf, sldf, name="", samples=[]):

        self.ndf = ndf
        self.sldf = sldf
        self.name = name

        # Get absolute path to the database.
        script_dir = os.path.dirname(__file__) #<-- absolute dir the script is in
        filename = "resq.db"
        self.abs_db_path = os.path.join(script_dir, "database", filename)

        master_pptx = "countries_comparison.pptx"
        self.master = os.path.normpath(os.path.join(script_dir, "backgrounds", master_pptx))

        def select_country(value):
            """ The function obtaining the country name from the package pytz based on the country code.

            :param value: the country code
            :type value: str
            :returns: the country name
            """
            if value == "UZB":
                value = 'UZ'
            country_name = pytz.country_names[value]
            return country_name


        # If country is used as site, the country name is selected from countries dictionary by country code. :
        if len(samples) > 0:
            self.nationally_countries = []
            for i in range(0, len(samples)):
                country = select_country(samples[i])
                self.nationally_countries.append(country)
        
        self._generate_graphs()

    def _generate_graphs(self):
        """ The function generating graphs into the presentation! """

        prs = Presentation(self.master) # Read the template presentation

        first_slide = prs.slides[0]
        shape = first_slide.shapes[5]
        text_frame = shape.text_frame

        first_slide_text = "Data Comparison"
    
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = first_slide_text

        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(20)
        font.color.rgb = RGBColor(250,250,250)

        # if (self.country_name in ['Ukraine', 'Poland'] and len(df) > 2):
        #    main_col = 'Site ID'
        # else:
        main_col = 'Site Name'

        ########################
        #### TOTAL PATIENTS ####
        ########################
        column_name = 'Total Patients'
        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=True)
        # combine them to one dataframe
        tmp_df = tmp_ndf.append(tmp_sldf)
        tmp_df = tmp_df.sort_values([column_name], ascending=True)

        title = 'Total number of cases - admission date in {}'.format(self.name)
        GenerateCountriesCompGraphs(ndf=tmp_df, sldf=None, presentation=prs, title=title, column_name=column_name)

        ########################
        #### TOTAL PATIENTS ####
        ########################
        column_name = 'Total Patients'
        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=True)
        # combine them to one dataframe
        tmp_df = tmp_ndf.append(tmp_sldf)
        tmp_df = tmp_df.sort_values([column_name], ascending=True)

        title = 'Total number of cases - admission date in {}'.format(self.name)
        GenerateCountriesCompGraphs(ndf=tmp_df, sldf=None, presentation=prs, title=title, column_name=column_name, samples=self.nationally_countries)
        
        ############################
        #### MEDIAN PATIENT AGE ####
        ############################
        column_name = 'Median patient age'
        title = 'Median patient age'
        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=True)
        
        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name)

        ###############
        # STROKE TYPE #
        ###############
        column_name = '% stroke type - ischemic stroke'
        legend = ['ischemic', 'transient ischemic attack', 'intracerebral hemorrhage', 'subarrachnoid hemorrhage', 'cerebral venous thrombosis', 'undetermined']

        tmp_ndf = self.ndf[[main_col, '% stroke type - ischemic stroke', '% stroke type - transient ischemic attack', '% stroke type - intracerebral hemorrhage', '% stroke type - subarrachnoid hemorrhage', '% stroke type - cerebral venous thrombosis', '% stroke type - undetermined stroke']]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)

        tmp_sldf = self.sldf[[main_col, '% stroke type - ischemic stroke', '% stroke type - transient ischemic attack', '% stroke type - intracerebral hemorrhage', '% stroke type - subarrachnoid hemorrhage', '% stroke type - cerebral venous thrombosis', '% stroke type - undetermined stroke']]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending = True)

        title = "Stroke type, ordered by % ischemic stroke"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #######################
        ### DEPARTMENT TYPE ###
        #######################
        column_name = '% department type - neurology'
        legend = ['neurology', 'neurosurgery', 'anesthesiology resuscitation critical care', 'internal medicine', 'geriatrics', 'other']

        tmp_ndf = self.ndf[[main_col, '% department type - neurology', '% department type - neurosurgery', '% department type - anesthesiology/resuscitation/critical care', '% department type - internal medicine', '% department type - geriatrics', '% department type - Other']]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)

        tmp_sldf = self.sldf[[main_col, '% department type - neurology', '% department type - neurosurgery', '% department type - anesthesiology/resuscitation/critical care', '% department type - internal medicine', '% department type - geriatrics', '% department type - Other']]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending = True)

        title = "Department type, ordered by % neurology"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ###################################
        ### HOSPITALIZATION DESTINATION ###
        ###################################
        column_name = '% patients hospitalized in stroke unit / ICU'
        legend = ['% stroke unit / ICU', '% monitored bed with telemetry', '% standard bed']

        tmp_ndf = self.ndf[[main_col, '% patients hospitalized in stroke unit / ICU', '% patients hospitalized in monitored bed with telemetry', '% patients hospitalized in standard bed']]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)

        tmp_sldf = self.sldf[[main_col, '% patients hospitalized in stroke unit / ICU', '% patients hospitalized in monitored bed with telemetry', '% patients hospitalized in standard bed']]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending = True)

        title = "Hospitalization type, ordered by % stroke unit"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################
        ### CT / MRI performed ###
        ##########################
        column_name = '% CT/MRI - performed'
        legend = ['% performed', '% performed within 1 hour after admission']

        tmp_ndf = self.ndf[[main_col, '% CT/MRI - performed', '% CT/MRI - Performed within 1 hour after admission']]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)

        tmp_sldf = self.sldf[[main_col, '% CT/MRI - performed', '% CT/MRI - Performed within 1 hour after admission']]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending = True)

        title = "% patients receiving CT / MRI"
        subtitle = "- Calculated out of number of IS + ICH + TIA + CVT -"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, legend=legend, number_of_series=len(legend), graph_type='grouped', subtitle=subtitle)

        ############################
        # RECANALIZATION TREATMENT #
        ############################
        column_name = '% recanalization procedures - IV tPa'
        legend = ['IV tPa', 'IV tPa + endovascular treatment', 'endovascular treatment', 'IV tPa + another centre for endovascular treatment']

        tmp_ndf = self.ndf[[main_col, '% patients recanalized', '% recanalization procedures - IV tPa', '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone', '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment']]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)

        tmp_sldf = self.sldf[[main_col, '% patients recanalized', '% recanalization procedures - IV tPa', '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone', '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment']]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending = True)

        title = "% patients receiving recanalization procedures"
        subtitle = "- Calculated out of number of IS -"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, legend=legend, number_of_series=len(legend), graph_type='stacked', subtitle=subtitle)

        ################
        # % median DTN #
        ################
        column_name = 'Median DTN (minutes)'
        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = False)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=False)

        title = "Median Door-to-Needle Time (DTN), in minutes"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name)

        #######################
        # dysphagia screening #
        #######################
        column_name = '% dysphagia screening - Guss test'
        column_names = ['% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre', '% dysphagia screening - Unable to test']
        legend = ['% GUSS test', '% Other test', '% Another centre', '% Unable to test']

        tmp_ndf = self.ndf[[main_col, '% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre', '% dysphagia screening - Unable to test']]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)

        tmp_sldf = self.sldf[[main_col, '% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre', '% dysphagia screening - Unable to test']]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending = True)

        title = "% patients screened for dysphagia, ordered by % GUSS test"
        subtitle = "- Calculated out of number of IS + ICH + CVT -"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, legend=legend, number_of_series=len(legend), graph_type='stacked', subtitle=subtitle)

        #################################
        # % ASSESSED FOR REHABILITATION #
        #################################
        column_name = '% patients assessed for rehabilitation - Yes'
        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=True)

        title = "% patients assessed for rehabilitation within 72 hrs after admission"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name)

        ##################################
        # CAROTID ARTERIES IMAGING - YES #
        ##################################
        column_name = '% carotid arteries imaging - Yes'

        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=True)

        title = "% patients receiving carotid arteries imaging"
        subtitle = "- Calculated out of number of IS + TIA -"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, subtitle=subtitle)

        ##################################
        # PRESCRIPTION OF ANTICOAGULANTS #
        ##################################
        # column_name = '% patients prescribed anticoagulants with aFib with CVT'
        column_name = '% patients prescribed anticoagulants with aFib'

        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=True)

        title = "Prescription of anticoagulants for patients with atrial fibrillation"
        subtitle = "- Calculated out of number of patients with atrial fibrillation discharge alive -"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, subtitle=subtitle)

        ###################################
        # PRESCRIPTION OF ANTITHROMBOTICS #
        ###################################
        column_name = '% patients prescribed antithrombotics with aFib'

        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=True)

        title = "Prescription of antithrombotics for patients with atrial fibrillation"
        subtitle = "- Calculated out of number of patients with atrial fibrillation discharge alive -"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, subtitle=subtitle)

        ##########################
        # DISCHARGE WITH STATINS #
        ##########################
        column_name = '% patients prescribed statins - Yes'

        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=True)

        title = "% patients discharged on statins"
        subtitle = "- Calculated out of number of IS + TIA -"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, subtitle=subtitle)
        
        #############################################################
        # % ANTIHYPERTENSIVE MEDICATION PRESCRIBED out of all cases #
        #############################################################
        column_name = '% prescribed antihypertensives - Yes'

        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=True)

        title = "% patients prescribed antihypertensives at discharge"
        subtitle = "- Calculated out of number of patients discharged alive -"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, column_name=column_name, subtitle=subtitle)

        ###########################################################
        # % RECOMMENDED TO A CEREBROVASCULAR EXPERT - RECOMMENDED #
        ###########################################################
        column_name = '% recommended to a cerebrovascular expert - Recommended'

        # Nationally df
        tmp_ndf = self.ndf[[main_col, column_name]]
        tmp_ndf = tmp_ndf.sort_values([column_name], ascending = True)
        # Site-level df
        tmp_sldf = self.sldf[[main_col, column_name]]
        tmp_sldf = tmp_sldf.sort_values([column_name], ascending=True)

        title = "% patients prescribed antihypertensives at discharge"
        subtitle = "- Calculated out of number of patients discharged alive -"

        GenerateCountriesCompGraphs(ndf=tmp_ndf, sldf=tmp_sldf, presentation=prs, title=title, subtitle=subtitle, column_name=column_name)

        # set pptx output name (for cz it'll be presentation_CZ.pptx)
        working_dir = os.getcwd()
        pptx = "comparison.pptx"
        presentation_path = os.path.normpath(os.path.join(working_dir, pptx))

        prs.save(presentation_path)


class GenerateCountriesCompGraphs:
    """ The class generating comparison graphs in presentation for nationally samples vs. site samples. 

    :param ndf: the calculated statistics for the national samples
    :type ndf: pandas dataframe
    :param presentation: the opened presentation document
    :type presentation: Presentation object
    :param column_name: the name of column to be included in the graph
    :type column_name: str
    :param title: the title of the slide
    :type title: str
    :param sldf: the calculated statistics for the site samples, can be `None` if Total Patients graph is generated
    :type sldf: pandas dataframe
    :param subtitle: the subtitle of the slide
    :type subtitle: str
    :param graph_type: the type of graph to be generated (normal barplot or stacked barplot)
    :type graph_type: str
    :param number_of_series: the number of columns included in the stacked barplot
    :type number_of_series: int
    :param legend: the list of values in the legend
    :type legend: list
    :param samples: the list of countries which should be displayed with different color
    :type samples: list of string
   
    """
    def __init__(self, ndf, presentation, column_name, title, sldf=None, subtitle="", graph_type=None, number_of_series=0, legend=None, samples=None):

        self.ndf = ndf
        self.sldf = sldf
        self.presentation = presentation
        self.ndf_title = "Nationally representative sample"
        self.sldf_title = "Site-level representative sample"
        self.title = title
        self.column_name = column_name
        self.number_of_series = number_of_series
        self.legend = legend
        self.subtitle = subtitle
        self.samples = samples

        self.font_name = 'Century Gothic'
        self.category_font_size = Pt(10)
        self.data_label_font_size = Pt(11)
        self.categories_column = 'Site Name'
            
        # Check type of graph
        if (graph_type == 'stacked'):
            self._create_stacked_barplot()
        elif (graph_type == 'grouped'):
            self._create_grouped_barplot()
        else:
            self._create_barplot()

    def _get_length_of_legend(self, legend):
        """ The function adjusting the number of letters in legend to quess the number of columns in the legend! """
        count = 0

        for i in legend:
            count = count + len(i)
        
        return count

    def _create_barplot(self):
        """ The function generating into the presentation the normal barplot. """

        maximum = 0

        # If graph is in %, set maximum valut to 100. 
        if '%' in self.title.lower():
            ndf_maximum = 100
            sldf_maximum = 100
        elif "total number of cases" in self.title.lower():
            ndf_maximum = round((max(self.ndf[self.column_name].tolist())),1)
        else:
            ndf_maximum = round((max(self.ndf[self.column_name].tolist())),1)
            sldf_maximum = round((max(self.sldf[self.column_name].tolist())),1)

        # Add slide to presentation (layout 11 is our custom layout where only title 'Agency FB', color: RGBColor(43, 88, 173)  and size:24 is set)
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        # Get title object
        title_placeholders = slide.shapes.title
        # Set title
        title_placeholders.text = self.title

        # If subtitle is not set, then delete placeholder for subtitle in the slide. 
        if self.subtitle == "":
            subtitle = slide.placeholders[1]
            sp = subtitle.element
            sp.getparent().remove(sp)
        else:
            subtitle = slide.placeholders[1]
            subtitle.text = self.subtitle

        # 1st chart (left side) - nationally sample
        chart_data = ChartData()
        chart_data.categories = self.ndf[self.categories_column].tolist()
        chart_data.add_series(self.column_name, self.ndf[self.column_name].tolist())

        if "total number of cases" in self.title.lower():
             # Add chart on slide
            specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
            }
        else:
            # Add chart on slide
            specs = {
                'height': Cm(16.5),
                'width': Cm(15.26),
                'left': Cm(0.5),
                'top': Cm(2)
                }

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        # Get series of chart
        series = chart.series[0]
        
        # If graphs for whole country are generated, set for bar with country with red color
        # else set to blue color (same color as title uses)
        if "total number of cases" in self.title.lower() and self.samples is not None:
            site_names = self.ndf[self.categories_column].tolist()
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                if (site_names[idx] in self.samples):
                    fill.fore_color.rgb = RGBColor(128,0,0)
                else:
                    fill.fore_color.rgb = RGBColor(43, 88, 173)
        else:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(43, 88, 173) 

        # Get plot 
        plot = chart.plots[0]
        # Set for each bar same color
        plot.vary_by_categories = False
        # Show data labels and set font
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = self.data_label_font_size
        data_labels.font.bold = True
        data_labels.font.name = self.font_name

        # Change color of graph title and set color gray
        if "total number of cases" not in self.title.lower():
            chart_text = chart.chart_title.text_frame
            chart_text.text = self.ndf_title
            chart_text.paragraphs[0].font.size = Pt(18)
            chart_text.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

        # Value for x-axis (change font size, name, and other things)
        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = self.category_font_size
        tick_labels.font.name = self.font_name

        # Don't show major gridlines
        value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
        value_axis.has_major_gridlines = False
        # Set range of axis
        value_axis.maximum_scale = ndf_maximum
        value_axis.minimum_scale = 0

        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        # Delete tick marks
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.major_unit = 1
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name

        if "total number of cases" not in self.title.lower():
            # 2nd graph (right side) - site-level samples
            chart_data = ChartData()
            chart_data.categories = self.sldf[self.categories_column].tolist()
            chart_data.add_series(self.column_name, self.sldf[self.column_name].tolist())

            # Add chart on slide
            specs = {
                'height': Cm(16.5),
                'width': Cm(15.26),
                'left': Cm(17.5),
                'top': Cm(2)
                }

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

            # Get series of chart
            series = chart.series[0]
            
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(43, 88, 173) 

            # Get plot 
            plot = chart.plots[0]
            # Set for each bar same color
            plot.vary_by_categories = False
            # Show data labels and set font
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = self.data_label_font_size
            data_labels.font.bold = True
            data_labels.font.name = self.font_name

            # Change color of graph title and set color gray
            chart_text = chart.chart_title.text_frame
            chart_text.text = self.sldf_title
            chart_text.paragraphs[0].font.size = Pt(18)
            chart_text.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

            # Value for x-axis (change font size, name, and other things)
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = self.category_font_size
            tick_labels.font.name = self.font_name

            # Don't show major gridlines
            value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
            value_axis.has_major_gridlines = False
            # Set range of axis
            value_axis.maximum_scale = sldf_maximum
            value_axis.minimum_scale = 0

            # Value for y-axis (change font size, name, and other things)
            category_axis = chart.category_axis
            # Delete tick marks
            category_axis.major_tick_mark = XL_TICK_MARK.NONE
            category_axis.major_unit = 1
            category_labels = category_axis.tick_labels
            category_labels.font.size = self.category_font_size
            category_labels.font.name = self.font_name
    
    def _create_stacked_barplot(self):
        """ The function generating into the presentation the stacked barplot. """

        # Calculate length of legend (in case that legend is too long, make smaller font size)
        count = self._get_length_of_legend(self.legend)

        # Get column names of dataframe
        column_names = self.ndf.columns.tolist()

        index = column_names.index(self.column_name)
        
        # Add new slide into presentation
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        title_placeholders = slide.shapes.title
        title_placeholders.text = self.title

        if self.subtitle == "":
            subtitle = slide.placeholders[1]
            sp = subtitle.element
            sp.getparent().remove(sp)
        else:
            subtitle = slide.placeholders[1]
            subtitle.text = self.subtitle


        # 1st dataframe (nationally sample)
        chart_data = ChartData()
        chart_data.categories = self.ndf[self.categories_column].tolist()
        # Add data in each category
        chart_data.add_series(self.legend[0], self.ndf[column_names[index]].tolist())
        if (self.number_of_series >= 2):
            chart_data.add_series(self.legend[1], self.ndf[column_names[index+1]].tolist())
        if (self.number_of_series >= 3):
            chart_data.add_series(self.legend[2], self.ndf[column_names[index+2]].tolist())
        if (self.number_of_series >= 4):
            chart_data.add_series(self.legend[3], self.ndf[column_names[index+3]].tolist())
        if (self.number_of_series >= 5):
            chart_data.add_series(self.legend[4], self.ndf[column_names[index+4]].tolist())
        if (self.number_of_series >= 6):
            chart_data.add_series(self.legend[5], self.ndf[column_names[index+5]].tolist())
        if (self.number_of_series >= 7):
            chart_data.add_series(self.legend[6], self.ndf[column_names[index+6]].tolist())
        if (self.number_of_series >= 8):
            chart_data.add_series(self.legend[7], self.ndf[column_names[index+7]].tolist())

        # Add chart on slide
        specs = {
            'height': Cm(16.5),
            'width': Cm(15.26),
            'left': Cm(0.5),
            'top': Cm(2)
            }
    
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        series = chart.series[0]
        # If graphs for whole country are generated, set for bar with country with red color
        # else set to blue color (same color as title uses)
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(43, 88, 173)   

        if (self.number_of_series >= 5):
            series = chart.series[4]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(80, 137, 188)  

        # Value for x-axis (change font size, name, and other things)
        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = Pt(11)
        tick_labels.font.name = self.font_name

        value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE

        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.dash_style = MSO_LINE.DASH
        value_axis.major_gridlines.format.line.width = Pt(0.5)
        value_axis.maximum_scale = 100

        category_axis = chart.category_axis
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name
        category_labels.tickLblSkip = 1

        # 2nd dataframe (nationally sample)
        # Get column names of dataframe
        column_names = self.sldf.columns.tolist()

        index = column_names.index(self.column_name)

        chart_data = ChartData()
        chart_data.categories = self.sldf[self.categories_column].tolist()
        # Add data in each category
        chart_data.add_series(self.legend[0], self.sldf[column_names[index]].tolist())
        if (self.number_of_series >= 2):
            chart_data.add_series(self.legend[1], self.sldf[column_names[index+1]].tolist())
        if (self.number_of_series >= 3):
            chart_data.add_series(self.legend[2], self.sldf[column_names[index+2]].tolist())
        if (self.number_of_series >= 4):
            chart_data.add_series(self.legend[3], self.sldf[column_names[index+3]].tolist())
        if (self.number_of_series >= 5):
            chart_data.add_series(self.legend[4], self.sldf[column_names[index+4]].tolist())
        if (self.number_of_series >= 6):
            chart_data.add_series(self.legend[5], self.sldf[column_names[index+5]].tolist())
        if (self.number_of_series >= 7):
            chart_data.add_series(self.legend[6], self.sldf[column_names[index+6]].tolist())
        if (self.number_of_series >= 8):
            chart_data.add_series(self.legend[7], self.sldf[column_names[index+7]].tolist())

        # Add chart on slide
        specs = {
            'height': Cm(16.5),
            'width': Cm(15.26),
            'left': Cm(17.5),
            'top': Cm(2)
            }
    
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        series = chart.series[0]
        # If graphs for whole country are generated, set for bar with country with red color
        # else set to blue color (same color as title uses)
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(43, 88, 173)   

        if (self.number_of_series >= 5):
            series = chart.series[4]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(80, 137, 188)  

        # Value for x-axis (change font size, name, and other things)
        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = Pt(11)
        tick_labels.font.name = self.font_name

        value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE

        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.dash_style = MSO_LINE.DASH
        value_axis.major_gridlines.format.line.width = Pt(0.5)
        value_axis.maximum_scale = 100

        category_axis = chart.category_axis
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name
        category_labels.tickLblSkip = 1

        # Set legend 
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = self.font_name

        if (count > 180 or 'antithrombotics prescribed' in self.title.lower()):
            chart.legend.font.size = Pt(11)
        else:
            chart.legend.font.size = Pt(12)

    def _create_grouped_barplot(self):
        """ The function generating into the presentation the grouped barplot. """

        # Calculate length of legend (in case that legend is too long, make smaller font size)
        count = self._get_length_of_legend(self.legend)

        # Get column names of dataframe
        column_names = self.ndf.columns.tolist()

        index = column_names.index(self.column_name)
        
        # Add new slide into presentation
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[11])
        title_placeholders = slide.shapes.title
        title_placeholders.text = self.title

        if self.subtitle == "":
            subtitle = slide.placeholders[1]
            sp = subtitle.element
            sp.getparent().remove(sp)
        else:
            subtitle = slide.placeholders[1]
            subtitle.text = self.subtitle

        # 1st dataframe (nationally sample)
        chart_data = ChartData()
        chart_data.categories = self.ndf[self.categories_column].tolist()
        # Add data in each category
        chart_data.add_series(self.legend[0], self.ndf[column_names[index]].tolist())
        chart_data.add_series(self.legend[1], self.ndf[column_names[index+1]].tolist())

        # Add chart on slide
        specs = {
            'height': Cm(16.5),
            'width': Cm(15.26),
            'left': Cm(0.5),
            'top': Cm(2)
            }
    
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        series = chart.series[0]
        # If graphs for whole country are generated, set for bar with country with red color
        # else set to blue color (same color as title uses)
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(43, 88, 173)   

        if (self.number_of_series >= 5):
            series = chart.series[4]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(80, 137, 188)  

        # Value for x-axis (change font size, name, and other things)
        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = Pt(11)
        tick_labels.font.name = self.font_name

        value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE

        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.dash_style = MSO_LINE.DASH
        value_axis.major_gridlines.format.line.width = Pt(0.5)
        value_axis.maximum_scale = 100

        category_axis = chart.category_axis
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name
        category_labels.tickLblSkip = 1

        # 2nd dataframe (nationally sample)
        # Calculate length of legend (in case that legend is too long, make smaller font size)
        count = self._get_length_of_legend(self.legend)

        # Get column names of dataframe
        column_names = self.sldf.columns.tolist()

        index = column_names.index(self.column_name)

        chart_data = ChartData()
        chart_data.categories = self.sldf[self.categories_column].tolist()
        # Add data in each category
        chart_data.add_series(self.legend[0], self.sldf[column_names[index]].tolist())
        chart_data.add_series(self.legend[1], self.sldf[column_names[index+1]].tolist())

        # Add chart on slide
        specs = {
            'height': Cm(16.5),
            'width': Cm(15.26),
            'left': Cm(17.5),
            'top': Cm(2)
            }
    
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

        series = chart.series[0]
        # If graphs for whole country are generated, set for bar with country with red color
        # else set to blue color (same color as title uses)
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(43, 88, 173)   

        if (self.number_of_series >= 5):
            series = chart.series[4]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(80, 137, 188)  

        # Value for x-axis (change font size, name, and other things)
        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = Pt(11)
        tick_labels.font.name = self.font_name

        value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE

        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.dash_style = MSO_LINE.DASH
        value_axis.major_gridlines.format.line.width = Pt(0.5)
        value_axis.maximum_scale = 100

        category_axis = chart.category_axis
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_labels = category_axis.tick_labels
        category_labels.font.size = self.category_font_size
        category_labels.font.name = self.font_name
        category_labels.tickLblSkip = 1

        # Set legend 
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = self.font_name


class GenerateYearsCompPresentation:
    """ The class creating presentation with graphs representing country comparison through all years in the dataset (eg. 2016, 2017, 2018, 2019). 

    :param df: the dataframe with calculated statistics per years or period
    :type df: pandas dataframe
    :param name: the name of the presentation
    :type name: str
    
    """

    def __init__(self, df, name):

        self.df = df
        self.name = name

        # Get absolute path to the database.
        script_dir = os.path.dirname(__file__) 

        master_pptx = "countries_comparison.pptx"
        self.master = os.path.normpath(os.path.join(script_dir, "backgrounds", master_pptx))

        self._generate_graphs()

    def _generate_graphs(self):
        """Generate graphs into presentation (pptx)."""

        prs = Presentation(self.master)

        first_slide = prs.slides[0]
        shape = first_slide.shapes[5]
        text_frame = shape.text_frame

        first_slide_text = "Data Comparison"
    
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = first_slide_text

        font = run.font
        font.name = 'Calibri Light'
        font.size = Pt(26)
        font.color.rgb = RGBColor(250,250,250)

        # if (self.country_name in ['Ukraine', 'Poland'] and len(df) > 2):
        #     main_col = 'Site ID'
        # else:
        main_col = 'Site Name'

        # main_col = "Site Name"
        years = ', '.join(map(str, self.df[main_col].tolist()))
        title = "Temporal trends - {}".format(years)

        titles = []
        graph_types = []
        legend = []
        # Total patients and median age graphst)
        column_name = 'Total Patients'
        df = self.df[[main_col, column_name]]
        titles.append("Total Patients")
        graph_types.append("normal")

        column_name = 'Median patient age'
        df1 = self.df[[main_col, column_name]]
        titles.append("Median patient age")
        graph_types.append("normal")
        
        GenerateYearsCompGraphs(df=df, df1=df1, presentation=prs, title=title, titles=titles, graph_types=graph_types)

        legends = []
        titles = []
        graph_types = []
        ### DEPARTMENT TYPE ###
        legend = ['neurology', 'neurosurgery', 'anesthesiology resuscitation critical care', 'internal medicine', 'geriatrics', 'other']
        legends.append(legend)
        titles.append("Department type (%)")
        graph_types.append("stacked")

        df = self.df[[main_col, '% department type - neurology', '% department type - neurosurgery', '% department type - anesthesiology/resuscitation/critical care', '% department type - internal medicine', '% department type - geriatrics', '% department type - Other']]

        ### HOSPITALIZATION DESTINATION ###
        legend = ['stroke unit', 'monitored bed with telemetry', 'standard bed']
        legends.append(legend)
        titles.append("Hospitalization type (%)")
        graph_types.append("stacked")

        df1 = self.df[[main_col, '% patients hospitalized in stroke unit / ICU', '% patients hospitalized in monitored bed with telemetry', '% patients hospitalized in standard bed']]
   
        ### REHABILIATION ###
        column_name = '% patients assessed for rehabilitation - Yes'
        df2 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("Assessed for rehabilitation (%)")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df, df1=df1, df2=df2, presentation=prs, title=title, titles=titles, graph_types=graph_types, legends=legends)

        legends = []
        titles = []
        graph_types = []

        # STROKE TYPE #
        legend = ['ischemic', 'transient ischemic attack', 'intracerebral hemorrhage', 'subarrachnoid hemorrhage', 'cerebral venous thrombosis', 'undetermined']
        legends.append(legend)
        titles.append("Stroke type (%)")
        graph_types.append("stacked")

        df = self.df[[main_col, '% stroke type - ischemic stroke', '% stroke type - transient ischemic attack', '% stroke type - intracerebral hemorrhage', '% stroke type - subarrachnoid hemorrhage', '% stroke type - cerebral venous thrombosis', '% stroke type - undetermined stroke']]

        # CT/MRI performed #
        column_name = '% CT/MRI - performed'
        df1 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("CT/MRI (%)")
        graph_types.append("normal")

        # patients recanalized
        column_name = '% patients recanalized'
        df2 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("% patients recanalized")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df, df1=df1, df2=df2, presentation=prs, title=title, titles=titles, graph_types=graph_types, legends=legends)

        legends = []
        titles = []
        graph_types = []

        # RECANALIZATION PROCEDURES #
        legend = ['IV tPa', 'IV tPa + endovascular treatment', 'endovascular treatment', 'IV tPa + another centre for endovascular treatment']
        legends.append(legend)
        titles.append("Recanalization procedures (%)")

        df = self.df[[main_col, '% recanalization procedures - IV tPa', '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone', '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment']]
        graph_types.append("stacked")

        # MEDIAN DTN #
        column_name = 'Median DTN (minutes)'
        df1 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("Median DTN (minutes)")
        graph_types.append("normal")

        # MEDIAN DTN #
        column_name = 'Median DTG (minutes)'
        df2 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("Median DTG (minutes)")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df, df1=df1, df2=df2, presentation=prs, title=title, titles=titles, graph_types=graph_types, legends=legends)

        legends = []
        titles = []
        graph_types = []
        # DYSPHAGIA SCREENING #
        column_name = '% dysphagia screening - Guss test'
        legend = ['GUSS test', 'Other test', 'Another centre']
        legends.append(legend)
        titles.append("Dysphagia screening (%)")
        graph_types.append("stacked")

        df = self.df[[main_col, '% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre']]

        # CAROTID ARTERIES IMAGING #
        column_name = '% carotid arteries imaging - Yes'
        df1 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("Carotid arteries imaging (%)")
        graph_types.append("normal")

        # % RECOMMENDED TO A CEREBROVASCULAR EXPERT - RECOMMENDED #
        column_name = '% recommended to a cerebrovascular expert - Recommended'
        df2 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("Recommended to a cerebrovascular \nexpert (%)")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df, df1=df1, df2=df2, presentation=prs, title=title, titles=titles, graph_types=graph_types, legends=legends)

        legends = []
        titles = []
        graph_types = []
        # PRESCRIBED ANTIPLATELETS #
        column_name = '% patients prescribed antiplatelets without aFib'
        df = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("% patients prescribed antiplatelets without aFib")
        graph_types.append("normal")

        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB #
        column_name = '% patients prescribed anticoagulants with aFib'
        df1 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("% patients prescribed anticoagulants with aFib")
        graph_types.append("normal")

        #  % PATIENTS PRESCRIBED ANTITHROMBOTICS WITH AFIB  #
        column_name = '% patients prescribed antithrombotics with aFib'
        df2 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("% patients prescribed antithrombotics with aFib")
        graph_types.append("normal")

        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB (HOME) #
        column_name = '% afib patients discharged home with anticoagulants'
        df3 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("% aFib patients discharged home with anticoagulants")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df, df1=df1, df2=df2, df3=df3, presentation=prs, title=title, titles=titles, graph_types=graph_types, legends=legends)

        legends = []
        titles = []
        graph_types = []
        # DISCHARGE WITH STATINS #
        column_name = '% patients prescribed statins - Yes'
        df = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("% patients prescribed statins")
        graph_types.append("normal")

        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB #
        column_name = '% prescribed antihypertensives - Yes'
        df1 = self.df[[main_col, column_name]]
        legend = []
        legends.append(legend)
        titles.append("% patients prescribed antihypertensives")
        graph_types.append("normal")

        GenerateYearsCompGraphs(df=df, df1=df1, presentation=prs, title=title, titles=titles, graph_types=graph_types, legends=legends)

        # set pptx output name (for cz it'll be presentation_CZ.pptx)
        working_dir = os.getcwd()
        pptx = self.name + ".pptx"
        presentation_path = os.path.normpath(os.path.join(working_dir, pptx))

        prs.save(presentation_path)

class GenerateYearsCompGraphs:
    """ The class generating graphs into presentation for country per years. If only one dataframe is provided than one graph is created on slide, two graphs is two dataframes are provided etc. 

    :param presentation: the opened presentation document
    :type presentation: Presentation object
    :param df: dataframe containing calculated statistics
    :type df: pandas dataframe
    :param df1: 2nd dataframe containing calculated statistics
    :type df1: pandas dataframe
    :param df2: 3rd dataframe containing calculated statistics
    :type df2: pandas dataframe
    :param df3: 4th dataframe containing calculated statistics
    :type df3: pandas dataframe
    :param title: the title of the slide
    :type title: str
    :param titles: the titles of each graph
    :type titles: list
    :param graph_types: the list of types of graphs
    :type graph_types: list
    :param legends: the nested list containing lists of legends
    :type legends: nested list
    :param outcome: `True` if outcome calculation should be included in the presentation
    :type outcome: bool
    """
    def __init__(self, presentation, df, df1=None, df2=None, df3=None, title="", titles=None, graph_types=None, legends=[], outcome=False):

        self.df = df
        self.df1 = df1
        self.df2 = df2
        self.df3 = df3
        self.presentation = presentation
        self.title = title
        self.titles = titles
        self.legends = legends
        self.num_graphs = 0
        self.graph_types = graph_types

        self.font_name = 'Century Gothic'
        self.category_font_size = Pt(8)
        self.data_label_font_size = Pt(9)

        if outcome:
            self.categories_column = 'Patient Group'
        else:
            self.categories_column = 'Site Name'

        if df1 is not None:
            self.num_graphs += 1
        if df2 is not None:
            self.num_graphs += 1
        if df3 is not None:
            self.num_graphs += 1

        # Add slide to presentation (layout 11 is our custom layout where only title 'Agency FB', color: RGBColor(43, 88, 173)  and size:24 is set)
        self.slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[12])
        # Get title object
        title_placeholders = self.slide.shapes.title
        # Set title
        title_placeholders.text = self.title.upper()

        subtitle = self.slide.placeholders[1]
        sp = subtitle.element
        sp.getparent().remove(sp)

        self.colors = {
                'blue': RGBColor(43, 88, 173),
                'violet': RGBColor(76, 70, 127),
                'orange': RGBColor(237, 145, 49),
                'green': RGBColor(146, 208, 80),
               # 'dark_blue': RGBColor(37, 94, 145),
                'yellow': RGBColor(255, 192, 0),
                'grey': RGBColor(165, 165, 165)
            }

        if self.num_graphs == 0:

            if outcome:
                # Add chart on slide
                specs = {
                    'height': Cm(16.5),
                    'width': Cm(32),
                    'left': Cm(0.5),
                    'top': Cm(2)
                } 
            else:
                specs = {
                    'height': Cm(16.5),
                    'width': Cm(32),
                    'left': Cm(0.7),
                    'top': Cm(2)
                }
            self._create_plot(df=df, title=titles[0], specs=specs, graph_type=graph_types[0], ix=0)

        # If number of graph is equal to 1 (2 graphs on one slide) set specs that each graph will be on one half of slide page
        if self.num_graphs == 1:
            # Add chart on slide
            specs = {
                'height': Cm(16.5),
                'width': Cm(15.26),
                'left': Cm(0.5),
                'top': Cm(2)
            } 

            # Add chart on slide
            specs1 = {
                'height': Cm(16.5),
                'width': Cm(15.26),
                'left': Cm(17.5),
                'top': Cm(2)
            }

            self._create_plot(df=df, title=titles[0], specs=specs, graph_type=graph_types[0], ix=0)
            self._create_plot(df=df1, title=titles[1], specs=specs1, graph_type=graph_types[1], ix=1)

        # If number of graph is equal to 2 (3 graphs on one slide) set specs that one graph will be on half of page, and two graphs on second half
        if self.num_graphs == 2:
            # Add chart on slide
            specs = {
                'height': Cm(16.5),
                'width': Cm(15.26),
                'left': Cm(0.5),
                'top': Cm(2)
            } 
            # Add chart on slide
            specs1 = {
                'height': Cm(8.25),
                'width': Cm(15.26),
                'left': Cm(17.5),
                'top': Cm(2)
            }
            # Add chart on slide
            specs2 = {
                'height': Cm(8.25),
                'width': Cm(15.26),
                'left': Cm(17.5),
                'top': Cm(10.25)
            }

            self._create_plot(df=df, title=titles[0], specs=specs, graph_type=graph_types[0], ix=0)
            self._create_plot(df=df1, title=titles[1], specs=specs1, graph_type=graph_types[1], ix=1)
            self._create_plot(df=df2, title=titles[2], specs=specs2, graph_type=graph_types[2], ix=2)

        # If number of graph is equal to 3 (4 graphs on one slide) set specs that each graph will be put on quarter of page
        if self.num_graphs == 3:
            # Add chart on slide
            specs = {
                'height': Cm(8.25),
                'width': Cm(15.26),
                'left': Cm(0.5),
                'top': Cm(2)
            } 
            specs1 = {
                'height': Cm(8.25),
                'width': Cm(15.26),
                'left': Cm(0.5),
                'top': Cm(10.25)
            } 
            # Add chart on slide
            specs2 = {
                'height': Cm(8.25),
                'width': Cm(15.26),
                'left': Cm(17.5),
                'top': Cm(2)
            }
            # Add chart on slide
            specs3 = {
                'height': Cm(8.25),
                'width': Cm(15.26),
                'left': Cm(17.5),
                'top': Cm(10.25)
            }
            
            self._create_plot(df=df, title=titles[0], specs=specs, graph_type=graph_types[0], ix=0)
            self._create_plot(df=df1, title=titles[1], specs=specs1, graph_type=graph_types[1], ix=1)
            self._create_plot(df=df2, title=titles[2], specs=specs2, graph_type=graph_types[2], ix=2)
            self._create_plot(df=df3, title=titles[3], specs=specs3, graph_type=graph_types[3], ix=3)
        
    def _set_transparency(self, transparency, elm):
        """ The function set the transparency of the row. 

        :param transparency: the transparency in %
        :type transparency: int
        :param elm: the element which transparency should be changed
        :type elm: format.line.color._xFill
        """
        a = str(100 - transparency) + '196'
        
        alpha = OxmlElement('a:alpha')
        alpha.set('val', a)
        elm.srgbClr.append(alpha)

    def _create_plot(self, df, title, specs, graph_type, legend=None, ix=0):   
        """ The function creating the new graph into the presentation based on the graph type. 
        
        :param df: the dataframe with data to be shown
        :type df: pandas dataframe
        :param title: the title of the graph
        :type title: str
        :param specs: the position settings
        :type specs: dictionary
        :param graph_type: the type of graph (normal or stacked)
        :type graph_type: str
        :param legend: the list of values in legend based on columns (only for stacked barplot)
        :type legend: list
        :param ix: the index which legend should be used 
        :type ix: int
        """

        if graph_type == "normal":
            # Get column names of dataframe
            column_names = df.columns.tolist()
            index = column_names.index(self.categories_column)   

            # 1st chart (left side) - nationally sample
            chart_data = ChartData()
            chart_data.categories = df[self.categories_column].tolist()
            chart_data.add_series(column_names[index+1], df[column_names[index+1]].tolist())

            chart = self.slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, specs['left'], specs['top'], specs['width'], specs['height'], chart_data).chart

            # Get series of chart
            series = chart.series[0]
            series.points[0].format.fill.solid()
            series.points[0].format.fill.fore_color.rgb = self.colors['blue']
            series.points[1].format.fill.solid()
            series.points[1].format.fill.fore_color.rgb = self.colors['orange']
            series.points[2].format.fill.solid()
            series.points[2].format.fill.fore_color.rgb = self.colors['green']
            
            # Get plot 
            plot = chart.plots[0]
            # Set for each bar same color
            plot.vary_by_categories = True
            # Show data labels and set font
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = self.data_label_font_size
            data_labels.font.bold = True
            data_labels.font.name = self.font_name

            chart_text = chart.chart_title.text_frame
            chart_text.text = title
            chart_text.paragraphs[0].font.name = self.font_name
            chart_text.paragraphs[0].font.size = Pt(14)
            chart_text.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

            # Value for x-axis (change font size, name, and other things)
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = self.category_font_size
            tick_labels.font.name = self.font_name

            # Don't show major gridlines
            #value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
            value_axis.has_major_gridlines = True
            value_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
            value_axis.major_gridlines.format.line.width = Pt(0.5)
            # Set range of axis
            value_axis.minimum_scale = 0
            values = df[column_names[index+1]].tolist()
            max_value = max(values)
            if '%' in title and max_value >= 90:
                value_axis.maximum_scale = 100
            else:
                value_axis.maximum_scale = math.ceil(max_value / 10.0) * 10
            value_axis.visible = False

            # Value for y-axis (change font size, name, and other things)
            category_axis = chart.category_axis
            # Delete tick marks
            category_axis.major_tick_mark = XL_TICK_MARK.NONE
            category_axis.major_unit = 1
            category_labels = category_axis.tick_labels
            category_labels.font.size = self.category_font_size
            category_labels.font.name = self.font_name

        # Create stacked barplot
        else: 
            # Get column names of dataframe
            column_names = df.columns.tolist()
            index = column_names.index(self.categories_column) + 1

            legend = self.legends[ix]
            number_of_series = len(legend)

            # 1st dataframe (nationally sample)
            chart_data = ChartData()
            chart_data.categories = df[self.categories_column].tolist()
            # Add data in each category
            chart_data.add_series(legend[0], df[column_names[index]].tolist())
            if (number_of_series >= 2):
                chart_data.add_series(legend[1], df[column_names[index+1]].tolist())
            if (number_of_series >= 3):
                chart_data.add_series(legend[2], df[column_names[index+2]].tolist())
            if (number_of_series >= 4):
                chart_data.add_series(legend[3], df[column_names[index+3]].tolist())
            if (number_of_series >= 5):
                chart_data.add_series(legend[4], df[column_names[index+4]].tolist())
            if (number_of_series >= 6):
                chart_data.add_series(legend[5], df[column_names[index+5]].tolist())
            if (number_of_series >= 7):
                chart_data.add_series(legend[6], df[column_names[index+6]].tolist())
            if (number_of_series >= 8):
                chart_data.add_series(legend[7], df[column_names[index+7]].tolist())

            chart = self.slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_STACKED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart

            series = chart.series[0]
            # Get series of chart
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = self.colors['blue']
            if (number_of_series >= 2):
                chart.series[1].format.fill.solid()
                chart.series[1].format.fill.fore_color.rgb = self.colors['orange']
            if (number_of_series >= 3):
                chart.series[2].format.fill.solid()
                chart.series[2].format.fill.fore_color.rgb = self.colors['green']
            if (number_of_series >= 4):
                chart.series[3].format.fill.solid()
                chart.series[3].format.fill.fore_color.rgb = self.colors['grey']
            if (number_of_series >= 5):
                chart.series[4].format.fill.solid()
                chart.series[4].format.fill.fore_color.rgb = self.colors['violet']
            if (number_of_series >= 6):
                chart.series[5].format.fill.solid()
                chart.series[5].format.fill.fore_color.rgb = self.colors['yellow']

            # Value for x-axis (change font size, name, and other things)
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = Pt(11)
            tick_labels.font.name = self.font_name

            value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE

            # Don't show major gridlines
            #value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
            value_axis.has_major_gridlines = True
            value_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
            value_axis.major_gridlines.format.line.width = Pt(0.5)
            # Set range of axis
            value_axis.minimum_scale = 0
            value_axis.maximum_scale = 100
            value_axis.visible = True

             # Set 100% transparency to value axis
            value_axis.format.line.color.rgb = RGBColor(0, 0, 0)
            solidFill = value_axis.format.line.color._xFill
            self._set_transparency(100, solidFill)


            category_axis = chart.category_axis
            category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
            solidFill = category_axis.format.line.color._xFill
            self._set_transparency(100, solidFill)
            
            category_axis.major_tick_mark = XL_TICK_MARK.NONE
            category_labels = category_axis.tick_labels
            category_labels.font.size = self.category_font_size
            category_labels.font.name = self.font_name
            category_labels.tickLblSkip = 1
        
            # Set legend 
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.name = self.font_name
            chart.legend.font.size = Pt(14)

            chart_text = chart.chart_title.text_frame
            chart_text.text = title
            chart_text.paragraphs[0].font.size = Pt(18)
            chart_text.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

        


