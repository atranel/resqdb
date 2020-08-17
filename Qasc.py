from resqdb.Connection import Connection
from resqdb.functions import save_file
from datetime import datetime
import pandas as pd
import logging
import os

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

class Qasc():

    def __init__(self):

        # Set logging
        debug =  f'debug_{datetime.now().strftime("%d-%m-%Y")}.log' 
        log_file = os.path.join(os.getcwd(), debug)
        logging.basicConfig(
            filename=log_file, 
            filemode='a', 
            format='%(asctime)s,%(msecs)d %(name)s %(levelname)s %(message)s', 
            datefmt='%H:%M:%S', 
            level=logging.DEBUG
            )
        logging.info('Start to generate QASC reports.')

        # Connect to database and get QASC data
        con = Connection(data='qasc')
        self.preprocessed_data = con.qasc_preprocessed_data
        # Get also data with studies 
        self.study_df = con.study_df

        # Get list of Site IDs from the preprocessed data
        site_ids = set(self.preprocessed_data['SITE_ID'].tolist())
        self.__site_ids = [x for x in site_ids if not x.endswith('_AUD')]
        del site_ids 

        self.table_font_size = Pt(11) # Set default font size for tables 

    @property
    def site_ids(self):
        return self.__site_ids

    @property
    def stats(self):
        return self.__stats

    @stats.setter
    def stats(self, df):
        self.__stats = df

    @property
    def site_id(self):
        return self.__site_id

    @site_id.setter
    def site_id(self, site_id):
        self.__site_id = site_id

    @property
    def pre_df(self):
        return self.__pre_df

    @pre_df.setter
    def pre_df(self, df):
        self.__pre_df = df

    @property
    def post_df(self):
        return self.__post_df

    @post_df.setter
    def post_df(self, df):
        self.__post_df = df

    @property
    def site_df(self):
        return self.__site_df

    @site_df.setter
    def site_df(self, df):
        self.__site_df = df

    @property
    def table_font_size(self):
        return self.__table_font_size

    @table_font_size.setter
    def table_font_size(self, value):
        self.__table_font_size = value

    def generate_reports(self, site_id):
        ''' Generate reports for the site ID. 
        
        :param site_id: the site ID of the hospital for which the report should be generated
        :type site_id: str
        '''
        self.site_id = site_id.upper() # Make provided site ID uppercased
        # check if provided site ID is in preprocessed data
        if self.site_id is not None and self.site_id in self.site_ids:
            # Filter dataframe for provided site ID
            self.site_df = self._filter_by_site(self.site_id)
            if self.site_df.empty:
                print(f"There are no data for this {site_id} hospital. The available sites are: {self.site_ids}.")
            else:
                # Obtain data for pre/post period
                self._pre_post_data()
                # Calculate pre stats and generate baselina reports that are made from pre dataset
                self.pre_stats = self.calculate_statistics(df=self.pre_df)
                self.generate_baseline_report(df=self.pre_stats)
                # If post dataset is not empty, generate also comparison reports otherwise stop the program
                if not self.post_df.empty:
                    self.post_stats = self.calculate_statistics(df=self.post_df)
                    self.generate_pre_post_report()
        else:
            print(f"There are no data for this site.")

    def _filter_by_site(self, site_id):
        ''' Return the filtered dataframe for site id filter on column SITE_ID. 
        
        :param site_id: site ID for which the report should be generated
        :type site_id: str
        :returns: filtered data
        '''
        return self.preprocessed_data.loc[self.preprocessed_data['SITE_ID'] == site_id].copy()

    def _filter_data(self, first_date, second_date):
        ''' Return the filtered dataframe where DATE_CREATED is between first and second date. 

        :param first_date: first date included in the filtration
        :type first_date: datetime
        :param second_date: second date included in the filtration
        :type second_date: datetime
        :returns: the filtered dataframe
        '''
        if self.site_df.empty or self.site_df is None:
            return self.preprocessed_data.loc[
                self.preprocessed_data['DATE_CREATED'].between(first_date, second_date)
            ].copy()
        else:
            return self.site_df.loc[
                self.preprocessed_data['DATE_CREATED'].between(first_date, second_date)
            ].copy()

    def _pre_post_data(self):
        ''' Set pre/post dataframe data. '''        
        # NOTE: for now the pre dataset is set to year 2019 and post dataset to 2020
        pre_date1 = datetime(2019, 1, 1)
        pre_date2 = datetime(2019, 12, 31)
        self.pre_df = self._filter_data(pre_date1, pre_date2)
        if (self.pre_df.empty):
            print(f"There are no data for pre phase for this {self.site_id} hospital.")
            exit()
        else:
            post_date1 = datetime(2020, 1, 1)
            post_date2 = datetime(2020, 12, 31)
            self.post_df = self._filter_data(post_date1, post_date2)
            if (self.post_df.empty):
                if self.pre_df.emtpy:
                    print(f"There are no data for post phase for this {self.site_id} hospital.")
                    exit()

    def _get_percentage_column_name(self, column_name):
        ''' Return value where # is replaced by %. 
        
        :params column_name: column name to be replaced
        :type column_name: string
        '''
        return column_name.replace('#', '%', 1)

    def _get_patients(self, stats, column_name, groups, out_of):
        ''' Return # and % of the group. 
        
        :param column_name: name of column to be creatd
        :type column_name: string
        :param groups: grouped dataframe
        :type groups: dataframe
        :param out_of: the name of column to be used as denominator
        :type out_of: string
        ''' 
        stats = stats.merge(groups[[self.main_col, column_name]], how='outer')
        # get percentages out of # n
        stats[self._get_percentage_column_name(column_name)] = stats.apply(
            lambda x: round(((x[column_name]/x[out_of]) * 100), 2) if x[out_of] > 0 else 0, axis=1)
        stats.fillna(0, inplace=True)
        return stats

    def _rename_column(self, df, prev_name, curr_name):
        ''' Rename the column in dataframe. If column not exists, create column. 
        
        :param df: the grouped dataframe
        :type df: dataframe
        :param prev_name: the name of column to be changed
        :type prev_name: str/int
        :param curr_name: new name of the column
        :type curr_name: str
        :returns: modified dataframe
        '''
        if not prev_name in df.columns:
            df[prev_name] = 0
        return df.rename(columns={prev_name: curr_name})    

    def calculate_statistics(self, df=None):
        ''' Calculate the statistics for the temperature, blood glucose and swallow screening. '''
        # Defina main column to be grouped by
        self.main_col = 'SITE_ID'

        if df is None:
            df = self.preprocessed_data.copy()

        # 1. Patients records entered
        stats = df.groupby([self.main_col]).size().to_frame('n').reset_index()

        # Calculate # of patients per stroke type
        groups = df.groupby([self.main_col, 'STROKE_TYPE']).size().unstack().reset_index().fillna(0)
        column_name = '# acute stroke'
        column_names = {
            1.0: '# subarrachnoid hemorrhage', 
            2.0: '# ischemic stroke', 
            3.0: '# intracerrebral hemorrhage'
        }
        # Check if all columns are in the grouped dataframe
        for key in column_names.keys():
            if key not in groups.columns:
                groups[key] = 0
        groups.rename(columns=column_names, inplace=True)
        groups[column_name] = groups[column_names.values()].sum(axis=1)
        stats = self._get_patients(stats=stats, column_name=column_name, groups=groups, out_of='n')

        ''' 
        2. Temperature monitring and treatment
        column_name = TEMP_MEASUREMENT
        question = Was temperature measured?
        type = checkbox
        answers:    1. at least four times on day one of admission, 
                    2. at least four times on day two of admission, 
                    3. at least four times on day three of admission, 
                    4. none of the above
        calculation: 1-3 are selected
        '''
        groups = df.groupby([self.main_col, 'TEMP_MEASUREMENT']).size().unstack().reset_index().fillna(0)
        # remove column with default values from the groups
        default = '1,2,3,4'
        if default in groups.columns:
            groups.drop([default], axis=1, inplace=True)  
        column_names = groups.columns

        include_columns = [x for x in column_names if '1' in x or '2' in x or '3' in x and '4' not in x]
        column_name = '# Temperature monitoring and treatment'
        groups[column_name] = groups[include_columns].sum(axis=1)
        stats = self._get_patients(stats=stats, column_name=column_name, groups=groups, out_of='n')

        include_columns = [x for x in column_names if '1' in x and '4' not in x]
        column_name = '# Temperature monitored at least four times per day - Day of admission'
        groups[column_name] = groups[include_columns].sum(axis=1)
        stats = self._get_patients(stats=stats, column_name=column_name, groups=groups, out_of='n')
       
        include_columns = [x for x in column_names if '2' in x and '4' not in x]
        column_name = '# Temperature monitored at least four times per day - Day two of admission'
        groups[column_name] = groups[include_columns].sum(axis=1)
        stats = self._get_patients(stats=stats, column_name=column_name, groups=groups, out_of='n')

        include_columns = [x for x in column_names if '3' in x and '4' not in x]
        column_name = '# Temperature monitored at least four times per day - Day three of admission'
        groups[column_name] = groups[include_columns].sum(axis=1)
        stats = self._get_patients(stats=stats, column_name=column_name, groups=groups, out_of='n')

        ''' 
        3. Temperature > 37.5°c recorded within 72 hours of admission
        column_name = FEVER
        question = In the first 72hrs following admission did the patient develop a fever ≥ 37.5 °C?
        type = select
        answers:    1. Yes,
                    2. No,
                    3. Unknown
        calculation: # of 1 selected
        '''
        groups = df.groupby([self.main_col, 'FEVER']).size().unstack().reset_index().fillna(0)
        column_name = '# Temperature > 37.5°c recorded within 72 hours of admission'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(stats=stats, column_name=column_name, groups=groups, out_of='n')

        ''' 
        3a. Paracetamol (or other anti-pyretic) given for first temperature > 37.5°c
        column_name = PARACETAMOL
        question = Was paracetamol (or other antipyretic) for the first elevated temperature (>37.5 °C) administered? 
        type = select
        answers:    1. Yes,
                    2. No,
                    3. Unknown
        calculation: # of 1 selected
        condition: this question is show if question 3 is answered "Yes"
        '''
        fever_df = df.loc[df['FEVER'] == 1].copy()
        groups = fever_df.groupby([self.main_col, 'PARACETAMOL']).size().unstack().reset_index().fillna(0)
        column_name = '# Paracetamol (or other anti-pyretic) given for first temperature > 37.5°C'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(
            stats=stats, 
            column_name=column_name, 
            groups=groups, 
            out_of='# Temperature > 37.5°c recorded within 72 hours of admission'
            )

        ''' 
        3b. Paracetamol (or other anti-pyretic) given with one hour from first temperature > 37.5°c #
        column_name = PARACETAMOL_1H
        question = Was paracetamol (or other antipyretic) for the first elevated temperature (>37.5 °C) administered within 1 hour? 
        type = select
        answers:    1. Yes,
                    2. No,
                    3. Unknown
        calculation: # of 1 selected
        condition: this question is show if question 3a is answered "Yes"
        '''
        first_temperature_df = df.loc[df['PARACETAMOL'] == 1].copy()
        groups = first_temperature_df.groupby([self.main_col, 'PARACETAMOL_1H']).size().unstack().reset_index().fillna(0)
        column_name = '# Paracetamol (or other anti-pyretic) given with one hour from first temperature > 37.5°C'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(
            stats=stats, 
            column_name=column_name, 
            groups=groups, 
            out_of='# Paracetamol (or other anti-pyretic) given for first temperature > 37.5°C'
            )

        ''' 
        4. Blood glucose monitoring and treatment
        column_name = GLUCOSE_LAB
        question = Was a venous blood glucose level sample collected and sent to laboratory?
        type = select
        answers:    1. Yes,
                    2. No,
                    3. Unknown
        calculation: # of 1 selected
        '''
        groups = df.groupby([self.main_col, 'GLUCOSE_LAB']).size().unstack().reset_index().fillna(0)
        column_name = '# Blood glucose monitoring and treatment'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(
            stats=stats, 
            column_name=column_name, 
            groups=groups, 
            out_of='n'
            )

        ''' 
        5. Blood Glucose Level (BGL) monitored > four times per day
        column_name = GLUCOSE_MONITOR
        question = Was a finger-prick blood glucose level recorded
        type = checkbox
        answers:    1. at least four times on day one of admission, 
                    2. at least four times on day two of admission, 
                    3. at least four times on day three of admission, 
                    4. none of the above
        calculation: 1-3 is selected
        '''
        monitoring_fever_hyperglycemia_df = df.loc[(df['TEMP_MEASUREMENT'].str.contains('1|2|3') & ~df['TEMP_MEASUREMENT'].str.contains('4')) | (df['GLUCOSE_MONITOR'].str.contains('1|2|3') & ~df['GLUCOSE_MONITOR'].str.contains('4'))]
        monitoring_fever_hyperglycemia_stats = monitoring_fever_hyperglycemia_df.groupby([self.main_col]).size().to_frame('# monitored for fever and/or hyperglycaemia four times a day').reset_index()
        stats = stats.merge(monitoring_fever_hyperglycemia_stats, how='outer')
        del monitoring_fever_hyperglycemia_df, monitoring_fever_hyperglycemia_stats

        groups = df.groupby([self.main_col, 'GLUCOSE_MONITOR']).size().unstack().reset_index().fillna(0)
        # remove column with default values from the groups
        default = '1,2,3,4'
        if default in groups.columns:
            groups.drop([default], axis=1, inplace=True)  
        column_names = groups.columns

        include_columns = [x for x in column_names if '1' in x and '4' not in x]
        column_name = '# Blood Glucose Level (BGL) monitored > four times per day - Day of admission'
        groups[column_name] = groups[include_columns].sum(axis=1)
        stats = self._get_patients(stats=stats, column_name=column_name, groups=groups, out_of='n')

        include_columns = [x for x in column_names if '2' in x and '4' not in x]
        column_name = '# Blood Glucose Level (BGL) monitored > four times per day - Day two of admission'
        groups[column_name] = groups[include_columns].sum(axis=1)
        stats = self._get_patients(stats=stats, column_name=column_name, groups=groups, out_of='n')

        include_columns = [x for x in column_names if '3' in x and '4' not in x]
        column_name = '# Blood Glucose Level (BGL) monitored > four times per day - Day three of admission'
        groups[column_name] = groups[include_columns].sum(axis=1)
        stats = self._get_patients(stats=stats, column_name=column_name, groups=groups, out_of='n')

        ''' 
        6. BGL ≥ 10mmol/L within 48 hours of admission
        column_name = GLUCOSE_LEVEL
        question = In the first 48 hours following ward admission did the patient develop a finger-prick glucose level of greater or equal to 10mmols/L?
        type = select
        answers:    1. Yes,
                    2. No,
                    3. Unknown
        calculation: # of 1 selected
        '''
        groups = df.groupby([self.main_col, 'GLUCOSE_LEVEL']).size().unstack().reset_index().fillna(0)
        column_name = '# BGL ≥ 10mmol/L within 48 hours of admission'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(
            stats=stats, 
            column_name=column_name, 
            groups=groups, 
            out_of='n'
            )

        ''' 
        6a. Insulin given for first BGL ≥ 10mmol/L
        column_name = INSULIN_ADMINISTRATION
        question = Was insulin for first elevated finger prick glucose (>=10mmol/L) administered?
        type = select
        answers:    1. Yes,
                    2. No,
                    3. Unknown
        calculation: # of 1 selected
        '''
        bgl_followed_df = df.loc[df['GLUCOSE_LEVEL'] == 1].copy()
        groups = bgl_followed_df.groupby([self.main_col, 'INSULIN_ADMINISTRATION']).size().unstack().reset_index().fillna(0)
        column_name = '# Insulin given for first BGL ≥ 10mmol/L'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(
            stats=stats, 
            column_name=column_name, 
            groups=groups, 
            out_of='# BGL ≥ 10mmol/L within 48 hours of admission'
            )


        ''' 
        6b. Insulin given within one hour from first BGL ≥ 10mmol/L #
        column_name = INSULIN_ADMINISTRATION_1H
        question = Was insulin for first elevated finger prick glucose (>=10mmol/L) administered within 1 hour?
        type = select
        answers:    1. Yes,
                    2. No,
                    3. Unknown
        calculation: # of 1 selected
        '''
        insulin_administration_df = df.loc[df['INSULIN_ADMINISTRATION'] == 1].copy()
        groups = insulin_administration_df.groupby([self.main_col, 'INSULIN_ADMINISTRATION_1H']).size().unstack().reset_index().fillna(0)
        column_name = '# Insulin given within one hour from first BGL ≥ 10mmol/L'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(
            stats=stats, 
            column_name=column_name, 
            groups=groups, 
            out_of='# Insulin given for first BGL ≥ 10mmol/L'
            )

        ''' 
        7. Swallow screening
        column_name = DYSPHAGIA
        question = Was a formal swallowing screen performed (i.e. not a test of gag reflex)?
        type = select
        answers:    1. Performed,
                    2. Not performed,
                    3. Not applicable
        calculation: # of 1 selected
        '''
        
        ''' 
        8. Formal swallow screen performed
        column_name = DYSPHAGIA
        question = Was a formal swallowing screen performed (i.e. not a test of gag reflex)?
        type = select
        answers:    1. Performed,
                    2. Not performed,
                    3. Not applicable
        calculation: # of 1 selected
        '''
        groups = df.groupby([self.main_col, 'DYSPHAGIA']).size().unstack().reset_index().fillna(0)
        column_name = '# Formal swallow screen performed'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(
            stats=stats, 
            column_name=column_name, 
            groups=groups, 
            out_of='n'
            )

        ''' 
        8a. Swallow screen performed within 24 hours #
        column_name = DYSPHAGIA_24H
        question = Did the patient receive a swallowing screen within 24 hours of admission to hospital?
        type = select
        answers:    1. Yes,
                    2. No,
                    3. Unknown
        calculation: # of 1 selected
        '''
        dysphagia_performed_df = df.loc[df['DYSPHAGIA'] == 1].copy()
        groups = insulin_administration_df.groupby(
            [self.main_col, 'DYSPHAGIA_24H']).size().unstack().reset_index().fillna(0)
        column_name = '# Swallow screen performed within 24 hours'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(
            stats=stats, 
            column_name=column_name, 
            groups=groups, 
            out_of='# Formal swallow screen performed'
            )

        ''' 
        9. Swallow screen or swallow assessment performed before being given oral medications #
        column_name = DYSPH_BEFORE_MED
        question = Was the swallow screen or swallow assessment performed before the patient was given oral medications?
        type = select
        answers:    1. Yes,
                    2. No,
                    3. Unknown
        calculation: # of 1 selected
        '''
        groups = df.groupby([self.main_col, 'DYSPH_BEFORE_MED']).size().unstack().reset_index().fillna(0)
        column_name = '# Swallow screen or swallow assessment performed before being given oral medications'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(
            stats=stats, 
            column_name=column_name, 
            groups=groups, 
            out_of='n'
            )

        ''' 
        10. Swallow screen or swallow assessment performed before being given oral food or fluids #
        column_name = DYSPH_BEFORE_FOOD
        question = Was the swallow screen or swallow assessment performed before the patient was given oral food or fluids?
        type = select
        answers:    1. Yes,
                    2. No,
                    3. Unknown
        calculation: # of 1 selected
        '''
        groups = df.groupby([self.main_col, 'DYSPH_BEFORE_FOOD']).size().unstack().reset_index().fillna(0)
        column_name = '# Swallow screen or swallow assessment performed before being given oral food or fluids'
        groups = self._rename_column(df=groups, prev_name=1.0, curr_name=column_name)
        stats = self._get_patients(
            stats=stats, 
            column_name=column_name, 
            groups=groups, 
            out_of='n'
            )

        # save calculated stats into csv file
        save_file(
            name=f'qasc_stats_{datetime.now().strftime("%Y-%m-%d")}.csv',
            data=stats,
            )

        return stats

    def _merge_cells(self, table, first_x, first_y, second_x, second_y):
        ''' Merge cells based on index. 
        
        :param table: the table to be modified
        :type table: table shape
        :param first_x: the 1st coordinate
        :type first_x: int
        :param first_y: the 2nd coordinate
        :type first_y: int
        :param second_x: 3rd coordinate
        :type second_x: int
        :param second_y: 4th coordinate
        :type second_y: int
        '''
        cell = table.cell(first_x, first_y)
        other_cell = table.cell(second_x, second_y)
        cell.merge(other_cell)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        return cell

    def _insert_values(self, table, row, values, bold=False):
        ''' Insert values into 'n' and '%' columns. 
        
        :param table: the table 
        :type table: table shape
        :param row: the number of row to add values
        :type row: int
        :param values: the list of values
        :type values: list
        :param bold: text should be bold
        :type bold: boolean
        '''
        for i in range(0, len(values)):
            cell = table.cell(row, i + 2) # Get cell in the 2nd row
            cell.text = values[i]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraphs in cell.text_frame.paragraphs:
                paragraphs.alignment = PP_ALIGN.CENTER
                for run in paragraphs.runs:
                    run.font.size = self.table_font_size
                    run.font.bold = bold

    def _insert_subrows(self, table, row, col, values):
        ''' Insert values into 'n' and '%' columns. 
        
        :param table: the table 
        :type table: table shape
        :param row: the number of row to add values
        :type row: int
        :param col: the number of column to add values
        :type col: int
        :param values: the list of values
        :type values: list
        :param bold: text should be bold
        :type bold: boolean
        '''
        for i in range(0, len(values)):
            cell = table.cell(row, col) # Get cell in the 2nd row
            cell.text = values[i]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraphs in cell.text_frame.paragraphs:
                paragraphs.alignment = PP_ALIGN.LEFT
                for run in paragraphs.runs:
                    run.font.size = self.table_font_size
            row += 1

    def _add_column_name(self, cell, text, alignment='left', bold=False, italic=False, level=False):
        ''' Insert text into row as column name. 
        
        :param cell: cell that should be modified
        :type cell: cell shape
        :param text: the text
        :type text: str
        :param bold: text should be bold
        :type bold: boolean
        :param italic: text should be italic
        :type italic: boolean
        :param aligment: how the text should be aligned (Default: left)
        :type alignment: str
        '''
        # Set alignment
        if alignment == 'left':
            alignment = PP_ALIGN.LEFT
        elif alignment == 'right':
            alignment = PP_ALIGN.RIGHT
        elif alignment == 'center':
            alignment = PP_ALIGN.CENTER

        cell.text = text
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = alignment
            # This will make text indented
            if level:
                paragraph.level = 1
            for run in paragraph.runs:
                run.font.size = self.table_font_size
                run.font.bold = bold
                run.font.italic = italic

    def _add_textbox(self, specs, slide, text, bold=False, italic=False, underline=False):
        ''' Add textbox with explanation text into the presentation. 
        
        :param specs: the position of the textbox
        :type specs: dict
        :param slide: slide to which textbox should be add
        :type slide: slide shape
        :param text: textbox's text
        :type text: str
        :param bold: make text bold
        :type bold: boolean (default: False)
        :param italic: make text italic
        :type italic: boolean (default: False)
        :param underline: make text underlined
        :type underline: boolean (default: False)
        '''
        txBox = slide.shapes.add_textbox(specs['left'], specs['top'], specs['width'], specs['height'])
        txBox.text_frame.clear()
        txBox.text_frame.word_wrap = True

        p = txBox.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.JUSTIFY
        run = p.add_run()
        run.font.size = Pt(11)
        run.font.underline = underline
        run.font.bold = bold
        run.font.italic = italic
        run.text = text

    def _set_transparency(self, transparency, elm):
        """ The function set the transparency of the row. 

        :param transparency: the transparency in %
        :type transparency: int
        :param elm: the element which transparency should be changed
        :type elm: format.line.color._xFill
        """
        a = str(100 - transparency) + '196'
        
        alpha = OxmlElement('a:alpha')
        alpha.set('val', a)
        elm.srgbClr.append(alpha)


    def _add_run(self, txtBox, text, bold=False, italic=False):
        ''' Add paragraph to the textbox. Each paragraph can have multiple runs. 
        
        :param txtBox: textbox to which add teh paragraph
        :type txtBox: textBox shape
        :param text: text to be insert into txtBox
        :type text: str
        :param bold: true if text should be boolean
        :type bold: boolean
        '''
        p = txtBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(11)
        run.font.bold = bold
        run.font.italic = italic

    def _get_column_values(self, column_name, baseline):
        ''' Return list of columns to be geenrated in the table based on the report type. 
        
        :param column_name: name of column to be included in the results
        :type column_name: str
        :param baseline: True if baseline report is generate
        :type baseline: boolean
        :returns: list of values
        '''
        if baseline:
            columns = [
                str(self.pre_stats[column_name].iloc[0]), 
                str(self.pre_stats[self._get_percentage_column_name(column_name)].iloc[0])]
        else:
            columns = [
                str(self.pre_stats[self._get_percentage_column_name(column_name)].iloc[0]),
                str(self.post_stats[self._get_percentage_column_name(column_name)].iloc[0])
            ]
        return columns

    def _create_table(self, slide, table_specs, title, trow, tcol, baseline=False):
        ''' Generate table in the report. 
        
        :param slide: the slide to which the table should be added
        :type slide: slide
        :param table_specs: specify a position of thetable
        :type table_specs: dict
        :param title: the name of the table
        :type title: str
        :param trow: number of rows of the table
        :type trow: int
        :param tcol: number of columns of the table
        :type tcol: int
        :param baseline: True if baseline report should be generated
        :type baseline: boolean
        '''
        # Add table to the slide, we need table with 21 rows and 3 columns
        shape = slide.shapes.add_table(
            trow, #rows 
            tcol, # cols
            table_specs['left'], 
            table_specs['top'],
            table_specs['width'],
            table_specs['height'])

        # Set look of the table
        # Change table style (https://github.com/scanny/python-pptx/issues/27)
        # This is a bit complicated how to change table style, I always do it that I select table style in powerpoint, save powerpoint and open xml file. There you can find style ID.
        style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
        tbl = shape._element.graphic.graphicData.tbl
        tbl[0][-1].text = style_id
        table = shape.table

        # Set column width based on the report type
        if baseline:
            widths = [
                int(table_specs['width'] * 0.5), 
                int(table_specs['width'] * 0.2), 
                int(table_specs['width'] * 0.15), 
                int(table_specs['width'] * 0.15)
                ]
        else:
            widths = [
                int(table_specs['width'] * 0.58), 
                int(table_specs['width'] * 0.2), 
                int(table_specs['width'] * 0.11), 
                int(table_specs['width'] * 0.11)
                ]
        for i in range(0, len(table.columns)):
            table.columns[i].width = widths[i]

        # Merge header row
        cell = self._merge_cells(table, 0, 0, 0, len(table.columns) - 1)
        self._add_column_name(cell, title, alignment='center', bold=True)

        # 2nd row
        nrow = 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        self._add_column_name(cell, "")
        # Set the header of table based on report type
        columns = ['n', '%'] if baseline else ['Pre n (%)', 'Post n (%)']
        self._insert_values(table=table, row=nrow, values=columns, bold=True)

        # Patient records entered
        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Patient records entered'
        self._add_column_name(cell, text)

        columns = [
            str(self.pre_stats['n'].iloc[0]), 
            '100'
            ] if baseline else [
                str(self.pre_stats['n'].iloc[0]),
                 str(self.post_stats['n'].iloc[0])
                 ]
        self._insert_values(table=table, row=nrow, values=columns)

        # Temperature monitoring and treatment
        nrow += 1
        cell = self._merge_cells(table, nrow, 0,nrow, 1)
        text = 'Temperature monitoring and treatment'
        self._add_column_name(cell, text, bold=True, italic=True)

        name = '# Temperature monitoring and treatment'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow + 2, 0)
        text = 'Temperature monitored at least four times \nper day #'
        self._add_column_name(cell, text, bold=True, italic=True)

        columns = [
            'Day of admission',
            'Day two of admission',
            'Day three of admission',
        ]
        self._insert_subrows(table=table, row=nrow, col=1, values=columns)

        name = '# Temperature monitored at least four times per day - Day of admission'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        name = '# Temperature monitored at least four times per day - Day two of admission'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        name = '# Temperature monitored at least four times per day - Day three of admission'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        # Temperature > 37.5°C recorded within 72 hours of admission
        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Temperature > 37.5°C recorded within 72 hours of admission'
        self._add_column_name(cell, text, alignment='left')

        name = '# Temperature > 37.5°c recorded within 72 hours of admission'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        # Paracetamol (or other anti-pyretic) given for first temperature > 37.5°C
        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Paracetamol (or other anti-pyretic) given for first temperature > 37.5°C'
        self._add_column_name(cell, text, alignment='left', level=True)

        name = '# Paracetamol (or other anti-pyretic) given for first temperature > 37.5°C'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        # Paracetamol (or other anti-pyretic) given with one hour from first temperature > 37.5°C #
        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Paracetamol (or other anti-pyretic) given with one hour from first temperature > 37.5°C #'
        self._add_column_name(cell, text, alignment='left', bold=True, level=True)

        name = '# Paracetamol (or other anti-pyretic) given with one hour from first temperature > 37.5°C'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        # Blood glucose monitoring and treatment
        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Blood glucose monitoring and treatment'
        self._add_column_name(cell, text, alignment='left', bold=True, italic=True)

        name = '# Blood glucose monitoring and treatment'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        # Blood Glucose Level (BGL) monitored > four times per day
        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow + 1, 0)
        text = 'Blood Glucose Level (BGL) monitored > four times per day'
        self._add_column_name(cell, text, alignment='left', bold=True)

        columns = [
            'Day of admission',
            'Day two of admission',
        ]
        self._insert_subrows(table=table, row=nrow, col=1, values=columns)
        
        name = '# Blood Glucose Level (BGL) monitored > four times per day - Day of admission'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        name = '# Blood Glucose Level (BGL) monitored > four times per day - Day two of admission'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'BGL ≥ 10mmol/L within 48 hours of admission'
        self._add_column_name(cell, text, alignment='left')

        name = '# BGL ≥ 10mmol/L within 48 hours of admission'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Insulin given for first BGL ≥ 10mmol/L'
        self._add_column_name(cell, text, alignment='left', bold=True, level=True)

        name = '# Insulin given for first BGL ≥ 10mmol/L'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Insulin given within one hour from first BGL ≥ 10mmol/L #'
        self._add_column_name(cell, text, alignment='left', bold=True, level=True)

        name = '# Insulin given within one hour from first BGL ≥ 10mmol/L'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Swallow screening'
        self._add_column_name(cell, text, alignment='left', bold=True, italic=True)

        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Formal swallow screen performed'
        self._add_column_name(cell, text, alignment='left')

        name = '# Formal swallow screen performed'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Swallow screen performed within 24 hours #'
        self._add_column_name(cell, text, alignment='left', bold=True, level=True)

        name = '# Swallow screen performed within 24 hours'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Swallow screen or swallow assessment performed before being given oral medications #'
        self._add_column_name(cell, text, alignment='left', bold=True)

        name = '# Swallow screen or swallow assessment performed before being given oral medications'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        nrow += 1
        cell = self._merge_cells(table, nrow, 0, nrow, 1)
        text = 'Swallow screen or swallow assessment performed before being given oral food or fluids #'
        self._add_column_name(cell, text, alignment='left', bold=True)

        name = '# Swallow screen or swallow assessment performed before being given oral food or fluids'
        columns = self._get_column_values(name, baseline)
        self._insert_values(table=table, row=nrow, values=columns)

        # Create iterator through cells in the table
        def iter_cells(table):
            for row in table.rows:
                for cell in row.cells:
                    yield cell

        # Set font size of all cells in the table based on the report type
        for cell in iter_cells(table):
            cell.text_frame.autosize = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = self.table_font_size

    def generate_baseline_report(self, df=None):
        ''' Generate baseline data summary feedback. 
        
        :param df: dataframe to be used for values (default: None)
        :type df: dataframe
        '''

        self.table_font_size = Pt(11) # Set default font size of the table in baseline report
        # Define template for the presentation
        master = os.path.normpath(os.path.join(os.path.dirname(__file__), 'backgrounds', 'qasc_baseline.pptx'))

        # Filter dataframe for site       
        hospital_name = self.study_df.loc[self.study_df['unique_identifier'] == self.site_id, 'facility_name'].iloc[0]

        # Create output filename containing qasc, current date and site ID
        output_file = f'qasc_{self.site_id}_{datetime.now().strftime("%Y-%m-%d")}.pptx'

        # set main text to be added into reports
        main_texts = [
            'Congratulations on completing your baseline audit. We have summarized the results for you in the table below. Please share these results with your team. These data can assist you when discussing the barriers and enablers to implementation of the FeSS clinical protocols at your hospital.', 
            'It is important to please let us know if there are problems with the data that can be explained further (eg. was there a question the people entering data may not have understood properly?)',
            'Please don’t hesitate to contact the NRI if you require clarification on any of the items above.'
        ]
        
        # Create new Presentaiton object
        prs = Presentation(master)

        # Get first slide
        first_slide = prs.slides[0]

        # Set specification of the table
        table_specs = {
            'height': Cm(18),
            'width': Cm(19),
            'left': Cm(1),
            'top': Cm(5)
        }

        # Create table
        self._create_table(
            slide=first_slide, 
            table_specs=table_specs, 
            title=f'Table 1: FeSS Management for {hospital_name}', 
            trow=21, 
            tcol=4,
            baseline=True)

        # Add the rest of explaining texts
        specs = {
            0: {
                'height': Cm(2),
                'width': Cm(19),
                'left': Cm(1),
                'top': Cm(3)
            }, 
            1: {
                'height': Cm(1),
                'width': Cm(19),
                'left': Cm(1),
                'top': Cm(25.5)
            },
            2: {
                'height': Cm(1),
                'width': Cm(19),
                'left': Cm(1),
                'top': Cm(27)
            },
        }

        for i in range(0, len(main_texts)):
            self._add_textbox(specs[i], first_slide, main_texts[i])

        # Create graph on the second slide
        second_slide = prs.slides.add_slide(prs.slide_layouts[0])
        graph_df = df[[
            '% Temperature monitored at least four times per day - Day of admission',
            '% Paracetamol (or other anti-pyretic) given with one hour from first temperature > 37.5°C',
            '% Blood Glucose Level (BGL) monitored > four times per day - Day of admission',
            '% Insulin given within one hour from first BGL ≥ 10mmol/L',
            '% Swallow screen performed within 24 hours',
        ]].copy()
        new_column_names = ["Temp (Day 1)", "Paracetamol (1hr)", "BGL's (Day 1)", "Insulin (1hr)", "Swallow screen (24hrs)"]
        graph_df.rename(columns=dict(zip(graph_df.columns, new_column_names)),inplace=True)

        column_name = 'Baseline audit'
        graph_df = graph_df.T.rename(columns={0: column_name})

        # Create chart data
        chart_data = ChartData()
        chart_data.categories = new_column_names
        chart_data.add_series(column_name, graph_df[column_name].tolist())     

        # Add chart on slide
        specs = {
            'height': Cm(10),
            'width': Cm(19),
            'left': Cm(1),
            'top': Cm(3)
            }
        chart = second_slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart         


        plot = chart.plots[0]
        # All bars with the same color
        plot.vary_by_categories = False

        # Set maximum to 100
        value_axis = chart.value_axis
        value_axis.maximum_scale = 100

        value_axis.major_gridlines.format.line.width = Pt(0.5)
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(206, 206, 206) # Set color to gray (A6A6A6)

        value_axis.format.line.color.rgb = RGBColor(0, 0, 0)
        solidFill = value_axis.format.line.color._xFill
        self._set_transparency(100, solidFill)

        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        # Set 100% transparency to category axis
        category_axis.format.line.color.rgb = RGBColor(206, 206, 206)
        solidFill = category_axis.format.line.color._xFill
        self._set_transparency(100, solidFill)

        # Set graph of title
        graph_title = f'Figure 1: FeSS Management {hospital_name} Hospital'
        chart_text = chart.chart_title.text_frame
        chart_text.text = graph_title
        chart_text.paragraphs[0].font.size = Pt(12)
        chart_text.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

        # Save presentation
        path = os.path.join(os.getcwd(), output_file)
        save_file(output_file)
        prs.save(path)

        
    def generate_pre_post_report(self):
        ''' Generate report with pre/post comparison. '''

        # Set smaller font size of the table
        self.table_font_size = Pt(9.5)
        
        # Define template
        master = os.path.normpath(os.path.join(os.path.dirname(__file__), 'backgrounds', 'qasc_comparison.pptx'))
   
        # Get hospital name based on study ID
        hospital_name = self.study_df.loc[self.study_df['unique_identifier'] == self.site_id, 'facility_name'].iloc[0]

        # Create output filename containing qasc, current date and site ID
        output_file = f'qasc_comp_{self.site_id}_{datetime.now().strftime("%Y-%m-%d")}.pptx'

        # Create Presentation object
        prs = Presentation(master)
        # Get first slide
        first_slide = prs.slides[0]

        # Add title 
        title_text = f'QASC Europe Project: Post-Intervention audit summary {hospital_name} Hospital'
        specs = {
            'height': Cm(1),
            'width': Cm(18),
            'left': Cm(0.6),
            'top': Cm(1.5),
        }
        self._add_textbox(specs, first_slide, title_text, bold=True, underline=True)

        # Add first paragraph with explanation and congratulations
        # A bit longer code because only some letters was made bold, so I had to created more runs in paragraph. 
        specs = {
            'height': Cm(2),
            'width': Cm(19.5),
            'left': Cm(0.6),
            'top': Cm(2),
        }
        txBox = first_slide.shapes.add_textbox(specs['left'], specs['top'], specs['width'], specs['height'])
        txBox.text_frame.clear()
        txBox.text_frame.word_wrap = True
        self._add_run(
            txBox, 
            'Congratulations on completing the QASC Europe project audits on the use of the FeSS (',
        )
        self._add_run(txBox, 'F', bold=True)
        self._add_run(txBox, 'ever, ')
        self._add_run(txBox, 'S', bold=True)
        self._add_run(txBox, 'ugar, and ')
        self._add_run(txBox, 'S', bold=True)
        self._add_run(txBox, 'wallowing) protocols for stroke patients. The summaries below reflect ')
        self._add_run(txBox, f"your hospital’s performance for the {self.pre_stats['n'].iloc[0]} stroke patients you reviewed for the baseline audit XX/XX/XXXX and the {self.post_stats['n'].iloc[0]} patients you reviewed during the post intervention period XX/XX/XXXX. ", bold=True)
        self._add_run(txBox, ' We present the number of patients audited (n) and the proportion of patients who met criteria (%).')

        # Specification of the table position
        table_specs = {
            'height': Cm(16),
            'width': Cm(19),
            'left': Cm(0.8),
            'top': Cm(4.0)
        }
        # Create table
        self._create_table(
            slide=first_slide, 
            table_specs=table_specs, 
            title=f'FeSS Management for {hospital_name} Hospital', 
            trow=21, 
            tcol=4, 
            baseline=False)

        # Add explanation text below table 
        text = "# Variables shown in bold above are the key recommendations to be followed in the QASC Europe project"
        specs = {
            'height': Cm(0.5),
            'width': Cm(18),
            'left': Cm(0.6),
            'top': Cm(20),
        }
        self._add_textbox(specs, first_slide, text)

        # Second slide
        second_slide = prs.slides[1]

        # Add temperature block
        column_name = '# Temperature monitoring and treatment'
        pre_temp = self.pre_stats[self._get_percentage_column_name(column_name)].iloc[0]
        post_temp = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]
        difference = "no difference" if post_temp <= pre_temp else "an increase"

        column_name = '# Temperature monitored at least four times per day - Day of admission'
        pre_temp_first = self.pre_stats[self._get_percentage_column_name(column_name)].iloc[0]
        post_temp_first = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]

        column_name = '# Temperature monitored at least four times per day - Day two of admission'
        pre_temp_second = self.pre_stats[self._get_percentage_column_name(column_name)].iloc[0]
        post_temp_second = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]

        column_name = '# Temperature monitored at least four times per day - Day three of admission'
        pre_temp_third = self.pre_stats[self._get_percentage_column_name(column_name)].iloc[0]
        post_temp_third = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]

        column_name = '# Temperature > 37.5°c recorded within 72 hours of admission'
        post_temp_rec_total = self.post_stats[column_name].iloc[0]

        column_name = '# Paracetamol (or other anti-pyretic) given with one hour from first temperature > 37.5°C'
        post_paracetamol = self.post_stats[column_name].iloc[0]
        post_paracetamol_perc = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]      

        text = "1. Temperature"
        specs = {
            'height': Cm(0.7),
            'width': Cm(28.5),
            'left': Cm(0.6),
            'top': Cm(1.5),
        }
        self._add_textbox(specs, second_slide, text, bold=True)

        proportion_text = f"i) From pre to post implementation there was {difference} in the proportion of patients who: were monitored four times a day for fever on Day 1 ({pre_temp_first}% vs {post_temp_first}% respectively), Day 2 ({pre_temp_second}% vs {post_temp_second}%) and Day 3 ({pre_temp_third}% vs {post_temp_third}%)."
        specs = {
            'height': Cm(1.2),
            'width': Cm(27.3),
            'left': Cm(1.8),
            'top': Cm(2),
        }
        self._add_textbox(specs, second_slide, proportion_text)

        paracetamol_text = f"ii) {post_paracetamol} out of {post_temp_rec_total} ({post_paracetamol_perc}%) patients with a temperature > 37.5°C received Paracetamol (or alternatively antipyretic) within one hour."
        specs = {
            'height': Cm(0.7),
            'width': Cm(27.3),
            'left': Cm(1.8),
            'top': Cm(3),
        }
        self._add_textbox(specs, second_slide, paracetamol_text)

        # Add blood glucose monitoring
        column_name = '# Blood glucose monitoring and treatment'
        pre_blood = self.pre_stats[self._get_percentage_column_name(column_name)].iloc[0]
        post_blood = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]
        difference = "no difference" if post_blood <= pre_blood else "an increase"

        column_name = '# Blood Glucose Level (BGL) monitored > four times per day - Day of admission'
        pre_blood_first = self.pre_stats[self._get_percentage_column_name(column_name)].iloc[0]
        post_blood_first = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]

        column_name = '# Blood Glucose Level (BGL) monitored > four times per day - Day two of admission'
        pre_blood_second = self.pre_stats[self._get_percentage_column_name(column_name)].iloc[0]
        post_blood_second = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]

        column_name = '# Insulin given for first BGL ≥ 10mmol/L'
        post_insulin_total = self.post_stats[column_name].iloc[0]

        column_name = '# Insulin given within one hour from first BGL ≥ 10mmol/L'
        post_insulin = self.post_stats[column_name].iloc[0]
        post_insulin_perc = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]      

        text = "2. Blood glucose monitoring"
        specs = {
            'height': Cm(0.7),
            'width': Cm(28.5),
            'left': Cm(0.6),
            'top': Cm(4.0),
        }
        self._add_textbox(specs, second_slide, text, bold=True)

        proportion_text = f"iii) {difference.capitalize()} in the proportion of patients observed four times daily for hyperglycaemia on Day 1 ({pre_blood_first}% vs {post_blood_first}% respectively) and Day 2 ({pre_blood_second}% vs {post_blood_second}%)."
        specs = {
            'height': Cm(1.2),
            'width': Cm(27.3),
            'left': Cm(1.8),
            'top': Cm(4.5),
        }
        self._add_textbox(specs, second_slide, proportion_text)

        paracetamol_text = f"iv) {post_insulin} out of {post_insulin_total} ({post_insulin_perc}%) patients with BGL's > 10mmol/L were treated with insulin."
        specs = {
            'height': Cm(0.7),
            'width': Cm(27.3),
            'left': Cm(1.8),
            'top': Cm(5.0),
        }
        self._add_textbox(specs, second_slide, paracetamol_text)

        # Add swallow screening
        column_name = '# Swallow screen performed within 24 hours'
        pre_swallow = self.pre_stats[self._get_percentage_column_name(column_name)].iloc[0]
        post_swallow = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]
        difference = "no difference" if post_swallow <= pre_swallow else "an increase"

        column_name = '# Swallow screen or swallow assessment performed before being given oral medications'
        post_oral_med = self.post_stats[column_name].iloc[0]
        post_oral_med_perc = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]      

        column_name = '# Swallow screen or swallow assessment performed before being given oral food or fluids'
        post_food = self.post_stats[column_name].iloc[0]
        post_food_perc = self.post_stats[self._get_percentage_column_name(column_name)].iloc[0]      

        text = "3. Swallow screening"
        specs = {
            'height': Cm(0.7),
            'width': Cm(28.5),
            'left': Cm(0.6),
            'top': Cm(6.0),
        }
        self._add_textbox(specs, second_slide, text, bold=True)

        swallow_text = f"v) {difference.capitalize()} in the proportion of patients who received a swallow screen within 24hrs ({pre_swallow}% vs {post_swallow}%)."
        specs = {
            'height': Cm(0.7),
            'width': Cm(27.3),
            'left': Cm(1.8),
            'top': Cm(6.5),
        }
        self._add_textbox(specs, second_slide, swallow_text)

        medication_text = f"vi) {post_oral_med} ({post_oral_med_perc}%) patients received oral medications before having their swallow screened and {post_food} ({post_food_perc}%) patients received oral food or fluid before having their swallow screened."
        specs = {
            'height': Cm(1.2),
            'width': Cm(27.3),
            'left': Cm(1.8),
            'top': Cm(7),
        }
        self._add_textbox(specs, second_slide, medication_text)

        # Summary box
        specs = {
            'height': Cm(2),
            'width': Cm(27.3),
            'left': Cm(0.6),
            'top': Cm(8.5),
        }
        txBox = second_slide.shapes.add_textbox(specs['left'], specs['top'], specs['width'], specs['height'])
        txBox.text_frame.clear()
        txBox.text_frame.word_wrap = True
        self._add_run(
            txBox, 
            'SUMMARY: ',
            bold=True
        )

        temp_monitoring = self.post_stats['# monitored for fever and/or hyperglycaemia four times a day'].iloc[0]
        acute_stroke_pts = self.post_stats['# acute stroke'].iloc[0]
        self._add_run(txBox, f"Monitoring for fever and/or hyperglycaemia four times a day for the first 48 - 72 hours of admission was recorded for {temp_monitoring} patient of the {acute_stroke_pts} acute stroke patients audited. A swallow screen or assessment was performed for {post_swallow} patient. Most of those patients who required treatment for fever and hyperglycaemia had Paracetamol (or other anti-pyretic) or Insulin administered within the recommended one-hour time period.", italic=True)

        # Add comparison graph
        data = {
            'Criterium': ["Temp (Day 1)", "BGL's (Day 1)", "Swallow screen first (24hrs)"], 
            'Pre': [pre_temp_first, pre_blood_first, pre_swallow],
            'Post': [post_temp_first, post_blood_first, post_swallow]}
        graph_data = pd.DataFrame(data=data)
        column_names = ["Pre", "Post"]

        chart_data = ChartData()
        chart_data.categories = graph_data['Criterium'].tolist()

        for val in column_names:
            chart_data.add_series(val, graph_data[val].tolist())

        # Add chart on slide
        specs = {
            'height': Cm(9.5),
            'width': Cm(15),
            'left': Cm(0.6),
            'top': Cm(10)
            }
        chart = second_slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, specs['left'],specs['top'], specs['width'],specs['height'], chart_data).chart         

        plot = chart.plots[0]
        # All bars with the same color
        plot.vary_by_categories = False

        # Set colors of each serie
        colors = [
            RGBColor(192, 80, 77), 
            RGBColor(79, 129, 189)
            ]    

        # Iterate over column names and set color of the series and add labels
        for i in range(0, len(column_names)):
            series = chart.series[i]
            fill = series.format.fill

            # Set colors of series 
            fill.solid()
            fill.fore_color.rgb = colors[i]

            # Show data labels and set font
            plot.has_data_labels = True
            # Set layout of labels
            data_labels = plot.data_labels
            data_labels.font.size = Pt(9)

        # Set maximum to 100
        value_axis = chart.value_axis
        value_axis.maximum_scale = 100

        # Show major gridlines
        value_axis.major_gridlines.format.line.width = Pt(0.5)
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(206, 206, 206) # Set color to gray (A6A6A6)

        value_axis.format.line.color.rgb = RGBColor(0, 0, 0)
        solidFill = value_axis.format.line.color._xFill
        self._set_transparency(100, solidFill)

        # Value for y-axis (change font size, name, and other things)
        category_axis = chart.category_axis
        # Set 100% transparency to category axis
        category_axis.format.line.color.rgb = RGBColor(206, 206, 206)
        solidFill = category_axis.format.line.color._xFill
        self._set_transparency(100, solidFill)

        # Set graph of title
        graph_title = f'Figure 1: Hospital {hospital_name} Pre/Post FeSS intervention'
        chart_text = chart.chart_title.text_frame
        chart_text.text = graph_title
        chart_text.paragraphs[0].font.size = Pt(12)
        chart_text.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

        # Set legend 
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)

        # Save presentation
        path = os.path.join(os.getcwd(), output_file)
        save_file(output_file)
        prs.save(path)