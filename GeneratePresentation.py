# -*- coding: utf-8 -*-
"""
File name: GeneratePresentation.py
Package: resq
Written by: Marie Jankujova - jankujova.marie@fnusa.cz on 11-2017
Version: v1.0
Version comment: The first version of script which calls class GenerateGraphs and create graphs into presentation. 
Description: This script is used to call class GenerateGraphs and produce graphs into presentation. 
"""


import pandas as pd
import sys
import os
from datetime import datetime, date
import sqlite3
import pytz
from resqdb.GenerateGraphs import GenerateGraphs, GenerateGraphsQuantiles, GenerateGraphsSites
import xlsxwriter
from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor

class GeneratePresentation:
    """ The class generating the general presentation for countries and sites. 

    :param df: the dataframe with calculated statistics
    :type df: pandas dataframe
    :param country: `True` if country is included in the statistics as site
    :type country: bool
    :param country_code: the country code
    :type country_code: str
    :param split_sites: `True` if presentation should be generated per sites seperately
    :type split_sites: bool
    :param site: the site code
    :type site: str
    :param report: the type of the report eg. quarter
    :type report: str
    :param quarter: the type of the period eg. Q1_2019
    :type quarter: str
    """

    def __init__(self, df, country=False, country_code=None, split_sites=False, site=None, report=None, quarter=None, country_name=None):

        self.df = df.drop_duplicates(subset=['Site ID', 'Total Patients'], keep='first')
        self.country_code = country_code
        self.report = report
        self.quarter = quarter

        #master_pptx = self.country_code + ".pptx"
        script_dir = os.path.dirname(__file__) 
        master_pptx = "master.pptx"
        self.master = os.path.normpath(os.path.join(script_dir, "backgrounds", master_pptx))

        # Connect to database and get country name according to country code.
        def select_country(value):
            """ The function obtaining the name of country from the package pytz based on the country code. 

            :param value: the country code
            :type value: str
            :returns: the country name
            """
            country_name = pytz.country_names[value]
            return country_name

        # If country is used as site, the country name is selected from countries dictionary by country code. :) 
        if (country):
            self.country_name = country_name
        else:
            self.country_name = country_name

        # If split_sites is True, then go through dataframe and generate graphs for each site (the country will be included as site in each file).
        site_ids = self.df['Site ID'].tolist()
        # Delete country name from side ids list.
        try:
            site_ids.remove(self.country_name)
        except:
            pass

        if site is not None:
            df = self.df[self.df['Site ID'].isin([site, self.country_name])].copy()
            self._generate_graphs(df=df, site_code=site)

        # Generate formatted statistics for all sites individualy + country as site is included
        if (split_sites) and site is None:
            for i in site_ids:
                df = self.df[self.df['Site ID'].isin([i, self.country_name])].copy()
                self._generate_graphs(df=df, site_code=i)
    
        # Produce formatted statistics for all sites + country as site
        if site is None:
            self._generate_graphs(df=self.df, site_code=country_code)

    def _generate_graphs(self, df, site_code=None):
        """ The function opening the presentation and generating graphs. 
        
        :param df: the dataframe with calculated statistic
        :type df: pandas dataframe
        :param site_code: the site ID
        :type site_code: str
        """
        
        prs = Presentation(self.master)

        first_slide = prs.slides[0]
        shape = first_slide.shapes[5]
        text_frame = shape.text_frame

        if self.country_name is None:
            first_slide_text = "Data Summary"
        else:
            first_slide_text = self.country_name + "\nData Summary"
        #first_slide_text = "\nData Summary"

        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = first_slide_text

        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(24)
        font.color.rgb = RGBColor(250,250,250)

        # if (self.country_name in ['Ukraine', 'Poland'] and len(df) > 2):
        #     main_col = 'Site ID'
        # else:
        main_col =  'Site Name'

        ########################
        #### TOTAL PATIENTS ####
        ########################
        column_name = 'Total Patients'
        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        #title = 'TOTAL PATIENTS'
        
        country_patients = str(max(tmp_df[column_name].tolist()))
        title = 'TOTAL PATIENTS (n = {})'.format(country_patients)
        if self.country_name is not None:
            tmp_df = tmp_df.loc[tmp_df[main_col] != self.country_name]
        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##########################
        ### MEDIAN PATIENT AGE ###
        ##########################
        column_name = 'Median patient age'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "MEDIAN PATIENT AGE"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###########################
        ### GENDER DISTRIBUTION ###
        ###########################
        column_name = '% patients female'
        legend = ['Female', 'Male']

        tmp_df = df[[main_col, '% patients female', '% patients male']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "GENDER DISTRIBUTION"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ############################################
        ### % RECURRENT STROKES OUT OF ALL CASES ###
        ############################################
        '''
        sheet_name = 'recurrent_stroke'
        column_name = '% recurrent stroke - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% RECURRENT STROKES out of all cases"

        GenerateGraphs(dataframe=tmp_df, title=title, column_name=column_name, sheet_name=sheet_name, writer=writer, workbook=workbook)
        '''

        #######################
        ### DEPARTMENT TYPE ###
        #######################
        column_name = '% department type - neurology'
        legend = ['neurology', 'neurosurgery', 'anesthesiology resuscitation critical care', 'internal medicine', 'geriatrics', 'other']

        tmp_df = df[[main_col, '% department type - neurology', '% department type - neurosurgery', '% department type - anesthesiology/resuscitation/critical care', '% department type - internal medicine', '% department type - geriatrics', '% department type - Other']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% DEPARTMENT TYPE ALLOCATION out of all cases" 

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ###################################
        ### HOSPITALIZATION DESTINATION ###
        ###################################
        column_name = '% patients hospitalized in stroke unit / ICU'
        legend = ['stroke unit', 'monitored bed with telemetry', 'standard bed']

        tmp_df = df[[main_col, '% patients hospitalized in stroke unit / ICU', '% patients hospitalized in monitored bed with telemetry', '% patients hospitalized in standard bed']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% HOSPITALIZATION DESTINATION out of all cases" 

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ###############
        # STROKE TYPE #
        ###############
        column_name = '% stroke type - ischemic stroke'
        legend = ['ischemic', 'transient ischemic attack', 'intracerebral hemorrhage', 'subarrachnoid hemorrhage', 'cerebral venous thrombosis', 'undetermined']

        tmp_df = df[[main_col, '% stroke type - ischemic stroke', '% stroke type - transient ischemic attack', '% stroke type - intracerebral hemorrhage', '% stroke type - subarrachnoid hemorrhage', '% stroke type - cerebral venous thrombosis', '% stroke type - undetermined stroke']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% STROKE TYPE out of all cases"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ######################
        # CONSIOUSNESS LEVEL #
        ######################
        column_name = 'alert_all_perc'
        legend = ['alert', 'drowsy', 'comatose']

        tmp_df = df[[main_col, 'alert_all_perc', 'drowsy_all_perc', 'comatose_all_perc']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CONSCIOUSNESS LEVEL for IS, ICH, CVT, SAH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        '''
        #######
        # GCS #
        #######
        column_names = ['% GCS - 15-13', '% GCS - 12-8', '% GCS - <8']
        column_name = '% GCS - 15-13'
        legend = ['GCS 15 - 13', 'GCS 12 - 8', 'GCS < 8']

        tmp_df = df[[main_col, '% GCS - 15-13', '% GCS - 12-8', '% GCS - <8']]
        tmp_df = tmp_df.sort_values(column_names, ascending = [True, True, True])

        title = "% GLASGOW COMA SCALE INTERVALS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')
        '''
        ###################
        # NIHSS PERFORMED #
        ###################
        column_name = '% NIHSS - Performed'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% NIHSS PERFORMED for IS, ICH, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ######################
        # NIHSS MEDIAN SCORE #
        ######################
        column_name = 'NIHSS median score'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "NIHSS median score"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ####################
        # CT/MRI performed #
        ####################
        column_name = '% CT/MRI - performed'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CT/MRI PERFORMED for IS, ICH, CVT, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################
        # CT/MRI PERFORMED WITHIN 1 HOUR #
        ##################################
        column_name = '% CT/MRI - Performed within 1 hour after admission'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CT/MRI PERFORMED WITHIN 1 HOUR AFTER ADMISSION for IS, ICH, CVT, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ####################
        # VASCULAR IMAGING #
        ####################
        sorted_by = ['vascular_imaging_cta_norm', 'vascular_imaging_mra_norm', 'vascular_imaging_dsa_norm', 'vascular_imaging_none_norm']
        column_name = 'vascular_imaging_cta_norm'
        legend = ['CTA', 'MRA', 'DSA', 'none']

        tmp_df = df[[main_col, 'vascular_imaging_cta_norm', 'vascular_imaging_mra_norm', 'vascular_imaging_dsa_norm', 'vascular_imaging_none_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% VASCULAR IMAGING PERFORMED for ICH, SAH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')


        ############################
        # RECANALIZATION TREATMENT #
        ############################
        column_name = '% recanalization procedures - IV tPa'
        legend = ['IV tPa', 'IV tPa + endovascular treatment', 'IV tPa + another centre for endovascular treatment']

        tmp_df = df[[main_col, '% recanalization procedures - IV tPa', '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% IV tPa for IS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')


        ####################################################
        # RECANALIZATION TREATMENT IN COMPREHENSIVE CENTRE #
        ####################################################
        sorted_by = ['% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone']
        column_name = '% recanalization procedures - IV tPa + endovascular treatment'
        legend = ['IV tPa + endovascular treatment', 'endovascular treatment']

        tmp_df = df[[main_col, '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% RECANALIZATION PROCEDURES IN COMPREHENSIVE CENTRES for IS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################################################
        # RECANALIZATION TREATMENT IN REFFERED TO ANOTHER CENTRE #
        ##########################################################
        sorted_by = ['% recanalization procedures - IV tPa + referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment and hospitalization continues at the referred to centre', '% recanalization procedures - Referred for endovascular treatment and patient is returned to the initial centre']
        column_name = '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment'
        legend = ['IV tPa + another centre for endovascular treatment', 'another centre for endovascular treatment', 'another centre for endovascular treatment and hospitalization continues', 'another centre for endovascular treatment and returned to the initial centre']

        tmp_df = df[[main_col, '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment and hospitalization continues at the referred to centre', '% recanalization procedures - Referred for endovascular treatment and patient is returned to the initial centre']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% PATIENTS TRANSFERRED TO ANOTHER CENTRE FOR RECANALIZATION PROCEDURES FROM PRIMARY CENTRE for IS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ############################
        # RECANALIZATION TREATMENT #
        ############################
        column_name = '% recanalization procedures - IV tPa'
        legend = ['IV tPa', 'IV tPa + endovascular treatment', 'endovascular treatment', 'IV tPa + another centre for endovascular treatment']

        tmp_df = df[[main_col, '% patients recanalized', '% recanalization procedures - IV tPa', '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone', '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment']]
        tmp_df = tmp_df.sort_values(['% patients recanalized'], ascending = True)

        title = "% RECANALIZATION PROCEDURES for IS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################
        # % patients recanalized #
        ##########################
        column_name = '% patients recanalized'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% of recanalized patients"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)


        ################
        # % median DTN #
        ################
        column_name = 'Median DTN (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "MEDIAN DOOR-TO-NEEDLE TIME (minutes) for thrombolyzed patients"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###############
        #  median DTG #
        ###############
        column_name = 'Median DTG (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "MEDIAN DOOR-TO-GROIN TIME (minutes) for patients receiving endovascular treatment in a comprehensive centre"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ####################
        #  median TBY DIDO #
        ####################
        column_name = 'Median TBY DIDO (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "MEDIAN DOOR-IN-DOOR-OUT TIME (minutes) for patients referred from a primary centre to another centre for recanalization therapy"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #######################
        # dysphagia screening #
        #######################
        column_name = '% dysphagia screening - Guss test'
        column_names = ['% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre']
        legend = ['GUSS test', 'Other test', 'Another centre']

        tmp_df = df[[main_col, '% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre']]
        tmp_df = tmp_df.sort_values(column_names, ascending = True)

        title = "% DYSPHAGIA SCREENING PERFORMED for IS, ICH, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ######################################
        #  dypshagia screening within 1 hour #
        ######################################
        column_name = '% dysphagia screening time - Within first 24 hours'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% DYSPHAGIA SCREENING TIME WITHIN FIRST 24 HOURS AFTER ADMISSION"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        
        ###########################
        # PUT ON VENTILATOR - YES #
        ###########################
        column_name = '% patients put on ventilator - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS PUT ON VENTILATOR for IS, ICH, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###################
        # HEMICRANEOCTOMY #
        ###################
        column_name = '% hemicraniectomy - Yes'
        legend = ['Yes', 'Referred to another centre']

        tmp_df = df[[main_col, '% hemicraniectomy - Yes', '% hemicraniectomy - Referred to another centre']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% HEMICRANIECTOMY PERFORMED for IS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ################
        # NEUROSURGERY #
        ################
        column_name = '% neurosurgery - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% NEUROSURGERY PERFORMED for ICH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #######################################
        # NEUROSURGERY TYPE (FIRST 3 OPTIONS) #
        #######################################
        sorted_by = ['% neurosurgery type - intracranial hematoma evacuation', '% neurosurgery type - external ventricular drainage', '% neurosurgery type - decompressive craniectomy']
        column_name = '% neurosurgery type - intracranial hematoma evacuation'
        legend = ['intracranial hematoma evacuation', 'external ventricular drainage', 'decompressive craniectomy']

        tmp_df = df[[main_col, '% neurosurgery type - intracranial hematoma evacuation', '% neurosurgery type - external ventricular drainage', '% neurosurgery type - decompressive craniectomy']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% NEUROSURGERY TYPE PERFORMED for ICH in comprehensive centres"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ###################################
        # NEUROSURGERY TYPE (LAST OPTION) #
        ###################################
        column_name = '% neurosurgery type - Referred to another centre'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS REFERRED TO ANOTHER CENTRE FOR NEUROSURGERY"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###################
        # BLEEDING REASON #
        ###################
        sorted_by = ['bleeding_arterial_hypertension_perc_norm', 'bleeding_aneurysm_perc_norm', 'bleeding_arterio_venous_malformation_perc_norm', 'bleeding_anticoagulation_therapy_perc_norm', 'bleeding_amyloid_angiopathy_perc_norm', 'bleeding_other_perc_norm']
        column_name = 'bleeding_arterial_hypertension_perc_norm'
        legend = ['arterial hypertension', 'aneurysm', 'arterio-venous malformation', 'anticoagulation therapy', 'amyloid angiopathy', 'other']

        tmp_df = df[[main_col, 'bleeding_arterial_hypertension_perc_norm', 'bleeding_aneurysm_perc_norm', 'bleeding_arterio_venous_malformation_perc_norm', 'bleeding_anticoagulation_therapy_perc_norm', 'bleeding_amyloid_angiopathy_perc_norm', 'bleeding_other_perc_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% BLEEDING REASON for ICH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################
        # INTERVENTION PERFORMED #
        ##########################
        sorted_by = ['intervention_endovascular_perc_norm', 'intervention_neurosurgical_perc_norm', 'intervention_other_perc_norm', 'intervention_referred_perc_norm', 'intervention_none_perc_norm']
        column_name = 'intervention_endovascular_perc_norm'
        legend = ['Endovascular (coiling)', 'Neurosurgical (clipping)', 'Other neurosurgical treatment', 'Patient referred to another centre', 'None']

        tmp_df = df[[main_col, 'intervention_endovascular_perc_norm', 'intervention_neurosurgical_perc_norm', 'intervention_other_perc_norm', 'intervention_referred_perc_norm', 'intervention_none_perc_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% INTERVENTION PERFORMED for SAH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################################
        # % PATIENTS ASSESSED FOR REHABILITATION #
        ##########################################
        column_name = '% patients assessed for rehabilitation - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% REHABILITATION ASSESSMENT for IS, ICH, CVT and SAH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###############################
        # VENOUS THROMBOSIS TREATMENT #
        ###############################
        sorted_by = ['vt_treatment_anticoagulation_perc_norm', 'vt_treatment_thrombectomy_perc_norm', 'vt_treatment_local_thrombolysis_perc_norm', 'vt_treatment_local_neurological_treatment_perc_norm']
        column_name = 'vt_treatment_anticoagulation_perc_norm'
        legend = ['anticoagulation', 'thrombectomy', 'local thrombolysis', 'neurosurgical treatment']

        tmp_df = df[[main_col, 'vt_treatment_anticoagulation_perc_norm', 'vt_treatment_thrombectomy_perc_norm', 'vt_treatment_local_thrombolysis_perc_norm', 'vt_treatment_local_neurological_treatment_perc_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% VENOUS THROMBOSIS TREATMENT for CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ################################
        # ATRIAL FIBRILLATION DETECTED #
        ################################
        sorted_by = ['% afib/flutter - Detected during hospitalization', '% afib/flutter - Newly-detected at admission', '% afib/flutter - Known']
        column_name = '% afib/flutter - Detected during hospitalization'
        legend = ['detected during hospitalization', 'newly-detected at admission', 'known aFib']

        tmp_df = df[[main_col, '% afib/flutter - Detected during hospitalization', '% afib/flutter - Newly-detected at admission', '% afib/flutter - Known']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% ATRIAL FIBRILLATION DETECTED for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ########################################
        # ATRIAL FIBRILLATION DETECTION METHOD #
        ########################################
        sorted_by = ['% afib detection method - Telemetry with monitor allowing automatic detection of aFib', '% afib detection method - Telemetry without monitor allowing automatic detection of aFib', '% afib detection method - Holter-type monitoring', '% afib detection method - EKG monitoring in an ICU bed with automatic detection of aFib', '% afib detection method - EKG monitoring in an ICU bed without automatic detection of aFib']
        column_name = '% afib detection method - Telemetry with monitor allowing automatic detection of aFib'
        legend = ['Telemetry with monitoring', 'Telemetry without monitoring', 'Holter-type monitoring', 'EKG monitoring in an ICU bed with automatic detection of aFib', 'EKG monitoring in an ICU bed without automatic detection of aFib']

        tmp_df = df[[main_col, '% afib detection method - Telemetry with monitor allowing automatic detection of aFib', '% afib detection method - Telemetry without monitor allowing automatic detection of aFib', '% afib detection method - Holter-type monitoring', '% afib detection method - EKG monitoring in an ICU bed with automatic detection of aFib', '% afib detection method - EKG monitoring in an ICU bed without automatic detection of aFib']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% ATRIAL FIBRILLATION DETECTION METHOD"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #######################################
        # AMBULATORY HEART RHYTHM RECOMMENDED #
        #######################################
        column_name = '% other afib detection method - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% AMBULATORY HEART RHYTHM RECOMMENDED for IS, TIA without AFib detection"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################
        # CAROTID ARTERIES IMAGING - YES #
        ##################################
        column_name = '% carotid arteries imaging - Yes'
        if df[column_name].values[0] != 'N/A':
            tmp_df = df[[main_col, column_name]]
            tmp_df = tmp_df.sort_values([column_name], ascending = True)

            title = "% CAROTID ARTERIES IMAGING PERFORMED for IS, TIA"

            GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##############################
        # ANTITHROMBOTICS PRESCRIBED #
        ##############################
        '''
        column_name = '% patients receiving antiplatelets'
        legend = ['Antiplatelets', 'Vitamin K', 'Dabigatran', 'Rivaroxaban', 'Apixaban', 'Edoxaban', 'LMWH or heparin in prophylactic dose', 'LMWH or heparin in anticoagulant dose']

        tmp_df = df[[main_col, '% patients prescribed antithrombotics', '% patients receiving antiplatelets', '% patients receiving Vit. K antagonist', '% patients receiving dabigatran', '% patients receiving rivaroxaban', '% patients receiving apixaban', '% patients receiving edoxaban', '% patients receiving LMWH or heparin in prophylactic dose', '% patients receiving LMWH or heparin in full anticoagulant dose']]

        tmp_df = tmp_df.sort_values(['% patients prescribed antithrombotics'], ascending = True)

        title = "% ANTITHROMBOTICS PRESCRIBED for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')
        '''
        ##################################################
        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB #
        ##################################################
        #column_name = '% patients prescribed anticoagulants with aFib'
        column_name = '% afib patients discharged with anticoagulants'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS WITH AFIB, PRESCRIBED ANTICOAGULANTS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #########################################################
        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB (HOME) #
        #########################################################
        column_name = '% afib patients discharged home with anticoagulants'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS DISCHARGED HOME WITH AFIB, PRESCRIBED ANTICOAGULANTS \nfor IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################################
        # PATIENTS PRESCRIBED ANTIPLATELETS WITHOUT AFIB #
        ##################################################
        column_name = '% patients prescribed antiplatelets without aFib'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS WITHOUT AFIB, PRESCRIBED ANTIPLATELETS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        
        ###################################
        # NOT PRESCRIBED, BUT RECOMMENDED #
        ###################################
        column_name = '% patients not prescribed antithrombotics, but recommended'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS NOT PRESCRIBED, BUT RECOMMENDED ANTITHROMBOTICS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ############
        # WITH CVT #
        ############
        ##############################
        # ANTITHROMBOTICS PRESCRIBED #
        ##############################
        column_name = '% patients receiving antiplatelets with CVT'
        legend = ['Antiplatelets', 'Vitamin K', 'Dabigatran', 'Rivaroxaban', 'Apixaban', 'Edoxaban', 'LMWH or heparin in prophylactic dose', 'LMWH or heparin in anticoagulant dose']

        tmp_df = df[[main_col, '% patients prescribed antithrombotics with CVT', '% patients receiving antiplatelets with CVT', '% patients receiving Vit. K antagonist with CVT', '% patients receiving dabigatran with CVT', '% patients receiving rivaroxaban with CVT', '% patients receiving apixaban with CVT', '% patients receiving edoxaban with CVT', '% patients receiving LMWH or heparin in prophylactic dose with CVT', '% patients receiving LMWH or heparin in full anticoagulant dose with CVT']]

        tmp_df = tmp_df.sort_values(['% patients prescribed antithrombotics with CVT'], ascending = True)

        title = "% ANTITHROMBOTICS PRESCRIBED for IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##################################################
        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB #
        ##################################################
        column_name = '% patients prescribed anticoagulants with aFib with CVT'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS WITH AFIB, PRESCRIBED ANTICOAGULANTS for IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################################
        # PATIENTS PRESCRIBED ANTIPLATELETS WITHOUT AFIB #
        ##################################################
        column_name = '% patients prescribed antiplatelets without aFib with CVT'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS WITHOUT AFIB, PRESCRIBED ANTIPLATELETS for IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        
        ###################################
        # NOT PRESCRIBED, BUT RECOMMENDED #
        ###################################
        column_name = '% patients not prescribed antithrombotics, but recommended with CVT'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS NOT PRESCRIBED, BUT RECOMMENDED ANTITHROMBOTICS \nfor IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##########################
        # DISCHARGE WITH STATINS #
        ##########################
        column_name = '% patients prescribed statins - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% DISCHARGED WITH STATINS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #############################
        # CAROTID STENOSIS DETECTED #
        #############################
        column_name = '% carotid stenosis - >70%'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CAROTID STENOSIS OF OVER 70 PERCENT for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################################
        # % CAROTID STENOSIS FOLLOWUP - YES, BUT PLANNED #
        ##################################################
        sorted_by = ['% carotid stenosis followup - Yes, but planned', '% carotid stenosis followup - Referred to another centre']
        column_name = '% carotid stenosis followup - Yes, but planned'
        legend = ['Yes or planned', 'Referred to another centre']

        tmp_df = df[[main_col, '% carotid stenosis followup - Yes, but planned', '% carotid stenosis followup - Referred to another centre']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% ENDARTERECTOMY OR ANGIOPLASTY / STENTING DONE OR PLANNED for IS, TIA with ICA STENOSIS > 70%"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #############################################################
        # % ANTIHYPERTENSIVE MEDICATION PRESCRIBED out of all cases #
        #############################################################
        column_name = '% prescribed antihypertensives - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% ANTIHYPERTENSIVE MEDICATION PRESCRIBED out of all cases"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ######################################################
        # % RECOMMENDED TO A SMOKING CESSATION PROGRAM - YES #
        ######################################################
        column_name = '% recommended to a smoking cessation program - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% RECOMMENDED TO A SMOKING CESSATION PROGRAM out of smokers"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###########################################################
        # % RECOMMENDED TO A CEREBROVASCULAR EXPERT - RECOMMENDED #
        ###########################################################
        column_name = '% recommended to a cerebrovascular expert - Recommended'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% RECOMMENDED TO A CEREBROVASCULAR EXPERT out of all cases"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #########################
        # DISCHARGE DESTINATION #
        #########################
        sorted_by = ['% discharge destination - Home', '% discharge destination - Transferred within the same centre', '% discharge destination - Transferred to another centre', '% discharge destination - Social care facility', '% discharge destination - Dead']
        column_name = '% discharge destination - Home'
        legend = ['home', 'transferred within the same centre', 'transferred to another centre', 'social care facility', 'dead']

        tmp_df = df[[main_col, '% discharge destination - Home', '% discharge destination - Transferred within the same centre', '% discharge destination - Transferred to another centre', '% discharge destination - Social care facility', '% discharge destination - Dead']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% DISCHARGE DESTINATION"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ############################################################################
        # DISCHARGE DESTINATION - DEPARTMENT TRANSFERRED TO WITHIN THE SAME CENTRE #
        ############################################################################
        sorted_by = ['% transferred within the same centre - Acute rehabilitation', '% transferred within the same centre - Post-care bed', '% transferred within the same centre - Another department']
        column_name = '% transferred within the same centre - Acute rehabilitation'
        legend = ['Acute rehabilitation', 'Post-care bed', 'Another department']

        tmp_df = df[[main_col, '% transferred within the same centre - Acute rehabilitation', '% transferred within the same centre - Post-care bed', '% transferred within the same centre - Another department']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% DISCHARGE DESTINATION - PATIENT TRANSFERRED WITHIN THE SAME CENTRE"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ####################################################
        # % DISCHARGE DESTINATION - CENTRE TRANSFERRED TO  #
        ####################################################
        sorted_by = ['% transferred to another centre - Stroke centre', '% transferred to another centre - Comprehensive stroke centre', '% transferred to another centre - Another hospital']
        column_name = '% transferred to another centre - Stroke centre'
        legend = ['Stroke centre', 'Comprehensive stroke centre', 'Another hospital']

        tmp_df = df[[main_col, '% transferred to another centre - Stroke centre', '% transferred to another centre - Comprehensive stroke centre', '% transferred to another centre - Another hospital']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% DISCHARGE DESTINATION - PATIENT TRANSFERRED TO ANOTHER CENTRE"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #################################################################################
        # % DISCHARGE DESTINATION - DEPARTMENT TRANSFERRED TO WITHIN TO ANOTHER CENTRE  #
        #################################################################################
        sorted_by = ['% department transferred to within another centre - Acute rehabilitation', '% department transferred to within another centre - Post-care bed', '% department transferred to within another centre - Neurology', '% department transferred to within another centre - Another department']
        column_name = '% department transferred to within another centre - Acute rehabilitation'
        legend = ['Acute rehabilitation', 'Post-care bed', 'Neurology', 'Another department']

        tmp_df = df[[main_col, '% department transferred to within another centre - Acute rehabilitation', '% department transferred to within another centre - Post-care bed', '% department transferred to within another centre - Neurology', '% department transferred to within another centre - Another department']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% DISCHARGE DESTINATION - PATIENT TRANSFERRED TO ANOTHER CENTRE (DEPARTMENT)"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ########################
        # MEDIAN DISCHARGE mRS #
        ########################
        column_name = 'Median discharge mRS'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "MEDIAN DISCHARGE MRS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###############################
        # MEDIAN HOSPITAL STAY (DAYS) #
        ###############################
        column_name = 'Median hospital stay (days)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "MEDIAN HOSPITAL STAY (DAYS)"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        # set pptx output name (for cz it'll be presentation_CZ.pptx)
        working_dir = os.getcwd()
        if site_code is None:
            pptx = self.report + "_" + self.quarter + ".pptx"
        else:
            pptx = self.report + "_" + site_code + "_" + self.quarter + ".pptx"
        presentation_path = os.path.normpath(os.path.join(working_dir, pptx))

        prs.save(presentation_path)

class GeneratePresentationQuantiles:
    """ The class generating the presentation with quantiles.

    :param df: the dataframe with calculated statistics
    :type df: pandas dataframe
    :param country: `True` if country is included in the statistics as site
    :type country: bool
    :param country_code: the country code
    :type country_code: str
    :param split_sites: `True` if presentation should be generated per sites seperately
    :type split_sites: bool
    :param site: the site code
    :type site: str
    :param report: the type of the report eg. quarter
    :type report: str
    :param quarter: the type of the period eg. Q1_2019
    :type quarter: str
    """

    def __init__(self, df, country=False, country_code=None, split_sites=False, site=None, report=None, quarter=None, country_name=None, mt_hospitals=None):

        self.df = df.drop_duplicates(subset=['Site ID', 'Total Patients'], keep='first')
        self.country_code = country_code
        self.report = report
        self.quarter = quarter
        self.mt_hospitals = mt_hospitals

        #master_pptx = self.country_code + ".pptx"
        script_dir = os.path.dirname(__file__) 
        master_pptx = "master.pptx"
        self.master = os.path.normpath(os.path.join(script_dir, "backgrounds", master_pptx))

        # Connect to database and get country name according to country code.
        def select_country(value):
            """ The function obtaining the name of country from the package pytz based on the country code. 

            :param value: the country code
            :type value: str
            :returns: the country name
            """
            country_name = pytz.country_names[value]
            return country_name

        # If country is used as site, the country name is selected from countries dictionary by country code. :) 
        if (country):
            self.country_name = country_name
        else:
            self.country_name = country_name

        # If split_sites is True, then go through dataframe and generate graphs for each site (the country will be included as site in each file).
        site_ids = self.df['Site ID'].tolist()
        # Delete country name from side ids list.
        try:
            site_ids.remove(self.country_name)
        except:
            pass

        if site is not None:
            df = self.df[self.df['Site ID'].isin([site, self.country_name])].copy()
            self._generate_graphs(df=df, site_code=site)

        # Generate formatted statistics for all sites individualy + country as site is included
        if (split_sites) and site is None:
            for i in site_ids:
                df = self.df[self.df['Site ID'].isin([i, self.country_name])].copy()
                print(i)
                self._generate_graphs(df=df, site_code=i)
    
        # Produce formatted statistics for all sites + country as site
        if site is None:
            self._generate_graphs(df=self.df, site_code=country_code)

    def _generate_graphs(self, df, site_code=None):
        """ The function opening the presentation and generating graphs. 
        
        :param df: the dataframe with calculated statistic
        :type df: pandas dataframe
        :param site_code: the site ID
        :type site_code: str
        """
        
        prs = Presentation(self.master)

        first_slide = prs.slides[0]
        shape = first_slide.shapes[5]
        text_frame = shape.text_frame

        if self.country_name is None:
            first_slide_text = "Data Summary"
        else:
            first_slide_text = self.country_name + "\nData Summary"
        #first_slide_text = "\nData Summary"

        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = first_slide_text

        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(24)
        font.color.rgb = RGBColor(250,250,250)

        # if (self.country_name in ['Ukraine', 'Poland'] and len(df) > 2):
        #     main_col = 'Site ID'
        # else:
        main_col =  'Site Name'

        ########################
        #### TOTAL PATIENTS ####
        ########################
        column_name = 'Total Patients'
        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        #title = 'TOTAL PATIENTS'
        
        country_patients = str(max(tmp_df[column_name].tolist()))
        title = 'TOTAL PATIENTS'
        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #####################################
        #### Department type - neurology ####
        #####################################
        column_name = '% department type - neurology'
        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        #title = 'TOTAL PATIENTS'
        
        country_patients = str(max(tmp_df[column_name].tolist()))
        title = '% DEPARTMENT TYPE ALLOCATION out of all cases - neurology'
        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)


        column_name = '% patients hospitalized in stroke unit / ICU'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% HOSPITALIZATION DESTINATION out of all cases – stroke unit"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)


        column_name = '% stroke type - ischemic stroke'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% STROKE TYPE out of all cases - ischemic"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% NIHSS - Performed'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% NIHSS PERFORMED for IS, ICH, CVT"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% CT/MRI - performed'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CT/MRI PERFORMED for IS, ICH, CVT, TIA"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% CT/MRI - Performed within 1 hour after admission'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CT/MRI PERFORMED WITHIN 1 HOUR AFTER ADMISSION for IS, ICH, CVT, TIA"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% patients recanalized'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% RECANALIZATION PROCEDURES for IS"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)
        
        # IVT PA
        column_name = '% IV tPa'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% IV tPa for IS"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        # TBY
        column_name = '% TBY'

        mt_hospitals_df = df.loc[df['Site ID'].isin(self.mt_hospitals) | df['Site ID'].isin([self.country_name, 'Q1', 'Q3'])]
        tmp_df = mt_hospitals_df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)
        

        title = "% MT for IS"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = 'Median DTN (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "MEDIAN DOOR-TO-NEEDLE TIME (minutes) for thrombolyzed patients"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)
        
        
        column_name = 'Median DTG (minutes)'

        tmp_df = mt_hospitals_df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "MEDIAN DOOR-TO-GROIN TIME (minutes) for patients receiving endovascular treatment in a comprehensive centre"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = 'Median TBY DIDO (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.loc[tmp_df[column_name] != 0]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "MEDIAN DOOR-IN-DOOR-OUT TIME (minutes) for patients referred from a primary centre to another centre for recanalization therapy"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% dysphagia screening - Guss test'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "% DYSPHAGIA SCREENING PERFORMED for IS, ICH, CVT – GUSS test"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% carotid arteries imaging - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CAROTID ARTERIES IMAGING PERFORMED for IS, TIA"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% recommended to a cerebrovascular expert - Recommended'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% RECOMMENDED TO A CEREBROVASCULAR EXPERT out of all cases"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% patients assessed for rehabilitation - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% REHABILITATION ASSESSMENT for IS, ICH, CVT and SAH"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)


        column_name = '% patients prescribed statins - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% DISCHARGED WITH STATINS for IS, TIA"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% prescribed antihypertensives - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "% ANTIHYPERTENSIVE MEDICATION PRESCRIBED out of all cases"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% afib patients discharged with anticoagulants'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS WITH AFIB, PRESCRIBED ANTICOAGULANTS for IS, TIA"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)


        column_name = '% patients detected for aFib'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS DETECTED FOR AFIB for IS, TIA"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)


        column_name = '% afib patients discharged home with anticoagulants'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS DISCHARGED HOME WITH AFIB, PRESCRIBED ANTICOAGULANTS for IS, TIA"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        column_name = '% patients prescribed antiplatelets without aFib with CVT'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS WITHOUT AFIB, PRESCRIBED ANTIPLATELETS for IS, TIA, CVT"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)



        column_name = 'Median discharge mRS'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "Median discharge mRS"

        GenerateGraphsQuantiles(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)


        
        # set pptx output name (for cz it'll be presentation_CZ.pptx)
        working_dir = os.getcwd()
        if site_code is None:
            pptx = self.report + "_" + self.quarter + "_quantiles.pptx"
        else:
            pptx = self.report + "_" + site_code + "_" + self.quarter + "_quantiles.pptx"
        presentation_path = os.path.normpath(os.path.join(working_dir, pptx))

        prs.save(presentation_path)

class GeneratePresentationSites(GeneratePresentation):

    def __init__(self):

        super(GeneratePresentationSites, self).__init__(*args, **kwargs)

    def _generate_graphs(self, df, site_code=None):
        """ The function opening the presentation and generating graphs. 
        
        :param df: the dataframe with calculated statistic
        :type df: pandas dataframe
        :param site_code: the site ID
        :type site_code: str
        """
        
        prs = Presentation(self.master)

        first_slide = prs.slides[0]
        shape = first_slide.shapes[5]
        text_frame = shape.text_frame

        if self.country_name is None:
            first_slide_text = "Data Summary"
        else:
            first_slide_text = self.country_name + "\nData Summary"
        #first_slide_text = "\nData Summary"

        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = first_slide_text

        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(24)
        font.color.rgb = RGBColor(250,250,250)

        main_col =  'Site Name'

        dfs = {}
        # Total patients
        column_name = 'Total Patients'
        title = 'TOTAL PATIENTS'
        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        df = {
            'column_name': column_name,
            'title': title,
            'df': tmp_df,
            'prs': prs,
            'country': self.country_name
        }

        dfs['df1'] = df

        # Median patient age
        column_name = 'Median patient age'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "MEDIAN PATIENT AGE"

        df = {
            'column_name': column_name,
            'title': title,
            'df': tmp_df,
            'prs': prs,
            'country': self.country_name
        }    
        dfs['df2'] = df  

        GenerateGraphsSites(data=df)

        dfs = {}
        # Gender distribution
        column_name = '% patients female'
        legend = ['Female', 'Male']

        tmp_df = df[[main_col, '% patients female', '% patients male']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "GENDER DISTRIBUTION"

        

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #######################
        ### DEPARTMENT TYPE ###
        #######################
        column_name = '% department type - neurology'
        legend = ['neurology', 'neurosurgery', 'anesthesiology resuscitation critical care', 'internal medicine', 'geriatrics', 'other']

        tmp_df = df[[main_col, '% department type - neurology', '% department type - neurosurgery', '% department type - anesthesiology/resuscitation/critical care', '% department type - internal medicine', '% department type - geriatrics', '% department type - Other']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% DEPARTMENT TYPE ALLOCATION out of all cases" 

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ###################################
        ### HOSPITALIZATION DESTINATION ###
        ###################################
        column_name = '% patients hospitalized in stroke unit / ICU'
        legend = ['stroke unit', 'monitored bed with telemetry', 'standard bed']

        tmp_df = df[[main_col, '% patients hospitalized in stroke unit / ICU', '% patients hospitalized in monitored bed with telemetry', '% patients hospitalized in standard bed']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% HOSPITALIZATION DESTINATION out of all cases" 

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ###############
        # STROKE TYPE #
        ###############
        column_name = '% stroke type - ischemic stroke'
        legend = ['ischemic', 'transient ischemic attack', 'intracerebral hemorrhage', 'subarrachnoid hemorrhage', 'cerebral venous thrombosis', 'undetermined']

        tmp_df = df[[main_col, '% stroke type - ischemic stroke', '% stroke type - transient ischemic attack', '% stroke type - intracerebral hemorrhage', '% stroke type - subarrachnoid hemorrhage', '% stroke type - cerebral venous thrombosis', '% stroke type - undetermined stroke']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% STROKE TYPE out of all cases"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ######################
        # CONSIOUSNESS LEVEL #
        ######################
        column_name = 'alert_all_perc'
        legend = ['alert', 'drowsy', 'comatose']

        tmp_df = df[[main_col, 'alert_all_perc', 'drowsy_all_perc', 'comatose_all_perc']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CONSCIOUSNESS LEVEL for IS, ICH, CVT, SAH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ###################
        # NIHSS PERFORMED #
        ###################
        column_name = '% NIHSS - Performed'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% NIHSS PERFORMED for IS, ICH, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ######################
        # NIHSS MEDIAN SCORE #
        ######################
        column_name = 'NIHSS median score'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "NIHSS median score"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ####################
        # CT/MRI performed #
        ####################
        column_name = '% CT/MRI - performed'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CT/MRI PERFORMED for IS, ICH, CVT, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################
        # CT/MRI PERFORMED WITHIN 1 HOUR #
        ##################################
        column_name = '% CT/MRI - Performed within 1 hour after admission'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CT/MRI PERFORMED WITHIN 1 HOUR AFTER ADMISSION for IS, ICH, CVT, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ####################
        # VASCULAR IMAGING #
        ####################
        sorted_by = ['vascular_imaging_cta_norm', 'vascular_imaging_mra_norm', 'vascular_imaging_dsa_norm', 'vascular_imaging_none_norm']
        column_name = 'vascular_imaging_cta_norm'
        legend = ['CTA', 'MRA', 'DSA', 'none']

        tmp_df = df[[main_col, 'vascular_imaging_cta_norm', 'vascular_imaging_mra_norm', 'vascular_imaging_dsa_norm', 'vascular_imaging_none_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% VASCULAR IMAGING PERFORMED for ICH, SAH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')


        ############################
        # RECANALIZATION TREATMENT #
        ############################
        column_name = '% recanalization procedures - IV tPa'
        legend = ['IV tPa', 'IV tPa + endovascular treatment', 'IV tPa + another centre for endovascular treatment']

        tmp_df = df[[main_col, '% recanalization procedures - IV tPa', '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% IV tPa for IS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')


        ####################################################
        # RECANALIZATION TREATMENT IN COMPREHENSIVE CENTRE #
        ####################################################
        sorted_by = ['% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone']
        column_name = '% recanalization procedures - IV tPa + endovascular treatment'
        legend = ['IV tPa + endovascular treatment', 'endovascular treatment']

        tmp_df = df[[main_col, '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% RECANALIZATION PROCEDURES IN COMPREHENSIVE CENTRES for IS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################################################
        # RECANALIZATION TREATMENT IN REFFERED TO ANOTHER CENTRE #
        ##########################################################
        sorted_by = ['% recanalization procedures - IV tPa + referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment and hospitalization continues at the referred to centre', '% recanalization procedures - Referred for endovascular treatment and patient is returned to the initial centre']
        column_name = '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment'
        legend = ['IV tPa + another centre for endovascular treatment', 'another centre for endovascular treatment', 'another centre for endovascular treatment and hospitalization continues', 'another centre for endovascular treatment and returned to the initial centre']

        tmp_df = df[[main_col, '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment', '% recanalization procedures - Referred to another centre for endovascular treatment and hospitalization continues at the referred to centre', '% recanalization procedures - Referred for endovascular treatment and patient is returned to the initial centre']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% PATIENTS TRANSFERRED TO ANOTHER CENTRE FOR RECANALIZATION PROCEDURES FROM PRIMARY CENTRE for IS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ############################
        # RECANALIZATION TREATMENT #
        ############################
        column_name = '% recanalization procedures - IV tPa'
        legend = ['IV tPa', 'IV tPa + endovascular treatment', 'endovascular treatment', 'IV tPa + another centre for endovascular treatment']

        tmp_df = df[[main_col, '% patients recanalized', '% recanalization procedures - IV tPa', '% recanalization procedures - IV tPa + endovascular treatment', '% recanalization procedures - Endovascular treatment alone', '% recanalization procedures - IV tPa + referred to another centre for endovascular treatment']]
        tmp_df = tmp_df.sort_values(['% patients recanalized'], ascending = True)

        title = "% RECANALIZATION PROCEDURES for IS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################
        # % patients recanalized #
        ##########################
        column_name = '% patients recanalized'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% of recanalized patients"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)


        ################
        # % median DTN #
        ################
        column_name = 'Median DTN (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "MEDIAN DOOR-TO-NEEDLE TIME (minutes) for thrombolyzed patients"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###############
        #  median DTG #
        ###############
        column_name = 'Median DTG (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "MEDIAN DOOR-TO-GROIN TIME (minutes) for patients receiving endovascular treatment in a comprehensive centre"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ####################
        #  median TBY DIDO #
        ####################
        column_name = 'Median TBY DIDO (minutes)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = False)

        title = "MEDIAN DOOR-IN-DOOR-OUT TIME (minutes) for patients referred from a primary centre to another centre for recanalization therapy"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #######################
        # dysphagia screening #
        #######################
        column_name = '% dysphagia screening - Guss test'
        column_names = ['% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre']
        legend = ['GUSS test', 'Other test', 'Another centre']

        tmp_df = df[[main_col, '% dysphagia screening - Guss test', '% dysphagia screening - Other test', '% dysphagia screening - Another centre']]
        tmp_df = tmp_df.sort_values(column_names, ascending = True)

        title = "% DYSPHAGIA SCREENING PERFORMED for IS, ICH, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ######################################
        #  dypshagia screening within 1 hour #
        ######################################
        column_name = '% dysphagia screening time - Within first 24 hours'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% DYSPHAGIA SCREENING TIME WITHIN FIRST 24 HOURS AFTER ADMISSION"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        
        ###########################
        # PUT ON VENTILATOR - YES #
        ###########################
        column_name = '% patients put on ventilator - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS PUT ON VENTILATOR for IS, ICH, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###################
        # HEMICRANEOCTOMY #
        ###################
        column_name = '% hemicraniectomy - Yes'
        legend = ['Yes', 'Referred to another centre']

        tmp_df = df[[main_col, '% hemicraniectomy - Yes', '% hemicraniectomy - Referred to another centre']]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% HEMICRANIECTOMY PERFORMED for IS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ################
        # NEUROSURGERY #
        ################
        column_name = '% neurosurgery - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% NEUROSURGERY PERFORMED for ICH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #######################################
        # NEUROSURGERY TYPE (FIRST 3 OPTIONS) #
        #######################################
        sorted_by = ['% neurosurgery type - intracranial hematoma evacuation', '% neurosurgery type - external ventricular drainage', '% neurosurgery type - decompressive craniectomy']
        column_name = '% neurosurgery type - intracranial hematoma evacuation'
        legend = ['intracranial hematoma evacuation', 'external ventricular drainage', 'decompressive craniectomy']

        tmp_df = df[[main_col, '% neurosurgery type - intracranial hematoma evacuation', '% neurosurgery type - external ventricular drainage', '% neurosurgery type - decompressive craniectomy']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% NEUROSURGERY TYPE PERFORMED for ICH in comprehensive centres"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ###################################
        # NEUROSURGERY TYPE (LAST OPTION) #
        ###################################
        column_name = '% neurosurgery type - Referred to another centre'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS REFERRED TO ANOTHER CENTRE FOR NEUROSURGERY"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###################
        # BLEEDING REASON #
        ###################
        sorted_by = ['bleeding_arterial_hypertension_perc_norm', 'bleeding_aneurysm_perc_norm', 'bleeding_arterio_venous_malformation_perc_norm', 'bleeding_anticoagulation_therapy_perc_norm', 'bleeding_amyloid_angiopathy_perc_norm', 'bleeding_other_perc_norm']
        column_name = 'bleeding_arterial_hypertension_perc_norm'
        legend = ['arterial hypertension', 'aneurysm', 'arterio-venous malformation', 'anticoagulation therapy', 'amyloid angiopathy', 'other']

        tmp_df = df[[main_col, 'bleeding_arterial_hypertension_perc_norm', 'bleeding_aneurysm_perc_norm', 'bleeding_arterio_venous_malformation_perc_norm', 'bleeding_anticoagulation_therapy_perc_norm', 'bleeding_amyloid_angiopathy_perc_norm', 'bleeding_other_perc_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% BLEEDING REASON for ICH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################
        # INTERVENTION PERFORMED #
        ##########################
        sorted_by = ['intervention_endovascular_perc_norm', 'intervention_neurosurgical_perc_norm', 'intervention_other_perc_norm', 'intervention_referred_perc_norm', 'intervention_none_perc_norm']
        column_name = 'intervention_endovascular_perc_norm'
        legend = ['Endovascular (coiling)', 'Neurosurgical (clipping)', 'Other neurosurgical treatment', 'Patient referred to another centre', 'None']

        tmp_df = df[[main_col, 'intervention_endovascular_perc_norm', 'intervention_neurosurgical_perc_norm', 'intervention_other_perc_norm', 'intervention_referred_perc_norm', 'intervention_none_perc_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% INTERVENTION PERFORMED for SAH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##########################################
        # % PATIENTS ASSESSED FOR REHABILITATION #
        ##########################################
        column_name = '% patients assessed for rehabilitation - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% REHABILITATION ASSESSMENT for IS, ICH, CVT and SAH"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###############################
        # VENOUS THROMBOSIS TREATMENT #
        ###############################
        sorted_by = ['vt_treatment_anticoagulation_perc_norm', 'vt_treatment_thrombectomy_perc_norm', 'vt_treatment_local_thrombolysis_perc_norm', 'vt_treatment_local_neurological_treatment_perc_norm']
        column_name = 'vt_treatment_anticoagulation_perc_norm'
        legend = ['anticoagulation', 'thrombectomy', 'local thrombolysis', 'neurosurgical treatment']

        tmp_df = df[[main_col, 'vt_treatment_anticoagulation_perc_norm', 'vt_treatment_thrombectomy_perc_norm', 'vt_treatment_local_thrombolysis_perc_norm', 'vt_treatment_local_neurological_treatment_perc_norm']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% VENOUS THROMBOSIS TREATMENT for CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ################################
        # ATRIAL FIBRILLATION DETECTED #
        ################################
        sorted_by = ['% afib/flutter - Detected during hospitalization', '% afib/flutter - Newly-detected at admission', '% afib/flutter - Known']
        column_name = '% afib/flutter - Detected during hospitalization'
        legend = ['detected during hospitalization', 'newly-detected at admission', 'known aFib']

        tmp_df = df[[main_col, '% afib/flutter - Detected during hospitalization', '% afib/flutter - Newly-detected at admission', '% afib/flutter - Known']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% ATRIAL FIBRILLATION DETECTED for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ########################################
        # ATRIAL FIBRILLATION DETECTION METHOD #
        ########################################
        sorted_by = ['% afib detection method - Telemetry with monitor allowing automatic detection of aFib', '% afib detection method - Telemetry without monitor allowing automatic detection of aFib', '% afib detection method - Holter-type monitoring', '% afib detection method - EKG monitoring in an ICU bed with automatic detection of aFib', '% afib detection method - EKG monitoring in an ICU bed without automatic detection of aFib']
        column_name = '% afib detection method - Telemetry with monitor allowing automatic detection of aFib'
        legend = ['Telemetry with monitoring', 'Telemetry without monitoring', 'Holter-type monitoring', 'EKG monitoring in an ICU bed with automatic detection of aFib', 'EKG monitoring in an ICU bed without automatic detection of aFib']

        tmp_df = df[[main_col, '% afib detection method - Telemetry with monitor allowing automatic detection of aFib', '% afib detection method - Telemetry without monitor allowing automatic detection of aFib', '% afib detection method - Holter-type monitoring', '% afib detection method - EKG monitoring in an ICU bed with automatic detection of aFib', '% afib detection method - EKG monitoring in an ICU bed without automatic detection of aFib']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% ATRIAL FIBRILLATION DETECTION METHOD"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #######################################
        # AMBULATORY HEART RHYTHM RECOMMENDED #
        #######################################
        column_name = '% other afib detection method - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% AMBULATORY HEART RHYTHM RECOMMENDED for IS, TIA without AFib detection"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################
        # CAROTID ARTERIES IMAGING - YES #
        ##################################
        column_name = '% carotid arteries imaging - Yes'
        if df[column_name].values[0] != 'N/A':
            tmp_df = df[[main_col, column_name]]
            tmp_df = tmp_df.sort_values([column_name], ascending = True)

            title = "% CAROTID ARTERIES IMAGING PERFORMED for IS, TIA"

            GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##############################
        # ANTITHROMBOTICS PRESCRIBED #
        ##############################
        '''
        column_name = '% patients receiving antiplatelets'
        legend = ['Antiplatelets', 'Vitamin K', 'Dabigatran', 'Rivaroxaban', 'Apixaban', 'Edoxaban', 'LMWH or heparin in prophylactic dose', 'LMWH or heparin in anticoagulant dose']

        tmp_df = df[[main_col, '% patients prescribed antithrombotics', '% patients receiving antiplatelets', '% patients receiving Vit. K antagonist', '% patients receiving dabigatran', '% patients receiving rivaroxaban', '% patients receiving apixaban', '% patients receiving edoxaban', '% patients receiving LMWH or heparin in prophylactic dose', '% patients receiving LMWH or heparin in full anticoagulant dose']]

        tmp_df = tmp_df.sort_values(['% patients prescribed antithrombotics'], ascending = True)

        title = "% ANTITHROMBOTICS PRESCRIBED for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')
        '''
        ##################################################
        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB #
        ##################################################
        #column_name = '% patients prescribed anticoagulants with aFib'
        column_name = '% afib patients discharged with anticoagulants'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS WITH AFIB, PRESCRIBED ANTICOAGULANTS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #########################################################
        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB (HOME) #
        #########################################################
        column_name = '% afib patients discharged home with anticoagulants'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS DISCHARGED HOME WITH AFIB, PRESCRIBED ANTICOAGULANTS \nfor IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################################
        # PATIENTS PRESCRIBED ANTIPLATELETS WITHOUT AFIB #
        ##################################################
        column_name = '% patients prescribed antiplatelets without aFib'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS WITHOUT AFIB, PRESCRIBED ANTIPLATELETS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        
        ###################################
        # NOT PRESCRIBED, BUT RECOMMENDED #
        ###################################
        column_name = '% patients not prescribed antithrombotics, but recommended'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS NOT PRESCRIBED, BUT RECOMMENDED ANTITHROMBOTICS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ############
        # WITH CVT #
        ############
        ##############################
        # ANTITHROMBOTICS PRESCRIBED #
        ##############################
        column_name = '% patients receiving antiplatelets with CVT'
        legend = ['Antiplatelets', 'Vitamin K', 'Dabigatran', 'Rivaroxaban', 'Apixaban', 'Edoxaban', 'LMWH or heparin in prophylactic dose', 'LMWH or heparin in anticoagulant dose']

        tmp_df = df[[main_col, '% patients prescribed antithrombotics with CVT', '% patients receiving antiplatelets with CVT', '% patients receiving Vit. K antagonist with CVT', '% patients receiving dabigatran with CVT', '% patients receiving rivaroxaban with CVT', '% patients receiving apixaban with CVT', '% patients receiving edoxaban with CVT', '% patients receiving LMWH or heparin in prophylactic dose with CVT', '% patients receiving LMWH or heparin in full anticoagulant dose with CVT']]

        tmp_df = tmp_df.sort_values(['% patients prescribed antithrombotics with CVT'], ascending = True)

        title = "% ANTITHROMBOTICS PRESCRIBED for IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ##################################################
        # % PATIENTS PRESCRIBED ANTICOAGULANTS WITH AFIB #
        ##################################################
        column_name = '% patients prescribed anticoagulants with aFib with CVT'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS WITH AFIB, PRESCRIBED ANTICOAGULANTS for IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################################
        # PATIENTS PRESCRIBED ANTIPLATELETS WITHOUT AFIB #
        ##################################################
        column_name = '% patients prescribed antiplatelets without aFib with CVT'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS WITHOUT AFIB, PRESCRIBED ANTIPLATELETS for IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        
        ###################################
        # NOT PRESCRIBED, BUT RECOMMENDED #
        ###################################
        column_name = '% patients not prescribed antithrombotics, but recommended with CVT'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% PATIENTS NOT PRESCRIBED, BUT RECOMMENDED ANTITHROMBOTICS \nfor IS, TIA, CVT"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##########################
        # DISCHARGE WITH STATINS #
        ##########################
        column_name = '% patients prescribed statins - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% DISCHARGED WITH STATINS for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #############################
        # CAROTID STENOSIS DETECTED #
        #############################
        column_name = '% carotid stenosis - >70%'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% CAROTID STENOSIS OF OVER 70 PERCENT for IS, TIA"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ##################################################
        # % CAROTID STENOSIS FOLLOWUP - YES, BUT PLANNED #
        ##################################################
        sorted_by = ['% carotid stenosis followup - Yes, but planned', '% carotid stenosis followup - Referred to another centre']
        column_name = '% carotid stenosis followup - Yes, but planned'
        legend = ['Yes or planned', 'Referred to another centre']

        tmp_df = df[[main_col, '% carotid stenosis followup - Yes, but planned', '% carotid stenosis followup - Referred to another centre']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% ENDARTERECTOMY OR ANGIOPLASTY / STENTING DONE OR PLANNED for IS, TIA with ICA STENOSIS > 70%"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #############################################################
        # % ANTIHYPERTENSIVE MEDICATION PRESCRIBED out of all cases #
        #############################################################
        column_name = '% prescribed antihypertensives - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% ANTIHYPERTENSIVE MEDICATION PRESCRIBED out of all cases"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ######################################################
        # % RECOMMENDED TO A SMOKING CESSATION PROGRAM - YES #
        ######################################################
        column_name = '% recommended to a smoking cessation program - Yes'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% RECOMMENDED TO A SMOKING CESSATION PROGRAM out of smokers"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###########################################################
        # % RECOMMENDED TO A CEREBROVASCULAR EXPERT - RECOMMENDED #
        ###########################################################
        column_name = '% recommended to a cerebrovascular expert - Recommended'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "% RECOMMENDED TO A CEREBROVASCULAR EXPERT out of all cases"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        #########################
        # DISCHARGE DESTINATION #
        #########################
        sorted_by = ['% discharge destination - Home', '% discharge destination - Transferred within the same centre', '% discharge destination - Transferred to another centre', '% discharge destination - Social care facility', '% discharge destination - Dead']
        column_name = '% discharge destination - Home'
        legend = ['home', 'transferred within the same centre', 'transferred to another centre', 'social care facility', 'dead']

        tmp_df = df[[main_col, '% discharge destination - Home', '% discharge destination - Transferred within the same centre', '% discharge destination - Transferred to another centre', '% discharge destination - Social care facility', '% discharge destination - Dead']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% DISCHARGE DESTINATION"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ############################################################################
        # DISCHARGE DESTINATION - DEPARTMENT TRANSFERRED TO WITHIN THE SAME CENTRE #
        ############################################################################
        sorted_by = ['% transferred within the same centre - Acute rehabilitation', '% transferred within the same centre - Post-care bed', '% transferred within the same centre - Another department']
        column_name = '% transferred within the same centre - Acute rehabilitation'
        legend = ['Acute rehabilitation', 'Post-care bed', 'Another department']

        tmp_df = df[[main_col, '% transferred within the same centre - Acute rehabilitation', '% transferred within the same centre - Post-care bed', '% transferred within the same centre - Another department']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% DISCHARGE DESTINATION - PATIENT TRANSFERRED WITHIN THE SAME CENTRE"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ####################################################
        # % DISCHARGE DESTINATION - CENTRE TRANSFERRED TO  #
        ####################################################
        sorted_by = ['% transferred to another centre - Stroke centre', '% transferred to another centre - Comprehensive stroke centre', '% transferred to another centre - Another hospital']
        column_name = '% transferred to another centre - Stroke centre'
        legend = ['Stroke centre', 'Comprehensive stroke centre', 'Another hospital']

        tmp_df = df[[main_col, '% transferred to another centre - Stroke centre', '% transferred to another centre - Comprehensive stroke centre', '% transferred to another centre - Another hospital']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% DISCHARGE DESTINATION - PATIENT TRANSFERRED TO ANOTHER CENTRE"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        #################################################################################
        # % DISCHARGE DESTINATION - DEPARTMENT TRANSFERRED TO WITHIN TO ANOTHER CENTRE  #
        #################################################################################
        sorted_by = ['% department transferred to within another centre - Acute rehabilitation', '% department transferred to within another centre - Post-care bed', '% department transferred to within another centre - Neurology', '% department transferred to within another centre - Another department']
        column_name = '% department transferred to within another centre - Acute rehabilitation'
        legend = ['Acute rehabilitation', 'Post-care bed', 'Neurology', 'Another department']

        tmp_df = df[[main_col, '% department transferred to within another centre - Acute rehabilitation', '% department transferred to within another centre - Post-care bed', '% department transferred to within another centre - Neurology', '% department transferred to within another centre - Another department']]
        tmp_df = tmp_df.sort_values(sorted_by, ascending = True)

        title = "% DISCHARGE DESTINATION - PATIENT TRANSFERRED TO ANOTHER CENTRE (DEPARTMENT)"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name, legend=legend, number_of_series=len(legend), graph_type='stacked')

        ########################
        # MEDIAN DISCHARGE mRS #
        ########################
        column_name = 'Median discharge mRS'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "MEDIAN DISCHARGE MRS"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        ###############################
        # MEDIAN HOSPITAL STAY (DAYS) #
        ###############################
        column_name = 'Median hospital stay (days)'

        tmp_df = df[[main_col, column_name]]
        tmp_df = tmp_df.sort_values([column_name], ascending = True)

        title = "MEDIAN HOSPITAL STAY (DAYS)"

        GenerateGraphs(dataframe=tmp_df, presentation=prs, title=title, column_name=column_name, country=self.country_name)

        # set pptx output name (for cz it'll be presentation_CZ.pptx)
        working_dir = os.getcwd()
        if site_code is None:
            pptx = self.report + "_" + self.quarter + ".pptx"
        else:
            pptx = self.report + "_" + site_code + "_" + self.quarter + ".pptx"
        presentation_path = os.path.normpath(os.path.join(working_dir, pptx))

        prs.save(presentation_path)

    









        


