from resqdb.Connection import Connection
from resqdb.functions import save_file

from datetime import datetime
import logging
import os
import sys
import json
import pandas as pd

from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.dml import MSO_LINE
from pptx.oxml.xmlchemy import OxmlElement

import xlsxwriter
from xlsxwriter.utility import xl_rowcol_to_cell, xl_col_to_name

class AfricaReport():
    ''' Generate reports for South Africa. 
    
    :param df: the raw data or preprocessed data
    :type df: DataFrame
    :param start_date: starting date included in filtered data
    :type start_date: datetime
    :param end_date: last date included in the filtered data
    :type end_date: datetime
    :param period_name: the name of the period
    :type period_name: str
    :param report_type: the name of the report
    :type report_type: str
    :param split: True if site reports should be generated as well
    :type split: boolean
    '''

    def __init__(self, df=None, start_date=None, end_date=None, period_name=None, report_type=None, site_reports=False, region_reports=False):
        
        # Set logging
        debug =  f'debug_{datetime.now().strftime("%d-%m-%Y")}.log' 
        log_file = os.path.join(os.getcwd(), debug)
        logging.basicConfig(
            filename=log_file, 
            filemode='a', 
            format='%(asctime)s,%(msecs)d %(name)s %(levelname)s %(message)s', 
            datefmt='%H:%M:%S', 
            level=logging.DEBUG
            )
        logging.info('Start to generate reports for South Africa.')

        # Get country code
        self.country_code = 'ZA'
        self.country_name = 'South Africa'
        self.period_name = period_name
        self.report_type = report_type
        self.site_reports = site_reports
        self.region_reports = region_reports
        self.region_name = None

        # Connect to database and get Africa data
        if df is None:
            con = Connection(data='africa')
            self.raw_data = con.preprocessed_data.copy()
            self.raw_data = self._preprocess_data(df=self.raw_data)
            logging.info('The preprocessed data were generated.')
        else:
            # Convert hospital date into datetime if data were read from csv
            date = df['HOSPITAL_DATE'].iloc[0]
            if '/' in date:
                dateForm = '%d/%m/%Y'
            else:
                dateForm = '%Y-%m-%d'

            self.raw_data = df.copy()
            columns = [x for x in self.raw_data.columns.tolist() if 'DATE' in x]
            for column in columns:
                self.raw_data[column] = pd.to_datetime(self.raw_data[column], format=dateForm, errors='ignore')


        # If start date and end date are defined, filter data by hospital date otherwise keep all data
        if start_date is None and end_date is None:
            self.preprocessed_data = self.raw_data
        else:
            self.preprocessed_data = self._filter_by_date(self.raw_data, start_date, end_date)
            logging.info('The data has been filter by date.')

        self._columns_to_be_deleted = []

        # Read regions mapping from the json file 
        path = os.path.join(os.path.dirname(__file__), 'tmp', 'south_africa_mapping.json')
        with open(path, 'r', encoding='utf-8') as json_file:
            self.regions = json.load(json_file)

        # Create REGION column in the dataframe based on the region in the SITE ID
        self.raw_data['REGION'] = self.raw_data.apply(
            lambda x: self._get_region(x['SITE_ID']), axis=1
        )

        # Add all data into dataframe again, this data will be set as country results, therefore we have to modify beofre appending SITE_ID, FACILITY_NAME and REGION
        country_df = self.preprocessed_data.copy()
        country_df['SITE_ID'] = self.country_name
        country_df['FACILITY_NAME'] = self.country_name
        country_df['REGION'] = self.country_name
        self.preprocessed_data = self.preprocessed_data.append(country_df, ignore_index=True)

        ###########################
        # Generate country report #
        # Calculate statistic
        self.calculate_statistics(self.preprocessed_data)
        # generate formatted statistic
        if self.report_type == 'all' and self.period_name == 'all':
            filename = self.country_code
        else:
            filename = f'{self.report_type}_{self.country_code}_{self.period_name}'
        self._generate_formatted_preprocessed_data(self.preprocessed_data, filename)
        self._generate_formatted_stats(self.stats, filename)
        # Generate presetation
        self._generate_presentation(self.stats, filename)
        logging.info('The country report has been generated.')

        if region_reports:
            region_preprocessed_data = self.preprocessed_data.copy()
            region_preprocessed_data['SITE_ID'] = region_preprocessed_data['REGION']
            region_preprocessed_data['FACILITY_NAME'] = region_preprocessed_data['REGION']

            self.calculate_statistics(region_preprocessed_data)
            if self.report_type == 'all' and self.period_name == 'all':
                filename = f'{self.country_code}_regions'
            else:
                filename = f'{self.report_type}_{self.country_code}_{self.period_name}_regions'
            self._generate_formatted_preprocessed_data(region_preprocessed_data, filename)
            self._generate_formatted_stats(self.stats, filename)
            # Generate presetation
            self._generate_presentation(self.stats, filename)
            logging.info('The country vs regions report has been generated.')

        if site_reports:
            site_ids = [x for x in set(self.preprocessed_data['SITE_ID'].tolist()) if x != self.country_name]
            for site_id in site_ids:
                self.region_name = self._get_region(site_id)
                # Filter data for site and country
                site_preprocessed_data = self.preprocessed_data.loc[
                    (self.preprocessed_data['SITE_ID'] == site_id) |
                    (self.preprocessed_data['SITE_ID'] == self.country_name)
                ].copy()
                site_name = site_preprocessed_data.loc[site_preprocessed_data['SITE_ID'] == site_id]['FACILITY_NAME'].iloc[0]
                # Append data for region to the site preprocessed data
                region_preprocessed_data = self.preprocessed_data.loc[
                    self.preprocessed_data['REGION'] == self.region_name
                ].copy()
                region_preprocessed_data['SITE_ID'] = self.region_name
                region_preprocessed_data['FACILITY_NAME'] = self.region_name
                site_preprocessed_data = site_preprocessed_data.append(
                    region_preprocessed_data, ignore_index=True
                )

                self.calculate_statistics(site_preprocessed_data)
                if self.report_type == 'all' and self.period_name == 'all':
                    filename = site_id
                else:
                    filename = f'{self.report_type}_{site_id}_{self.period_name}'
                self._generate_formatted_preprocessed_data(site_preprocessed_data, filename, exclude_country=True)
                self._generate_formatted_stats(self.stats, filename)
                # Generate presetation
                self._generate_presentation(self.stats, filename, site_name)
                logging.info(f'The site report for {site_id} has been generated.')

    @property
    def stats(self):
        return self._stats

    @stats.setter
    def stats(self, df):
        self._stats = df

    @property
    def columns_to_be_deleted(self):
        return self._columns_to_be_deleted

    @columns_to_be_deleted.setter
    def columns_to_be_deleted(self, value):
        self._column_to_be_deleted = value

    @property
    def country_code(self):
        return self._country_code

    @country_code.setter
    def country_code(self, value):
        self._country_code = value

    @property
    def country_name(self):
        return self._country_name

    @country_name.setter
    def country_name(self, value):
        self._country_name = value

    @property
    def region_name(self):
        return self._region_name

    @region_name.setter
    def region_name(self, value):
        self._region_name = value

    def _get_region(self, site_id):
        ''' Get region name based on code in Site ID. 

        :param site_id: the site ID
        :type site_id: str
        :returns: str
        '''
        region_code = site_id.split('_')[1]
        if region_code in self.regions:
            return self.regions[region_code] 
        else:
            return 'Demo'


    def _filter_by_date(self, df, start_date, end_date):
        ''' Filter data by DISCHARGE DATE where discharge date is between start and end date. 

        :param start: first date to be included
        :type start: datetime
        :param end: last date to be included
        :type end: datetime
        '''
        if isinstance(start_date, datetime):
            start_date = start_date.date()            
        if isinstance(end_date, datetime):
            end_date = end_date.date()

        start_date = pd.Timestamp(start_date)
        end_date = pd.Timestamp(end_date)

        df = df.loc[df['HOSPITAL_DATE'].between(start_date, end_date, inclusive=True)]
        if df.empty:
            sys.exit("The dataframe filtered for period is empty! Please, select a different period!")
        return df

    def _merge_stats(self, df):
        ''' Merge calculated column with the stats dataframe. 
        
        :param df: the dataframe that should be merged with the stats
        :type df: DataFrame
        '''
        self.stats = self.stats.merge(df, how="outer")
        self.stats.fillna(0, inplace=True)

    def _get_numbers(self, df, column_name, denominator):
        ''' Get numbers of patients for dataframe grouped by SITE_ID column. 
        The % column is calculated as well. 

        :param df: the filtered dataframe
        :type df: DataFram
        :param column_name: the name of the column to be created
        :type column_name: str
        :param denominator: the name of column that was already calculated and is in the stats
        :type denominator: str
        '''
        import re

        grouped_df = df.groupby(['SITE_ID']).size().reset_index(name=column_name)
        self._merge_stats(grouped_df)

        column_name_perc = re.sub('^#', '%', column_name)
        self.stats[column_name_perc] = self.stats.apply(lambda x: round(
            ((
                x[column_name]/x[denominator]
            ) * 100), 2
        ) if x[denominator] > 0 else 0, axis=1)
        

    def _get_median(self, df, column, new_column):
        ''' Calculate median for a column grouped by SITE_ID. 
        
        :param df: the dataframe
        :type df: DataFrame
        :param column: the name of column for which the median should be calculated
        :type column: str
        :param new_column: the name of column that will be created and merged with stats
        :type new_column: str
        '''
        median_df = df[['SITE_ID', column]].groupby(['SITE_ID']).median().reset_index().rename(columns={column: new_column})
        self._merge_stats(median_df)

    def _get_iqr(self, df, column, new_column):
        ''' Get IQR for a column grouped by SITE_ID. 

        :param df: the dataframe
        :type df: DataFrame
        :param column: the name of column for which the median should be calculated
        :type column: str
        :param new_column: the name of column that will be created and merged with stats
        :type new_column: str
        '''
        from scipy.stats import iqr
        iqr_df = df.groupby(['SITE_ID']).agg({column: iqr}).reset_index().rename(columns={column: new_column})
        self._merge_stats(iqr_df)

    def _get_total_patients(self, df, column_name, to_be_deleted=False):
        ''' Get number of patients in dataframe grouped by SITE_ID. 

        :param df: the dataframe
        :type df: DataFrame
        :param column_name: the name of column to be calculated
        :type column_name: str
        :param to_be_deleted: add column to the list of column names that will be at the end deleted from the results
        :type to_be_deleted: bool
        '''
        if to_be_deleted:
            self.columns_to_be_deleted.append(column_name)

        tmp_df = df.groupby(['SITE_ID']).size().reset_index(name=column_name)
        self._merge_stats(tmp_df)

    def calculate_statistics(self, df=None):
        ''' Calculate statistics for the South Africa. 
        
        :param df: the preprocessed data that can be filtered (default: None)
        :type df: DataFrame
        '''
        # Check if argument is provided
        if df is None:
            df = self.preprocessed_data.copy()

        self.stats = df.groupby(['SITE_ID', 'FACILITY_NAME']).size().reset_index(name='Total Patients')
        # Get patients with stroke
        tmp_df = df.loc[df['STROKE_TYPE'] != 6].copy() 
        self._get_numbers(
            tmp_df, 
            'Total number of stroke patients treated at the hospital for period of assessment', 
            'Total Patients'
            )
        
        # Get median of age per all patients
        self._get_median(df, 'AGE', 'Median patient age')
       
        # Get # and % of males
        tmp_df = df.loc[df['GENDER'] == 1].copy()
        self._get_numbers(tmp_df, '# Male', 'Total Patients')
        
        # Get # and % of females
        tmp_df = df.loc[df['GENDER'] == 2].copy()
        self._get_numbers(tmp_df, '# Female', 'Total Patients')

        # Get number of IS, TIA, ICH, CVT
        is_tia_ich_cvt = df.loc[df['STROKE_TYPE'].isin([1,2,3,5])].copy()
        self._get_total_patients(is_tia_ich_cvt, 'is_tia_ich_cvt_pts', True)

        # Get # CT or MRI done for stroke patients
        tmp_df = is_tia_ich_cvt.loc[is_tia_ich_cvt['CT_MRI_OVERALL'] == 1].copy()
        self._get_numbers(tmp_df, '# CT or MRI done for stroke patients - Yes', 'is_tia_ich_cvt_pts')

        # Get # CT or MRI not done for stroke patients
        tmp_df = is_tia_ich_cvt.loc[is_tia_ich_cvt['CT_MRI_OVERALL'] == 2].copy()
        self._get_numbers(tmp_df, '# CT or MRI done for stroke patients - No', 'is_tia_ich_cvt_pts')

        # Overall outcome for all patients at time of discharge for period of assessment (as per NIHSS score) - can be calculated only if admission nihss and discharge nihss was performed
        nihss_done = df.loc[(df['NIHSS'] == 2) & (df['D_NIHSS'] == 2)].copy()
        self._get_total_patients(nihss_done, 'nihss_done_pts', True)

        # Get overall outcome improved
        tmp_df = nihss_done.loc[(df['OVERALL_OUTCOME'] == 1)].copy()
        self._get_numbers(tmp_df, '# of patients that improved', 'nihss_done_pts')

        # Get overall outcome whose clinical condition did not cange
        tmp_df = nihss_done.loc[(df['OVERALL_OUTCOME'] == 2)].copy()
        self._get_numbers(tmp_df, '# of patients whose clinical condition did not change', 'nihss_done_pts')

        # Get overall outcome of patients who deteriorated in hospital
        tmp_df = nihss_done.loc[(df['OVERALL_OUTCOME'] == 3)].copy()
        self._get_numbers(tmp_df, '# of patients who deteriorated in hospital', 'nihss_done_pts')

        # Stroke complications (all patients)
        # Get patients with stroke and not referred for recanalization procedures
        stroke_not_referred = df.loc[
            (df['STROKE_TYPE'] != 6) & 
            (df['REFERRED_FOR_RECAN'] == 2) &
            (df['HEMICRANIECTOMY'] != 3)].copy()
        self._get_total_patients(stroke_not_referred, 'stroke_not_referred_pts', True)

        def check_if_value_selected(selected_values, value):
            values = selected_values.split(',')
            if value in values:
                return 1
            else:
                return 2
        
        # # of patients with no complications
        stroke_not_referred['COMPLICATIONS'] = stroke_not_referred.apply(
            lambda x: check_if_value_selected(x['POST_STROKE_COMPLICATIONS'], '1'), 
            axis=1
        )
        tmp_df = stroke_not_referred.loc[stroke_not_referred['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with no complications', 'stroke_not_referred_pts')

        # # of patients with pneumonia
        stroke_not_referred['COMPLICATIONS'] = stroke_not_referred.apply(
            lambda x: check_if_value_selected(x['POST_STROKE_COMPLICATIONS'], '3'), 
            axis=1
        )
        tmp_df = stroke_not_referred.loc[stroke_not_referred['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with pneumonia', 'stroke_not_referred_pts')

        # # of patients with DVT
        stroke_not_referred['COMPLICATIONS'] = stroke_not_referred.apply(
            lambda x: check_if_value_selected(x['POST_STROKE_COMPLICATIONS'], '4'), 
            axis=1
        )
        tmp_df = stroke_not_referred.loc[stroke_not_referred['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with DVT', 'stroke_not_referred_pts')

        # # of patients with pulmonary embolus
        stroke_not_referred['COMPLICATIONS'] = stroke_not_referred.apply(
            lambda x: check_if_value_selected(x['POST_STROKE_COMPLICATIONS'], '5'), 
            axis=1
        )
        tmp_df = stroke_not_referred.loc[stroke_not_referred['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with pulmonary embolus', 'stroke_not_referred_pts')

        # # of patients with worsening of stroke
        stroke_not_referred['COMPLICATIONS'] = stroke_not_referred.apply(
            lambda x: check_if_value_selected(x['POST_STROKE_COMPLICATIONS'], '6'), 
            axis=1
        )
        tmp_df = stroke_not_referred.loc[stroke_not_referred['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with worsening of stroke', 'stroke_not_referred_pts')

        # # of patients with drip sespis
        stroke_not_referred['COMPLICATIONS'] = stroke_not_referred.apply(
            lambda x: check_if_value_selected(x['POST_STROKE_COMPLICATIONS'], '10'), 
            axis=1
        )
        tmp_df = stroke_not_referred.loc[stroke_not_referred['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with drip sepsis', 'stroke_not_referred_pts')

        # # of patients with UTI
        stroke_not_referred['COMPLICATIONS'] = stroke_not_referred.apply(
            lambda x: check_if_value_selected(x['POST_STROKE_COMPLICATIONS'], '7'), 
            axis=1
        )
        tmp_df = stroke_not_referred.loc[stroke_not_referred['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with UTI', 'stroke_not_referred_pts')

        # # of patients with pressure sore
        stroke_not_referred['COMPLICATIONS'] = stroke_not_referred.apply(
            lambda x: check_if_value_selected(x['POST_STROKE_COMPLICATIONS'], '9'), 
            axis=1
        )
        tmp_df = stroke_not_referred.loc[stroke_not_referred['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with pressure sore', 'stroke_not_referred_pts')

        # # of patients with death
        stroke_not_referred['COMPLICATIONS'] = stroke_not_referred.apply(
            lambda x: check_if_value_selected(x['POST_STROKE_COMPLICATIONS'], '8'), 
            axis=1
        )
        tmp_df = stroke_not_referred.loc[stroke_not_referred['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with death (overall mortality)', 'stroke_not_referred_pts')

        # # of patients with other stroke complications
        stroke_not_referred['COMPLICATIONS'] = stroke_not_referred.apply(
            lambda x: check_if_value_selected(x['POST_STROKE_COMPLICATIONS'], '11'), 
            axis=1
        )
        tmp_df = stroke_not_referred.loc[stroke_not_referred['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with other post-stroke complications', 'stroke_not_referred_pts')

        # NIHSS score on arrival 
        nihss = df.loc[df['NIHSS'] == 2].copy() # get NIHSS performed
        self._get_median(nihss, 'NIHSS_SCORE', 'NIHSS score on arrival (median)')
        self._get_iqr(nihss, 'NIHSS_SCORE', 'NIHSS score on arrival (IQR)')

        # NIHSS score on discharge
        discharge_nihss = df.loc[df['D_NIHSS'] == 2].copy() # get NIHSS performed
        self._get_median(discharge_nihss, 'D_NIHSS_SCORE', 'NIHSS score on discharge (median)')
        self._get_iqr(discharge_nihss, 'D_NIHSS_SCORE', 'NIHSS score on discharge (IQR)')

        # Calculate Discharge MRS score median and iQR
        stroke_not_referred_or_not_dead = df.loc[
            (df['HEMICRANIECTOMY'] != 3) & (df['DISCHARGE_DESTINATION'] != 5)].copy()
        self._get_median(
            stroke_not_referred_or_not_dead, 
            'DISCHARGE_MRS_SCORE', 
            'modified Ranking Score on discharge (median)')
        self._get_iqr(
            stroke_not_referred_or_not_dead, 
            'DISCHARGE_MRS_SCORE', 
            'modified Ranking Score on discharge (IQR)')

        # Calculate overall stroke pateint mortality for period of assessment
        dead = df.loc[(df['STROKE_TYPE'] != 6) & (df['DISCHARGE_DESTINATION'] == 5)].copy()
        self._get_numbers(tmp_df, '# overall stroke patient mortality for period of assessment', 'Total number of stroke patients treated at the hospital for period of assessment')

        # Reperfusion therapy
        intravenous_lysis = df.loc[(df['STROKE_TYPE'] == 1) & (df['RECANALIZATION_PROCEDURES'].isin([2,3,4,5,6]))].copy()    
        self._get_total_patients(intravenous_lysis, 'Total number stroke patients treated with intravenous lysis for period of assessment', False)

        # Total number of patients that arrived within 3 hours after onset 
        arrival_within_180min = df.loc[df['TIME_ONSET_TO_ED_ARRIVAL'] <= 180].copy()
        self._get_total_patients(arrival_within_180min, 'Total number of patients that arrived at hospital within 3 hours of symptom onset for period of assessment', False)

        # # of all patients that received IVT or TBY - Yes
        ischemic = df.loc[df['STROKE_TYPE'] == 1].copy()
        self._get_total_patients(ischemic, 'ischemic', True)

        ivt_or_tby = df.loc[(df['IVT_DONE'] == 1) | (df['TBY_DONE'] == 1)].copy() 
        self._get_numbers(ivt_or_tby, '# of all patients that received Intravenous thrombolysis or mechanical thrombectomy - Yes', 'ischemic')

        ischemic_and_not_admitted_for_recan = df.loc[(df['STROKE_TYPE'] == 1) & (df['RECANALIZATION_PROCEDURES'] != 7)].copy()
        self._get_total_patients(ischemic_and_not_admitted_for_recan, 'ischemic_and_not_admitted_for_recan_pts', True)

        recan_not_done = df.loc[(df['STROKE_TYPE'] == 1) & (df['RECANALIZATION_PROCEDURES'] == 1)].copy()
        self._get_numbers(recan_not_done, '# of all patients that received Intravenous thrombolysis or mechanical thrombectomy - No', 'ischemic_and_not_admitted_for_recan_pts')

        # Symptom onset to needle time (median, IQR)
        ivt_done = df.loc[df['IVT_DONE'] == 1].copy()
        self._get_total_patients(ivt_done, 'ivt_done_pts', True)

        self._get_median(ivt_done, 'ONSET_TO_NEEDLE_TIME', 'Symptom onset to needle time (median)')
        self._get_iqr(ivt_done, 'ONSET_TO_NEEDLE_TIME', 'Symptom onset to needle time (IQR)')

        # Door to needle time (Median, IQR)
        self._get_median(ivt_done, 'DOOR_TO_NEEDLE_TIME', 'Door to needle time (median)')
        self._get_iqr(ivt_done, 'DOOR_TO_NEEDLE_TIME', 'Door to needle time (IQR)')

        # Outcome after intravenous lysis (as per NIHSS score)
        ivt_nihss_done = ivt_done.loc[(ivt_done['NIHSS'] == 2) & (ivt_done['D_NIHSS'] == 2)].copy()
        self._get_total_patients(ivt_nihss_done, 'ivt_nihss_done_pts', True)

        # Get overall outcome improved
        tmp_df = ivt_nihss_done.loc[(df['OVERALL_OUTCOME'] == 1)].copy()
        self._get_numbers(tmp_df, '# of patients that improved after intravenous lysis', 'ivt_nihss_done_pts')

        # Get overall outcome whose clinical condition did not cange
        tmp_df = ivt_nihss_done.loc[(df['OVERALL_OUTCOME'] == 2)].copy()
        self._get_numbers(tmp_df, '# of patients whose clinical condition did not change after intravenous lysis', 'ivt_nihss_done_pts')

        # Get overall outcome of patients who deteriorated in hospital
        tmp_df = ivt_nihss_done.loc[(df['OVERALL_OUTCOME'] == 3)].copy()
        self._get_numbers(tmp_df, '# of patients who deteriorated in hospital after intravenous lysis', 'ivt_nihss_done_pts')

        # Get number of patients that died after lysis
        ivt_done_dead = ivt_done.loc[ivt_done['DISCHARGE_DESTINATION'] == 5].copy()
        self._get_numbers(ivt_done_dead, '# of patients that died after lysis', 'ivt_done_pts')

        # Total number of stroke patients treated with mechanical thrombectomy for period of assessment
        tby_done = df.loc[df['TBY_DONE'] == 1].copy()
        self._get_total_patients(tby_done, 'tby_done_pts', True)

        self._get_numbers(tby_done, '# of all patients that received Mechanical thrombectomy - Yes', 'ischemic')

        mt_not_received = df.loc[(df['STROKE_TYPE'] == 1) & (df['TBY_DONE'] != 1)].copy()
        self._get_numbers(tby_done, '# of all patients that received Mechanical thrombectomy - No', 'ischemic')

        # Symptoms onset to arterial puncture time (median, IQR)
        self._get_median(tby_done, 'ONSET_TO_GROIN_PUNCTURE_TIME', 'Symptom onset to arterial puncture time (median)')
        self._get_iqr(tby_done, 'ONSET_TO_GROIN_PUNCTURE_TIME', 'Symptom onset to arterial puncture time (IQR)')

        # Symptoms door to arterial puncture time (median, IQR)
        self._get_median(tby_done, 'DOOR_TO_GROIN_PUNCTURE_TIME', 'Door to arterial puncture time (median)')
        self._get_iqr(tby_done, 'DOOR_TO_GROIN_PUNCTURE_TIME', 'Door to arterial puncture time (IQR)')

         # Outcome after mechanical thrombectomy (as per NIHSS score)
        tby_nihss_done = tby_done.loc[(tby_done['NIHSS'] == 2) & (tby_done['D_NIHSS'] == 2)].copy()
        self._get_total_patients(tby_nihss_done, 'tby_nihss_done_pts', True)

        # Get overall outcome improved
        tmp_df = tby_nihss_done.loc[(df['OVERALL_OUTCOME'] == 1)].copy()
        self._get_numbers(tmp_df, '# of patients that improved after mechanical thrombectomy', 'tby_nihss_done_pts')

        # Get overall outcome whose clinical condition did not cange
        tmp_df = tby_nihss_done.loc[(df['OVERALL_OUTCOME'] == 2)].copy()
        self._get_numbers(tmp_df, '# of patients whose clinical condition did not change after mechanical thrombectomy', 'tby_nihss_done_pts')

        # Get overall outcome of patients who deteriorated in hospital
        tmp_df = tby_nihss_done.loc[(df['OVERALL_OUTCOME'] == 3)].copy()
        self._get_numbers(tmp_df, '# of patients who deteriorated in hospital after mechanical thrombectomy', 'tby_nihss_done_pts')

         # Get number of patients that died after mechanical thrombectomy
        tby_done_dead = tby_done.loc[tby_done['DISCHARGE_DESTINATION'] == 5].copy()
        self._get_numbers(tby_done_dead, '# of patients that died after mechanical thrombectomy', 'tby_done_pts')

        # Complications of reperfusion therapy (lysis and/or thrombectomy - expressed as a percentage of the number that received reperfusion therapy)
        ischemic_recan = df.loc[(df['STROKE_TYPE'] == 1) & (df['RECANALIZATION_PROCEDURES'] != 1)].copy()
        self._get_total_patients(ischemic_recan, 'ischemic_recan_pts', True)

        # # #  of patients with non-fatal symptomatic intracranial haemorrhage (SICH)
        ischemic_recan['COMPLICATIONS'] = ischemic_recan.apply(
            lambda x: check_if_value_selected(x['REPERFUSION_COMPLICATIONS'], '3'), 
            axis=1
        )
        tmp_df = ischemic_recan.loc[ischemic_recan['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with non-fatal symptomatic intracranial haemorrhage (SICH)', 'ischemic_recan_pts')

        # # # of patients with fatal SICH
        ischemic_recan['COMPLICATIONS'] = ischemic_recan.apply(
            lambda x: check_if_value_selected(x['REPERFUSION_COMPLICATIONS'], '4'), 
            axis=1
        )
        tmp_df = ischemic_recan.loc[ischemic_recan['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with fatal SICH', 'ischemic_recan_pts')

        # # of patients with remote cerebral haemorrhage
        ischemic_recan['COMPLICATIONS'] = ischemic_recan.apply(
            lambda x: check_if_value_selected(x['REPERFUSION_COMPLICATIONS'], '5'), 
            axis=1
        )
        tmp_df = ischemic_recan.loc[ischemic_recan['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with remote cerebral haemorrhage', 'ischemic_recan_pts')

        # # of patients with vessel perforation
        ischemic_recan['COMPLICATIONS'] = ischemic_recan.apply(
            lambda x: check_if_value_selected(x['REPERFUSION_COMPLICATIONS'], '6'), 
            axis=1
        )
        tmp_df = ischemic_recan.loc[ischemic_recan['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with vessel perforation', 'ischemic_recan_pts')

        # # of patients with procedure related dissection
        ischemic_recan['COMPLICATIONS'] = ischemic_recan.apply(
            lambda x: check_if_value_selected(x['REPERFUSION_COMPLICATIONS'], '7'), 
            axis=1
        )
        tmp_df = ischemic_recan.loc[ischemic_recan['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with procedure related dissection', 'ischemic_recan_pts')

        # # of patients with embolization to different vascular territory
        ischemic_recan['COMPLICATIONS'] = ischemic_recan.apply(
            lambda x: check_if_value_selected(x['REPERFUSION_COMPLICATIONS'], '8'), 
            axis=1
        )
        tmp_df = ischemic_recan.loc[ischemic_recan['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with embolization to different vascular territory', 'ischemic_recan_pts')

        # # of patients with groin haematoma requiring transfusion
        ischemic_recan['COMPLICATIONS'] = ischemic_recan.apply(
            lambda x: check_if_value_selected(x['REPERFUSION_COMPLICATIONS'], '9'), 
            axis=1
        )
        tmp_df = ischemic_recan.loc[ischemic_recan['COMPLICATIONS'] == 1].copy()
        self._get_numbers(tmp_df, '# of patients with groin haematoma requiring transfusion', 'ischemic_recan_pts')

        # # of patients who died (mortality),
        ischemic_recan_dead = ischemic_recan.loc[ischemic_recan['DISCHARGE_DESTINATION'] == 5].copy()
        self._get_numbers(ischemic_recan_dead, '# of patients who died (mortality)', 'ischemic_recan_pts')

        # of patients with failure to achieve reperfusion (<TICI 2b or 3)
        reperfusion_failure = tby_done.loc[tby_done['TICI_SCORE'].isin([1,2,3,7])].copy()
        self._get_numbers(reperfusion_failure, '# of patients with failure to achieve reperfusion (<TICI 2b or 3)', 'tby_done_pts')
        
        ### Post-Acute Treatment received for period of assessment (all patients)
        # % of patients who had a Dysphagia screen done
        is_ich_cvt_not_referred = df.loc[(df['STROKE_TYPE'].isin([1,2,5])) & (~df['RECANALIZATION_PROCEDURES'].isin([3,5]))].copy()
        self._get_total_patients(
            is_ich_cvt_not_referred, 'is_ich_cvt_not_referred_pts', True
        )

        dysphagia_done = is_ich_cvt_not_referred.loc[is_ich_cvt_not_referred['DYSPHAGIA_SCREENING'].isin([1,2])].copy()
        self._get_numbers(dysphagia_done, '# of patients who had a Dysphagia screen done', 'is_ich_cvt_not_referred_pts')

        
        is_ich_sah_cvt_not_referred = df.loc[(df['STROKE_TYPE'].isin([1,2,4,5])) & (df['REFERRED_FOR_RECAN'] == 2) & (df['HEMICRANIECTOMY'] != 3)].copy()
        self._get_total_patients(
            is_ich_sah_cvt_not_referred, 'is_ich_sah_cvt_not_referred_pts', True
        )

        # % of patients who had physiotherapy assessment
        physiotherapy_done = is_ich_sah_cvt_not_referred.loc[
            is_ich_sah_cvt_not_referred['PHYSIOTHERAPIST_EVALUATION'].isin([1,2,3])
        ].copy()
        self._get_numbers(physiotherapy_done, '# of patients who had physiotherapy assessment', 'is_ich_sah_cvt_not_referred_pts')
        del physiotherapy_done

        # % of patients who had Occupational therapy assessment
        occup_therapy_done = is_ich_sah_cvt_not_referred.loc[
            is_ich_sah_cvt_not_referred['OCCUP_PHYSIOTHERAPIST_EVALUATION'].isin([1,2,3])
        ].copy()
        self._get_numbers(occup_therapy_done, '# of patients who had Occupational therapy assessment', 'is_ich_sah_cvt_not_referred_pts')
        del occup_therapy_done

        # % of patients who had Speech and Language Therapy assessment
        speech_therapy_done = is_ich_sah_cvt_not_referred.loc[
            is_ich_sah_cvt_not_referred['SPEECHTHERAPIST_EVALUATION'].isin([1])
        ].copy()
        self._get_numbers(speech_therapy_done, '# of patients who had Speech and Language Therapy assessment', 'is_ich_sah_cvt_not_referred_pts')
        del speech_therapy_done

        # % of patients who had Cardiac arrhythmia screening
        is_tia_not_referred = df.loc[(df['STROKE_TYPE'].isin([1,3])) & (df['REFERRED_FOR_RECAN'] == 2)].copy()
        self._get_total_patients(
            is_tia_not_referred, 'is_tia_not_referred_pts', True
        )

        afib_screening = is_tia_not_referred.loc[is_tia_not_referred['AFIB_FLUTTER'].isin([3,4])].copy()
        self._get_numbers(afib_screening, '# of patients who had Cardiac arrhythmia screening', 'is_tia_not_referred_pts')
        del afib_screening

        ## Medical treatment received
        all_not_referred = df.loc[(df['REFERRED_FOR_RECAN'] == 2) & (df['HEMICRANIECTOMY'] != 3)].copy()
        self._get_total_patients(
            all_not_referred, 'all_not_referred_pts', True
        )

        # % of patients who received antihypertensive medication
        antihypertensive_prescribed = all_not_referred.loc[
            all_not_referred['ANTIHYPERTENSIVE'] == 1
        ].copy()
        self._get_numbers(antihypertensive_prescribed, '# of patients who received antihypertensive medication', 'all_not_referred_pts')
        del antihypertensive_prescribed

        # % of patients who received a statin
        statin_prescribed = stroke_not_referred.loc[
            stroke_not_referred['STATIN'] == 1
        ].copy()
        self._get_numbers(statin_prescribed, '# of patients who received a statin', 'stroke_not_referred_pts')
        del statin_prescribed

        # % of patients who received an antiplatelet agent
        def get_antithrombotics_prescribed(values, selected_values):
            ''' Return 1 if antiplatelets/anticoagulants were prescribed else return 2.
            
            :param string: the string with values
            :type string: str
            :returns: list
            '''
            selected_values_list = selected_values.split(',')
            res = [x for x in selected_values_list if int(x) in values] 
            if len(res) > 0:
                return 1
            else: 
                return 2
            
        antiplatelets = [1,2,3,4,5,6]
        all_not_referred['ANTIPLATELETS'] = all_not_referred.apply(
            lambda x: get_antithrombotics_prescribed(
                antiplatelets, x['ANTITHROMBOTICS']), 
                axis=1
        )
        antiplatelets_prescribed = all_not_referred.loc[
            all_not_referred['ANTIPLATELETS'] == 1
        ].copy()
        self._get_numbers(antiplatelets_prescribed, '# of patients who received an antiplatelet agent', 'all_not_referred_pts')
        
        # % of patients who received anticoagulation
        anticoagulation = [8,9,10,11,12,13,14]
        all_not_referred['ANTICOAGULANTS'] = all_not_referred.apply(
            lambda x: get_antithrombotics_prescribed(
                anticoagulation, x['ANTITHROMBOTICS']), 
                axis=1
        )
        anticoagulants_prescribed = all_not_referred.loc[
            all_not_referred['ANTICOAGULANTS'] == 1
        ].copy()
        self._get_numbers(anticoagulants_prescribed, '# of patients who received anticoagulation', 'all_not_referred_pts')
        
        ## Surgical treatment
        # % of patients who received hemicraniectomy
        ischemic_not_referred = df.loc[
            (df['STROKE_TYPE'] == 1) &
            (df['REFERRED_FOR_RECAN'] == 2)
        ].copy()
        self._get_total_patients(
            ischemic_not_referred, 'ischemic_not_referred_pts', True
        )

        hemicraniectomy = ischemic_not_referred.loc[
            ischemic_not_referred['HEMICRANIECTOMY'] == 1
        ].copy()
        self._get_numbers(hemicraniectomy, '# of patients who received hemicraniectomy', 'ischemic_not_referred_pts')
        del hemicraniectomy

        # % of patients who received carotid stenosis
        is_tia_not_referred_hemi = is_tia_not_referred.loc[
            is_tia_not_referred['HEMICRANIECTOMY'] != 3
        ].copy()
        self._get_total_patients(
            is_tia_not_referred_hemi, 'is_tia_not_referred_hemi_pts', True
        )

        carotid_stenosis = is_tia_not_referred_hemi.loc[
            is_tia_not_referred_hemi['CAROTID_STENOSIS'].isin([1,2])
        ].copy()
        self._get_numbers(carotid_stenosis, '# of patients who received carotid stenosis', 'is_tia_not_referred_hemi_pts')

        # % of patients who received clot surgery
        cartoid_stenosis_followup = carotid_stenosis.loc[
            carotid_stenosis['CAROTID_STENOSIS_FOLLOWUP'] == 1
        ].copy()
        self._get_numbers(cartoid_stenosis_followup, '# of patients who received clot surgery', '# of patients who received carotid stenosis')

        self.stats.rename(columns={
            'SITE_ID': 'Site ID',
            'FACILITY_NAME': 'Facility Name'
        }, inplace=True)

    def __get_ct_mri_overall(self, ct_mri):
        ''' Return 1 if CT/MRI was done else return 2. 
            
        :param ct_mri: the value of CT/MRI for ischemic stroke
        :type ct_mri: int
        :returns: int
        '''
        if ct_mri in [1,2,3,4,5,6]:
            return 1
        else:
            return 2
    
    def __get_overall_outcome(self, nihss, discharge_nihss):
        ''' Return 1 if nihss > discharge_nihss, 2 if nihss == discharge_nihss and 3 if nihss < discharge_nihss. 
            
        :param nihss: nihss score before stroke
        :type nihss: int
        :param discharge_nihss: nihss score at discharge
        :type discharge_nihss: int
        :returns: int
        '''
        if nihss > discharge_nihss:
            return 1
        elif nihss < discharge_nihss:
            return 3
        else:
            return 2

    def __get_mrs_score(self, selected_mrs):
        ''' Get mRS score from the dropdown index. 
        
        :param selected_mrs: the index from the dropdown
        :type selected_mrs: int
        :returns: converted score
        '''
        if selected_mrs == 1:
            return -2
        else:
            return selected_mrs - 2

    def __get_timestamp(self, date, time):
        ''' Convert date and time to timestamp. 

        :param date: date 
        :type date: str
        :param time: time
        :type time: str
        :returns: timestamp
        '''
        combine = f'{date} {time}'
        dateFormat = '%Y-%m-%d %H:%M:%S'
        return datetime.strptime(combine, dateFormat)

    def __get_minutes(self, start, end):
        ''' Get difference between two dates in minutes. 

        :param start: starting date
        :type start: datetime
        :param end: ending date
        :type end: datetime
        :returns: int
        '''
        minutes_diff = (end - start).total_seconds() / 60.0
        return minutes_diff
        
    def _preprocess_data(self, df):
        ''' Preprocess data. 
        
        :param df: the data from connection
        :type df: DataFrame
        :returns: DataFrame
        '''
        import numpy as np

        df['IVT_DONE'] = df.apply(
            lambda x: 1 if x['STROKE_TYPE'] == 1 and x['RECANALIZATION_PROCEDURES'] in [2,3,4,5,6] else 2,
            axis=1
        )
        df['TBY_DONE'] = df.apply(
            lambda x: 1 if x['STROKE_TYPE'] == 1 and x['RECANALIZATION_PROCEDURES'] in [4,8,9] else 2,
            axis=1
        )
        # Get CT_MRI overall for all stroke types except undetermined or sah
        df['CT_MRI_OVERALL'] = df.apply(
            lambda x: self.__get_ct_mri_overall(
                x['CT_MRI'] if x['STROKE_TYPE'] == 1 else x['CT_MRI_OTHER']), 
                axis=1)
        
        # Get overall outcome
        df['OVERALL_OUTCOME'] = df.apply(
            lambda x: self.__get_overall_outcome(
                x['NIHSS_SCORE'], x['D_NIHSS_SCORE']
                ) if x['NIHSS'] == 2 and x['D_NIHSS'] == 2 else np.nan,
            axis = 1
        )

        # Get if patients has been referred for recanaliztion
        df['REFERRED_FOR_RECAN'] = df.apply(
            lambda x: 1 if x['STROKE_TYPE'] == 1 and x['RECANALIZATION_PROCEDURES'] in [3, 5, 9] else 2,
            axis=1 
        )

        # Get mRS score
        df['DISCHARGE_MRS_SCORE'] = df.apply(
            lambda x: self.__get_mrs_score(x['DISCHARGE_MRS']), axis=1
        )

        # Get onset timestamp (onset date + onset time)
        df['ONSET_TIMESTAMP'] = df.apply(
            lambda x: self.__get_timestamp(x['ONSET_DATE'], x['ONSET_TIME']), 
            axis=1)

        # Get hospital timestamp (hospital date + hospital time)
        df['HOSPITAL_TIMESTAMP'] = df.apply(
            lambda x: self.__get_timestamp(x['HOSPITAL_DATE'], x['HOSPITAL_TIME']), 
            axis=1)
        
        # Get IVT timestamps in one column
        df['BOLUS_TIMESTAMP'] = np.nan
        df['BOLUS_TIMESTAMP'] = df.apply(
            lambda x: self.__get_timestamp(
                x['IVT_ONLY_IVT_DATE'], x['IVT_ONLY_BOLUS_TIME']
                ) if x['RECANALIZATION_PROCEDURES'] == 2 else x['BOLUS_TIMESTAMP'], 
                axis=1)
        df['BOLUS_TIMESTAMP'] = df.apply(
            lambda x: self.__get_timestamp(
                x['IVT_ONLY_REFER_ALL_IVT_DATE'], x['IVT_ONLY_REFER_ALL_BOLUS_TIME']
                ) if x['RECANALIZATION_PROCEDURES'] == 3 else x['BOLUS_TIMESTAMP'], 
                axis=1)
        df['BOLUS_TIMESTAMP'] = df.apply(
            lambda x: self.__get_timestamp(
                x['IVT_TBY_IVT_DATE'], x['IVT_TBY_BOLUS_TIME']
                ) if x['RECANALIZATION_PROCEDURES'] == 4 else x['BOLUS_TIMESTAMP'], 
                axis=1)
        df['BOLUS_TIMESTAMP'] = df.apply(
            lambda x: self.__get_timestamp(
                x['IVT_TBY_REFER_ALL_IVT_DATE'], x['IVT_TBY_REFER_ALL_BOLUS_TIME']
                ) if x['RECANALIZATION_PROCEDURES'] == 5 else x['BOLUS_TIMESTAMP'], 
                axis=1)
        df['BOLUS_TIMESTAMP'] = df.apply(
            lambda x: self.__get_timestamp(
                x['IVT_TBY_REFER_LIM_IVT_DATE'], x['IVT_TBY_REFER_LIM_BOLUS_TIME']
                ) if x['RECANALIZATION_PROCEDURES'] == 6 else x['BOLUS_TIMESTAMP'], 
                axis=1)

        # Calculate symptom onset to needle time in minutes
        df['ONSET_TO_NEEDLE_TIME'] = df.apply(lambda x: self.__get_minutes(x['ONSET_TIMESTAMP'], x['BOLUS_TIMESTAMP']) if x['IVT_DONE'] == 1 else np.nan, axis=1)   

        # Calculate door to needle time in minutes
        df['DOOR_TO_NEEDLE_TIME'] = df.apply(lambda x: self.__get_minutes(x['HOSPITAL_TIMESTAMP'], x['BOLUS_TIMESTAMP']) if x['IVT_DONE'] == 1 else np.nan, axis=1)   
        
        # Get MT timstamps in one column
        df['GROIN_PUNCTURE_TIMESTAMP'] = np.nan
        df['GROIN_PUNCTURE_TIMESTAMP'] = df.apply(
            lambda x: self.__get_timestamp(
                x['IVT_TBY_MT_DATE'], x['IVT_TBY_GROIN_PUNCTURE_TIME']
                ) if x['RECANALIZATION_PROCEDURES'] == 4 else x['GROIN_PUNCTURE_TIMESTAMP'], 
                axis=1)
        df['GROIN_PUNCTURE_TIMESTAMP'] = df.apply(
            lambda x: self.__get_timestamp(
                x['TBY_ONLY_MT_DATE'], x['TBY_ONLY_GROIN_PUNCTURE_TIME']
                ) if x['RECANALIZATION_PROCEDURES'] == 8 else x['GROIN_PUNCTURE_TIMESTAMP'], 
                axis=1)
        df['GROIN_PUNCTURE_TIMESTAMP'] = df.apply(
            lambda x: self.__get_timestamp(
                x['TBY_REFER_ALL_MT_DATE'], x['TBY_REFER_ALL_GROIN_PUNCTURE_TIME']
                ) if x['RECANALIZATION_PROCEDURES'] == 9 else x['GROIN_PUNCTURE_TIMESTAMP'], 
                axis=1)

        # Calculate symptom onset to arterial puncture time in minutes
        df['ONSET_TO_GROIN_PUNCTURE_TIME'] = df.apply(lambda x: self.__get_minutes(x['ONSET_TIMESTAMP'], x['GROIN_PUNCTURE_TIMESTAMP']) if x['TBY_DONE'] == 1 else np.nan, axis=1)   

        # Calculate door to arterial puncture time in minutes
        df['DOOR_TO_GROIN_PUNCTURE_TIME'] = df.apply(lambda x: self.__get_minutes(x['HOSPITAL_TIMESTAMP'], x['GROIN_PUNCTURE_TIMESTAMP']) if x['TBY_DONE'] == 1 else np.nan, axis=1)   

        # Get one column for TICI score
        df['TICI_SCORE'] = np.nan
        df['TICI_SCORE'] = df.apply(
            lambda x: x['IVT_TBY_TICI_SCORE'] if x['RECANALIZATION_PROCEDURES'] == 4 and x['STROKE_TYPE'] == 1 else x['TICI_SCORE'],
            axis=1
        )
        df['TICI_SCORE'] = df.apply(
            lambda x: x['TBY_ONLY_TICI_SCORE'] if x['RECANALIZATION_PROCEDURES'] == 8 and x['STROKE_TYPE'] == 1 else x['TICI_SCORE'],
            axis=1
        )
        df['TICI_SCORE'] = df.apply(
            lambda x: x['TBY_REFER_ALL_TICI_SCORE'] if x['RECANALIZATION_PROCEDURES'] == 9 and x['STROKE_TYPE'] == 1 else x['TICI_SCORE'],
            axis=1
        )
        
        df.fillna(0, inplace=True)
        return df

    def _add_group_header_format(self, workbook, color):
        ''' Add format to workbook with defined color for headers. 
        
        :param workbook: workbook to be updated
        :type workbook: Workbook
        :param color: the color in the hex format
        :type color: str
        :returns: Format
        '''
        format1 = workbook.add_format({
            'bold': 2,
            'border': 0,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': color
        })
        return format1

    def _add_group_text(self, worksheet, formatting, group_name, start_index, end_index):
        ''' Function that will add the group text with provided formatting. 
        
        :param worksheet: Worksheet to which the group text should be added
        :type worksheet: WorkSheet
        :param formatting: Format object that will be used for formatting of merged cells
        :type formatting: Format
        :param group_name: the name of group
        :type group_name: str
        :param start_index: the index of the start column
        :type start_index: int
        :param end_index: the index of the ending column
        :type end_index: int
        '''
        from xlsxwriter.utility import xl_rowcol_to_cell

        start_cell = xl_rowcol_to_cell(0, start_index)
        end_cell = xl_rowcol_to_cell(0, end_index)
        worksheet.merge_range(f'{start_cell}:{end_cell}', group_name, formatting)

    

    def __get_header(self, x):
        ''' Get the header as dictionary. 
        
        :param x: the name of column
        :type x: str
        '''
        tmp = {}
        tmp['header'] = x
        return tmp

    def _generate_formatted_preprocessed_data(self, df, filename, exclude_country=False):
        ''' Generate formatted preprocessed data. 

        :param df: the preprocessed data (filtered)
        :type df: DataFrame
        :param filename: the name of file without suffix
        :type filename: str
        :param exclude_country: if True the data for country that are appended in the beginning will be excluded, applies for site preprocessed data
        :type exclude_country: bool
        '''
        if exclude_country:
            df = df.loc[df['SITE_ID'] != self.country_name]

        workbook = xlsxwriter.Workbook(f'{filename}_preprocessed_data.xlsx')
        worksheet = workbook.add_worksheet('Preprocessed data')

        # Get number of columns and rows
        ncol = len(df.columns)
        nrow = len(df) + 1
        worksheet.set_column(0, ncol, 30)

        columns = df.columns.tolist()
        values = df.values.tolist()

        headers = [self.__get_header(name) for name in columns]

        options = {
            'data': values,
            'header_row': True,
            'columns': headers,
            'style': 'Table Style Light 8'
        }

        worksheet.add_table(1, 0, nrow, ncol - 1, options)
        workbook.close()

    
    def _generate_formatted_stats(self, df, filename):
        ''' Generate formatted statisics. 

        :param df: the calculated statisitcs
        :type df: DataFrame
        :param filename: the name of file without suffix
        :type filename: str
        '''
        # Save data into csv
        df.rename(columns={'SITE_ID': 'Site ID'}, inplace=True)
        save_file(f'{filename}.csv', data=df, index=False)

        # Remove temporary columns from csv before the data are saved into excel file
        for column in self.columns_to_be_deleted:
            if column in df.columns:
                del df[column]

        # Create new workbook
        workbook = xlsxwriter.Workbook(f'{filename}.xlsx', {'strings_to_numbers': True})
        worksheet = workbook.add_worksheet('Statistics')

        # Get number of columns and rows
        ncol = len(df.columns)
        nrow = len(df) + 1
        # Set width of columns
        worksheet.set_column(0, ncol, 30)
        
        # Get column names and values
        columns = df.columns.tolist()
        values = df.values.tolist()

        headers = [self.__get_header(name) for name in columns]

        options = {
            'data': values,
            'header_row': True,
            'columns': headers,
            'style': 'Table Style Light 8'
        }

        worksheet.add_table(1, 0, nrow, ncol - 1, options)

        gender_format = self._add_group_header_format(workbook, "#477187")
        self._add_group_text(
            worksheet, gender_format, 'GENDER', columns.index('# Male'), columns.index('% Female')
        )

        ct_mri_format = self._add_group_header_format(workbook, '#AA8739')
        self._add_group_text(
            worksheet, 
            ct_mri_format, 
            'CT or MRI done for stroke patients', 
            columns.index('# CT or MRI done for stroke patients - Yes'),
            columns.index('% CT or MRI done for stroke patients - No')
        )

        overall_outcome_format = self._add_group_header_format(workbook, '#D4B86A')
        self._add_group_text(
            worksheet, 
            overall_outcome_format, 
            'Overall outcome for all patients at time of discharge for period of assessment (as per NIHSS score)', 
            columns.index('# of patients that improved'),
            columns.index('% of patients who deteriorated in hospital')
        )

        stroke_complications_format = self._add_group_header_format(workbook, '#D4A46A')
        self._add_group_text(
            worksheet, 
            stroke_complications_format,
            'Stroke complications (all patients)',
            columns.index('# of patients with no complications'),
            columns.index('% of patients with other post-stroke complications')
        )

        nihss_on_arrival_format = self._add_group_header_format(
            workbook, '#D4916A'
        )
        self._add_group_text(
            worksheet,
            nihss_on_arrival_format,
            'NIHSS score on arrival',
            columns.index('NIHSS score on arrival (median)'),
            columns.index('NIHSS score on arrival (IQR)')
        )

        nihss_on_discharge_format = self._add_group_header_format(
            workbook, '#7F4C91'
        )
        self._add_group_text(
            worksheet,
            nihss_on_discharge_format,
            'NIHSS score on discharge',
            columns.index('NIHSS score on discharge (median)'),
            columns.index('NIHSS score on discharge (IQR)')
        )

        mrs_on_discharge_format = self._add_group_header_format(
            workbook, '#D4BA6A'
        )
        self._add_group_text(
            worksheet,
            mrs_on_discharge_format,
            'modified Ranking Score on discharge',
            columns.index('modified Ranking Score on discharge (median)'),
            columns.index('modified Ranking Score on discharge (IQR)')
        )

        intravenous_lysis_format = self._add_group_header_format(
            workbook, '#565595'
        )
        self._add_group_text(
            worksheet,
            intravenous_lysis_format,
            'Reperfusion therapy',
            columns.index('Total number stroke patients treated with intravenous lysis for period of assessment'),
            columns.index('Door to needle time (IQR)')
        )
        
        intravenou_lysis_outcome_format = self._add_group_header_format(
            workbook, '#468B78'
        )
        self._add_group_text(
            worksheet,
            intravenou_lysis_outcome_format,
            'Outcome after intravenous lysis (as per NIHSS score)',
            columns.index('# of patients that improved after intravenous lysis'),
            columns.index('% of patients that died after lysis')
        )

        mechanical_thrombectomy_format = self._add_group_header_format(
            workbook, '#B9D6C1'
        )
        self._add_group_text(
            worksheet,
            mechanical_thrombectomy_format,
            'Total number of stroke patients treated with mechanical thrombectomy for period of assessment',
            columns.index('# of all patients that received Mechanical thrombectomy - Yes'),
            columns.index('Door to arterial puncture time (IQR)')
        )

        mechanical_thrombectomy_outcome_format = self._add_group_header_format(
            workbook, '#BEBCBC'
        )
        self._add_group_text(
            worksheet,
            mechanical_thrombectomy_outcome_format,
            'Outcome after mechanical thrombectomy (as per NIHSS score)',
            columns.index('# of patients that improved after mechanical thrombectomy'),
            columns.index('% of patients that died after mechanical thrombectomy')
        )

        complications_format = self._add_group_header_format(
            workbook, '#C5D068'
        )
        self._add_group_text(
            worksheet,
            complications_format,
            'Complications of reperfusion therapy (lysis and/or thrombectomy - expressed as a percentage of the number that received reperfusion therapy)',
            columns.index('# of patients with non-fatal symptomatic intracranial haemorrhage (SICH)'),
            columns.index('% of patients with failure to achieve reperfusion (<TICI 2b or 3)')
        )

        post_acute_treatment_format = self._add_group_header_format(
            workbook, '#AA8739'
        )
        self._add_group_text(
            worksheet,
            post_acute_treatment_format,
            'Post-Acute Treatment received for period of assessment (all patients)',
            columns.index('# of patients who had a Dysphagia screen done'),
            columns.index('% of patients who had Cardiac arrhythmia screening')
        )

        medical_treatment_format = self._add_group_header_format(
            workbook, '#277650'
        )
        self._add_group_text(
            worksheet,
            medical_treatment_format,
            'Medical treatment received',
            columns.index('# of patients who received antihypertensive medication'),
            columns.index('% of patients who received anticoagulation')
        )

        surgical_treatment_format = self._add_group_header_format(
            workbook, '#AA5039'
        )
        self._add_group_text(
            worksheet,
            surgical_treatment_format,
            'Surgical treatment',
            columns.index('# of patients who received hemicraniectomy'),
            columns.index('% of patients who received clot surgery')
        )

        # Hide all columns with '#' in the name and keep only percentage values
        for column in columns:
            if column.startswith('#'):
                index = columns.index(column)
                workbook_index = xl_col_to_name(index)
                worksheet.set_column(
                    f'{workbook_index}:{workbook_index}', None, None, {'hidden': True})

        workbook.close()
 
    def _generate_presentation(self, df, filename, site_name=None):
        ''' Generate formatted statisics. 

        :param df: the calculated statisitcs
        :type df: DataFrame
        '''
        script_dir = os.path.dirname(__file__)
        master = os.path.normpath(os.path.join(script_dir, 'backgrounds', 'master.pptx'))

        prs = Presentation(master)

        # Add title to the main page
        first_slide = prs.slides[0]
        shape = first_slide.shapes[5]
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        # Set title of first slide, if site report the title contains the site name
        if site_name is None:
            run.text = f'{self.country_name}\nData Summary'
        else:
            run.text = f'{site_name}\nData Summary'
        font = run.font
        font.name = 'Centruy Gothic'
        font.size = Pt(24)
        font.color.rgb = RGBColor(250,250,250)

        # Total number of patients
        column = 'Total Patients'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        total_patients = graph_df.loc[graph_df['Facility Name'] == self.country_name, column].iloc[0]
        graph_df = graph_df.loc[graph_df['Facility Name'] != self.country_name].copy()
        title = f'Total Patients (n = {total_patients})'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # Gender distribution
        graph_df = df[['Facility Name', '% Male', '% Female']].copy()
        graph_df = graph_df.sort_values(['% Male'], ascending=True)
        legend = ['Male', 'Female']
        self._create_graph(
            prs, graph_df, title='Gender distribution (%)', graph_type='stacked', show_value_axis=True, legend=legend
        )

        # Total number of patients
        column = '% CT or MRI done for stroke patients - Yes'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'CT or MRI done for stroke patients (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)
        
        # Overall outcome for all patients at time of discharge
        graph_df = df[[
            'Facility Name', 
            '% of patients with no complications', 
            '% of patients with pneumonia', 
            '% of patients with DVT',
            '% of patients with pulmonary embolus',
            '% of patients with worsening of stroke',
            '% of patients with drip sepsis',
            '% of patients with UTI',
            '% of patients with pressure sore',
            '% of patients with death (overall mortality)',
            '% of patients with other post-stroke complications'
        ]].copy()
        graph_df = graph_df.sort_values(['% of patients with no complications'], ascending=True)
        legend = [
            'No complications',
            'Pneumonia',
            'DVT',
            'Pulmonary embolus',
            'Worsening of stroke',
            'Drip sepsis',
            'UTI',
            'Pressure sore',
            'Death (overall mortality)',
            'Other'
        ]
        self._create_graph(
            prs, 
            graph_df, 
            title='Stroke complications (all patients) (%)', 
            graph_type='stacked', 
            show_value_axis=True, 
            legend=legend
        )

        # NIHSS score on arival (median)
        column = 'NIHSS score on arrival (median)'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'NIHSS score on arrival (median)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # NIHSS score on discharge (median)
        column = 'NIHSS score on discharge (median)'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'NIHSS score on discharge (median)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # modified Ranking Score on discharge (median)
        column = 'modified Ranking Score on discharge (median)'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'modified Ranking Score on discharge (median)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)
        
        # Total number stroke patients treated with intravenous lysis for period of assessment
        column = 'Total number stroke patients treated with intravenous lysis for period of assessment'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        total_patients = graph_df.loc[graph_df['Facility Name'] == self.country_name, column].iloc[0]
        graph_df = graph_df.loc[graph_df['Facility Name'] != self.country_name].copy()
        title = f'# of stroke patients treated with intravenous lysis (n = {total_patients})'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # Total number of patients that arrived at hospital within 3 hours of symptom onset for period of assessment
        column = 'Total number of patients that arrived at hospital within 3 hours of symptom onset for period of assessment'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        total_patients = graph_df.loc[graph_df['Facility Name'] == self.country_name, column].iloc[0]
        graph_df = graph_df.loc[graph_df['Facility Name'] != self.country_name].copy()
        title = f'# of patients that arrived at hospital within 3 hours of symptom onset (n = {total_patients})'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # % of all patients that received Intravenous thrombolysis or mechanical thrombectomy - Yes
        column = '% of all patients that received Intravenous thrombolysis or mechanical thrombectomy - Yes'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = f'Intravenous thrombolysis or mechanical thrombectomy received (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # Symptom onset to needle time (median)
        column = 'Symptom onset to needle time (median)'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Symptom onset to needle time (median)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)
        
        # Door to needle time (median)
        column = 'Door to needle time (median)'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Door to needle time (median)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # Outcome after intravenous lysis (as per NIHSS score)
        graph_df = df[[
            'Facility Name', 
            '% of patients that improved after intravenous lysis', 
            '% of patients whose clinical condition did not change after intravenous lysis', 
            '% of patients who deteriorated in hospital after intravenous lysis']].copy()
        graph_df = graph_df.sort_values(['% of patients that improved after intravenous lysis'], ascending=True)
        legend = ['Improved', 'Did not change', 'Deteriorated']
        self._create_graph(
            prs, graph_df, title='Outcome after intravenous lysis (as per NIHSS score) (%)', graph_type='stacked', show_value_axis=True, legend=legend
        )

        # % of all patients that received Mechanical thrombectomy - Yes
        column = '% of all patients that received Mechanical thrombectomy - Yes'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Mechanical thrombectomy received (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)
        		
        # Symptom onset to arterial puncture time (median)
        column = 'Symptom onset to arterial puncture time (median)'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Symptom onset to arterial puncture time (median)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)
        
        # Door to arterial puncture time (median)
        column = 'Door to arterial puncture time (median)'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Door to arterial puncture time (median)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # Outcome after mechanical thrombectomy (as per NIHSS score)
        graph_df = df[[
            'Facility Name', 
            '% of patients that improved after mechanical thrombectomy', 
            '% of patients whose clinical condition did not change after mechanical thrombectomy', 
            '% of patients who deteriorated in hospital after mechanical thrombectomy']].copy()
        graph_df = graph_df.sort_values(['% of patients that improved after mechanical thrombectomy'], ascending=True)
        legend = ['Improved', 'Did not change', 'Deteriorated']
        self._create_graph(
            prs, graph_df, title='Outcome after mechanical thrombectomy (as per NIHSS score) (%)', graph_type='stacked', show_value_axis=True, 
            legend=legend
        )

        # Complications of reperfusion therapy
        graph_df = df[[
            'Facility Name', 
            '% of patients with non-fatal symptomatic intracranial haemorrhage (SICH)', 
            '% of patients with fatal SICH', 
            '% of patients with remote cerebral haemorrhage',
            '% of patients with vessel perforation',
            '% of patients with procedure related dissection',
            '% of patients with embolization to different vascular territory',
            '% of patients with groin haematoma requiring transfusion',
            '% of patients who died (mortality)',
            '% of patients with failure to achieve reperfusion (<TICI 2b or 3)',
        ]].copy()
        graph_df = graph_df.sort_values(
            ['% of patients with non-fatal symptomatic intracranial haemorrhage (SICH)'], ascending=True)
        legend = [
            'Non-fatal symptomatic intracranial haemorrhage (SICH)',
            'Fatal SICH',
            'Remote cerebral haemorrhage',
            'Vessel perforation',
            'Procedure related dissection',
            'Embolization to different vascular territory',
            'Groin haematoma requiring transfusion',
            'Died (mortality)',
            'Failure to achieve reperfusion (<TICI 2b or 3)'
        ]
        self._create_graph(
            prs, 
            graph_df, 
            title='Complications of reperfusion therapy (%)', 
            graph_type='stacked', 
            show_value_axis=True, 
            legend=legend
        )
               
        # % of patients who had a Dysphagia screen done
        column = '% of patients who had a Dysphagia screen done'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Dysphagia screening done (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)         			

        # % of patients who had physiotherapy assessment
        column = '% of patients who had physiotherapy assessment'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Physiotherapy assessment (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)    

        # % of patients who had Occupational therapy assessment
        column = '% of patients who had Occupational therapy assessment'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Occupational therapy assessment (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)   

        # % of patients who had Speech and Language Therapy assessment
        column = '% of patients who had Speech and Language Therapy assessment'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Speech and Language therapy assessment (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)   

        # % of patients who had Cardiac arrhythmia screening
        column = '% of patients who had Cardiac arrhythmia screening'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Cardiac arrhythmia screening (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # % of patients who received antihypertensive medication
        column = '% of patients who received antihypertensive medication'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Antihypertensive medication received (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # % of patients who received a statin
        column = '% of patients who received a statin'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Statin received (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # % of patients who received an antiplatelet agent
        column = '% of patients who received an antiplatelet agent'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'An antiplatelet agent received (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # % of patients who received anticoagulation
        column = '% of patients who received anticoagulation'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Anticoagulation received (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # % of patients who received hemicraniectomy
        column = '% of patients who received hemicraniectomy'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Hemicraniectomy received (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # % of patients who received carotid stenosis
        column = '% of patients who received carotid stenosis'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Carotid stenosis received (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        # % of patients who received clot surgery
        column = '% of patients who received clot surgery'
        graph_df = df[['Facility Name', column]].copy()
        graph_df = graph_df.sort_values([column], ascending=True)
        title = 'Clot surgery received (%)'
        self._create_graph(prs, graph_df, title=title, show_value_axis=False)

        working_dir = os.getcwd()
        pptx = f'{filename}.pptx'
        prs.save(os.path.normpath(os.path.join(working_dir, pptx)))

    def _create_graph(self, presentation, df, title, graph_type='barplot', show_value_axis=True, legend=None):
        ''' Generate graph based on setting. 
        
        :param df: the dataframe containing columns to be displayed
        :type df: DataFrame
        :param presentation: the presentation object
        :type presentation: Presentation
        :param title: the title of the graph
        :type title: str
        :param graph_type: the type of the graph
        :type graph_type: str
        :param legend: list of legend
        :type legend: list
        '''
        # create dictioanry of columns with index 
        colors = {
            0: RGBColor(43, 88, 173), # dark blue
            1: RGBColor(237, 125, 49), # orange
            2: RGBColor(165, 165, 165), # gray
            3: RGBColor(255, 192, 0), # yellow
            4: RGBColor(136, 106, 159),
            5: RGBColor(98, 153, 62), # green
            6: RGBColor(151, 185, 224), # light blue
            7: RGBColor(241, 167, 138), # beige     
            8: RGBColor(199, 124, 169),
            9: RGBColor(117, 231, 118)
        }

        font_name = 'Century Gothic'
        category_column = 'Facility Name'
        # Add new slide to the presentation
        slide = presentation.slides.add_slide(presentation.slide_layouts[11])
        # Get title object
        title_placeholders = slide.shapes.title
        # Set title
        title_placeholders.text = title

        # Get list of column names
        column_names = df.columns.tolist()
        # Get rest of columns without category column
        index = column_names.index(category_column) + 1
        series_columns = column_names[index:]

        # Create ChartData object and set categories
        chart_data = ChartData()
        categories = df[category_column].tolist()
        chart_data.categories = categories

        # Define specification where the chart will be placed
        specs = {
            'height': Cm(16.5),
            'width': Cm(32),
            'left': Cm(0.7),
            'top': Cm(2)
        }

        if graph_type == 'barplot':
            # Add series to the graph
            chart_data.add_series(series_columns[0], df[series_columns[0]].tolist())
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, 
                specs['left'],
                specs['top'], 
                specs['width'],
                specs['height'], 
                chart_data
            ).chart        

            # Get plot
            plot = chart.plots[0]
            # Set for each bar same color
            plot.vary_by_categories = False
            # Show data labels 
            plot.has_data_labels = True
            # Change gap width
            plot.gap_width = 100

            series = chart.series[0]
            # If there is more then 2 categories in the dataframe, the country bar and the region bar will be colored with different color to be distinguished, else the color will be blue
            if (len(df) > 2):
                for idx, point in enumerate(series.points):
                    fill = point.format.fill
                    fill.solid()
                    if (categories[idx] == self.country_name):
                        fill.fore_color.rgb = RGBColor(128,0,0)
                    elif (
                        self.site_reports and 
                        categories[idx] == self.region_name and 
                        self.region_name is not None
                    ):
                        fill.fore_color.rgb = RGBColor(124,124,124)
                    else:
                        fill.fore_color.rgb = RGBColor(43, 88, 173)
            else:
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(43, 88, 173) 

            # Show values at bar
            value_axis = chart.value_axis
            value_axis.visible = show_value_axis
            value_axis.has_major_gridlines = False

            if show_value_axis:
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)
                tick_labels.font.name = font_name

                value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
                # Set range of axis
                if '%' in title:
                    maximum = 100
                else:
                    maximum = round((max(df[series_columns[1]].tolist())), 1)
                value_axis.maximum_scale = maximum
                value_axis.minimum_scale = 0

            category_axis = chart.category_axis
            category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
            solidFill = category_axis.format.line.color._xFill
            self.__set_transparency(100, solidFill)

            # Delete tick marks
            category_axis.major_tick_mark = XL_TICK_MARK.NONE
            category_axis.major_unit = 1
            category_labels = category_axis.tick_labels
            category_labels.font.size = Pt(10)
            category_labels.font.name = font_name

        else:
            # If more series should be shown, add them together with coressponding legend label
            for idx, col in enumerate(series_columns):
                chart_data.add_series(legend[idx], df[col].tolist())      

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_STACKED, 
                specs['left'],
                specs['top'], 
                specs['width'],
                specs['height'], 
                chart_data
            ).chart
        
            # Get plot
            plot = chart.plots[0]
            plot.gap_width = 100

            series = chart.series
            for s_idx, serie in enumerate(series):
                # Set color ofthe first series to dark blue
                fill = serie.format.fill
                fill.solid()
                fill.fore_color.rgb = colors[s_idx]

                if (len(df) > 2):
                    # ---add an `a:alpha` child element---
                    solidFill = fill.fore_color._xFill
                    self.__set_transparency(30, solidFill)

                    # Change color of borders of series and transparency
                    serie.format.line.color.rgb = colors[s_idx]
                    solidFill = serie.format.line.color._xFill
                    self.__set_transparency(70, solidFill)

                    # Remove transparency from country point 
                    for idx, point in enumerate(serie.points):
                        if (categories[idx] == self.country_name):
                            point.format.line.color.rgb = colors[s_idx]
                            # Get fill of point for country
                            fill = point.format.fill
                            fill.solid()
                            fill.fore_color.rgb = colors[s_idx]
                       

            value_axis = chart.value_axis
            value_axis.has_major_gridlines = False
            value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
            tick_labels = value_axis.tick_labels
            tick_labels.font.size = Pt(10)
            tick_labels.font.name = font_name
           
            if (len(df) > 2):
                value_axis.has_major_gridlines = True
                value_axis.major_gridlines.format.line.width = Pt(0.5)
                value_axis.major_gridlines.format.line.color.rgb = RGBColor(166, 166, 166) # Set color to gray (A6A6A6)

                # Set 100% transparency to value axis
                value_axis.format.line.color.rgb = RGBColor(0, 0, 0)
                solidFill = value_axis.format.line.color._xFill
                self.__set_transparency(100, solidFill)

            value_axis.maximum_scale = 100
            value_axis.minimum_scale = 0

            # Value for y-axis (change font size, name, and other things)
            category_axis = chart.category_axis
            
            # Set 100% transparency to category axis
            category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
            solidFill = category_axis.format.line.color._xFill
            self.__set_transparency(100, solidFill)

            category_axis.major_tick_mark = XL_TICK_MARK.NONE
            category_labels = category_axis.tick_labels
            category_labels.font.size = Pt(10)
            category_labels.font.name = font_name
            category_labels.tickLblSkip = 1

            # Set legend 
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.TOP
            chart.legend.include_in_layout = False
            chart.legend.font.name = font_name



    def __set_transparency(self, transparency, elm):
        """ The function set the transparency of the row. 

        :param transparency: the transparency in %
        :type transparency: int
        :param elm: the element which transparency should be changed
        :type elm: format.line.color._xFill
        """
        a = str(100 - transparency) + '196'
        alpha = OxmlElement('a:alpha')
        alpha.set('val', a)
        elm.srgbClr.append(alpha)









        
        
        

        

        


